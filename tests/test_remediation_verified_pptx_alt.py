"""The 1.1.1 pptx alt-text lane, proved end to end (WCAG 1.1.1 Non-text Content).

The third lane to meet the REMEDIATION-VERIFIED bar, after 2.4.4 on docx and xlsx, and the same
five properties are asserted: the original document trips the finding, an approval changes the
document, a REAL re-scan verifies it, unrelated content survives, and a failed engine earns no
credit. Nothing but the blob store is patched — `handlers._apply_approved_values` runs the
production seam through `proposals.verify_residual_scs` to `scanner.analyse_and_assess`.

WHY THE APPROVED VALUE IS AUTHORED HERE RATHER THAN PROPOSED. The other two lanes take their
value from the deterministic proposer (`derive_link_text`), so the test can drive the whole chain
including the draft. 1.1.1's proposer is an AI VISION model: on a host without one, a
proposer-driven test would skip, and a lane that can only be proved where a model happens to be
available is not proved. A reviewer authoring alt text from scratch is also the real workflow for
any image a model cannot describe — `handlers` reads the approved value the same way whichever
produced it, so the lane under test is identical.

WHAT THIS CLAIMS, and it is narrower than "the image is now accessible". A `descr` a screen
reader will announce is present, the deck survived, and ACP's own criterion stops firing.
Whether the sentence DESCRIBES the picture is a human judgement no detector makes — 1.1.1 is a
REVIEW-lane criterion for exactly that reason (ADR 0016), and a clean scan here resolves to
REVIEW, never to PASS. This file proves the write and the verification, not the quality of the
words.
"""
from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

pytest.importorskip("pptx")

FILE = "review.pptx"
SID = "rv-pptx-111"
SLIDE = "ppt/slides/slide1.xml"

TITLE = "Q3 accessibility review"
BODY = "Unrelated body copy that must survive the write."
APPROVED = "A bar chart of Q3 findings, grouped by severity."
SECOND = "A photograph of the accessibility working group."

# A 1x1 PNG — the smallest thing that is genuinely a picture part, so the deck exercises the real
# image relationship rather than a shape pretending to be one.
_PNG = (b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00"
        b"\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4"
        b"\x00\x00\x00\x00IEND\xaeB`\x82")


def _deck(*, pictures: int = 1, alts: tuple[str | None, ...] = (None,)) -> bytes:
    """A deck with `pictures` images, each given the corresponding entry of `alts` as its descr.

    Two pictures is not padding: it is what makes
    `test_a_partial_write_is_not_credited_because_the_criterion_still_fails` able to fail. A lane
    that credited on "the writer wrote something" rather than on the re-scan would pass a
    one-image deck and be wrong about every real one.
    """
    from pptx import Presentation
    from pptx.util import Inches

    img = Path(tempfile.mkdtemp()) / "i.png"
    img.write_bytes(_PNG)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = TITLE
    for n in range(pictures):
        pic = slide.shapes.add_picture(str(img), Inches(1 + 2 * n), Inches(2), Inches(1.5), Inches(1.5))
        alt = alts[n] if n < len(alts) else None
        if alt:
            pic._element._nvXxPr.cNvPr.set("descr", alt)
    box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1))
    box.text_frame.text = BODY

    out = Path(tempfile.mkdtemp()) / FILE
    prs.save(out)
    return out.read_bytes()


def _assess(data: bytes) -> set[str]:
    """The SCs a REAL assessment reports — the same call the production re-verification makes."""
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / FILE).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), FILE, detect_pii=False)
    return {sc for i in (fd or {}).get("issues", []) if (sc := _extract_sc(i.get("wcag", "")))}


def _locators(data: bytes) -> list[str]:
    """The locators the DETECTOR emits, read off a real scan rather than constructed here — a
    hand-built locator that happens to match is not evidence the two agree."""
    from scanner import analyse_and_assess
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / FILE).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), FILE, detect_pii=False)
    return [i["locator"] for i in (fd or {}).get("issues", [])
            if "1.1.1" in (i.get("wcag") or "") and i.get("locator")]


def _slide_xml(data: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(SLIDE).decode("utf-8")


class _Blob:
    """The only thing patched in this module. Stores bytes verbatim; decides nothing."""

    def __init__(self, data: bytes):
        self.data, self.uploads = data, []

    def download_remediated(self, owner, sid, f):
        return self.data

    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data
        self.uploads.append((f, mime))
        return "http://b/2"


@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rv.db")
    return store_mod.Store()


def _seed(store, values: dict[str, str]) -> int:
    """A scanned + remediated deck with one 1.1.1 card, and `values` approved on it."""
    store.init_scan_run(SID, "drive", 1, "2026-08-31T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": FILE, "engine": "office", "status": "pass", "score": 60, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "d1",
        "issues": [{"ruleId": "PPTX_IMAGE_NO_ALT", "wcag": "1.1.1 Non-text Content",
                    "severity": "SERIOUS", "locator": loc} for loc in values],
    }, "2026-08-31T00:00:00Z")
    store.record_remediation(SID, FILE, drive_write_url="http://d/1", blob_url="http://b/1")
    item_id = store.enqueue_proposals(SID, FILE, "1.1.1", [
        {"locator": loc, "before": "(no alt text)", "proposed_value": "", "rationale": "r",
         "source": "reviewer"} for loc in values
    ], rule_name="Non-text Content")
    store.update_hitl_item(item_id, "approved", None, None)
    store.approve_proposal_values(item_id, list(values.values()))
    return item_id


def _run_lane(monkeypatch, store, blob):
    """The production handler, with the re-scan UNPATCHED."""
    import core
    import handlers
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    handlers._apply_approved_values({"scan_id": SID, "file": FILE}, {})


@pytest.fixture(scope="module")
def undescribed() -> bytes:
    return _deck(pictures=1, alts=(None,))


# ── 1. the finding ────────────────────────────────────────────────────────────

def test_a_real_assessment_reports_1_1_1_on_an_undescribed_picture(undescribed):
    assert "1.1.1" in _assess(undescribed)


def test_a_described_picture_is_not_flagged(undescribed):
    """The control. Without it, a detector that flagged every picture would satisfy the test
    above and the whole file would be measuring nothing."""
    assert "1.1.1" not in _assess(_deck(pictures=1, alts=(APPROVED,)))


def test_the_finding_carries_a_locator_the_writer_can_resolve(undescribed):
    """The join between detection and writing. These are different modules with different
    parsers, and a locator the writer cannot resolve makes the lane unreachable however good
    both halves are on their own."""
    from apply_alt import parse_locator
    locs = _locators(undescribed)
    assert locs, "the finding carries no locator, so no approval can be aimed at it"
    assert parse_locator(locs[0]) == (SLIDE, "Picture 2"), locs[0]


# ── 2. approval → write → re-scan → credit, through the real path ─────────────

@pytest.fixture()
def applied(store, monkeypatch, undescribed):
    blob = _Blob(undescribed)
    _seed(store, {_locators(undescribed)[0]: APPROVED})
    _run_lane(monkeypatch, store, blob)
    return blob, store


def test_the_saved_deck_carries_the_approved_description(applied):
    blob, _ = applied
    assert f'descr="{APPROVED}"' in _slide_xml(blob.data)


def test_unrelated_content_survives(applied):
    blob, _ = applied
    xml = _slide_xml(blob.data)
    assert TITLE in xml and BODY in xml


def test_the_deck_still_opens(applied):
    """Through python-pptx, which had no part in writing the change."""
    from pptx import Presentation
    blob, _ = applied
    p = Path(tempfile.mkdtemp()) / FILE
    p.write_bytes(blob.data)
    assert zipfile.ZipFile(p).testzip() is None
    prs = Presentation(str(p))
    assert any(sh.has_text_frame and TITLE in sh.text_frame.text
               for sh in prs.slides[0].shapes)


def test_a_second_real_assessment_no_longer_reports_1_1_1(applied):
    """THE claim: a fresh assessment of the SAVED bytes, not the writer's return value."""
    blob, _ = applied
    assert "1.1.1" not in _assess(blob.data)


def test_the_row_is_credited_and_the_copy_is_stored(applied):
    blob, store = applied
    assert store.count_unapplied_approved_values(SID, FILE) == 0
    assert blob.uploads


# ── 3. where the lane must NOT credit ────────────────────────────────────────

def test_a_partial_write_is_not_credited_because_the_criterion_still_fails(store, monkeypatch):
    """Two undescribed pictures, one approved. The write succeeds and the deck genuinely
    improves — and 1.1.1 still fails, because the other picture is still silent.

    This is the control that separates "the writer wrote something" from "the criterion
    cleared". A lane crediting on the write would mark the file compliant here and publish a
    deck that still fails the criterion it was certified against.
    """
    deck = _deck(pictures=2, alts=(None, None))
    locs = _locators(deck)
    assert len(locs) == 2, f"expected two undescribed pictures, got {locs}"

    blob = _Blob(deck)
    _seed(store, {locs[0]: APPROVED})
    _run_lane(monkeypatch, store, blob)

    assert "1.1.1" in _assess(blob.data), (
        "the fixture no longer fails after describing one of two pictures — this control cannot "
        "distinguish a withheld credit from a cleared one")
    assert store.count_unapplied_approved_values(SID, FILE) == 1, (
        "the value was credited even though the criterion still fails on re-scan")
    assert store.mark_file_compliant_if_reviewed(SID, FILE) is False
    assert not blob.uploads, "an uncleared write was published as the corrected copy"


def test_describing_both_pictures_does_clear_it(store, monkeypatch):
    """The other half of the pair, so the test above is known to be about COUNT rather than
    about the lane being broken for multi-image decks."""
    deck = _deck(pictures=2, alts=(None, None))
    locs = _locators(deck)
    blob = _Blob(deck)
    _seed(store, {locs[0]: APPROVED, locs[1]: SECOND})
    _run_lane(monkeypatch, store, blob)

    assert "1.1.1" not in _assess(blob.data)
    assert store.count_unapplied_approved_values(SID, FILE) == 0
    assert blob.uploads


def test_an_already_described_deck_is_left_alone(store, monkeypatch):
    deck = _deck(pictures=1, alts=(APPROVED,))
    assert "1.1.1" not in _assess(deck)
    blob = _Blob(deck)
    _seed(store, {})
    _run_lane(monkeypatch, store, blob)
    assert blob.data == deck and not blob.uploads


# ── 4. a broken engine earns nothing ─────────────────────────────────────────

@pytest.mark.parametrize("name,script,timeout", [
    ("cannot be launched", None, None),
    ("exits non-zero", "#!/bin/sh\necho boom >&2\nexit 9\n", None),
    ("hangs past the timeout", "#!/bin/sh\nsleep 30\n", "2"),
])
def test_a_broken_office_analyser_never_credits_this_lane_either(monkeypatch, undescribed,
                                                                 name, script, timeout):
    """The guarantee #1058 established for docx, asserted for this lane rather than assumed to
    transfer: the fail-open lived in one shared seam, so a regression there would take every
    lane with it at once.

    The re-scan must return a real answer — the first-party pptx checks run after the .NET call
    in their own try/except — and never None, which `_apply_one_value_kind` reads as "cleared".
    """
    import stat as _stat

    import scanner
    if script is None:
        monkeypatch.setattr(scanner, "DOTNET", "/nonexistent/dotnet", raising=False)
    else:
        fake = Path(tempfile.mkdtemp()) / "dotnet"
        fake.write_text(script)
        fake.chmod(fake.stat().st_mode | _stat.S_IEXEC)
        monkeypatch.setattr(scanner, "DOTNET", str(fake), raising=False)
    if timeout:
        monkeypatch.setenv("ACP_OFFICE_CLI_TIMEOUT", timeout)

    from proposals import verify_residual_scs
    residual = verify_residual_scs(undescribed, FILE)
    assert residual is not None, (
        f"an office CLI that {name} made the re-scan return None — every approved value on this "
        f"lane would be credited on a scan that never happened")
    assert "1.1.1" in residual

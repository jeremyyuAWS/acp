"""The locators the proposers actually mint must reach the images they name.

remediate_office reaches an image in two ways, and mints a different locator for each:

  * the decorative-inference branch never opens the image, and names the ELEMENT
    (`ppt/slides/slide1.xml#Picture 3`);
  * every branch that reads the image's BYTES — the vision alt draft, the xlsx chart-data
    draft, the evidence thumbnail — got there through the blip's relationship id via
    `_image_bytes_for`, and mints THAT (`ppt/slides/slide1.xml#rId2`).

apply_alt understood only the first. An rId locator looked for an element named "rId2", found
none, and was reported unresolved — so the reviewer's approved description was stored, counted
by store.count_unapplied_approved_values as content the document owed, handed to the writer,
dropped, and the row never marked applied. The file could never certify or reach Publish. That
is the whole vision alt-text lane: the flagship AI-drafts / human-approves path for Office.

It stayed invisible for two reasons, and the tests here close both. An unresolved locator is
only logged, never raised — correct behaviour for a genuinely missing element, and silence for
a systematically wrong lookup. And tests/test_apply_approved_values.py seeds NAME locators, so
no test ever fed apply_alt what the proposers emit; the first test below takes its locators from
the real proposer instead of writing one by hand, which is the only way this class of drift gets
caught rather than re-created.
"""
from __future__ import annotations
import io
import random
import re
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

import apply_alt  # noqa: E402

SLIDE = "ppt/slides/slide1.xml"
DOC = "word/document.xml"
DRAWING = "xl/drawings/drawing1.xml"
FILE = "deck.pptx"
SID = "s-rid"


def _pkg(part: str, xml: str) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(part, xml)
        z.writestr("docProps/core.xml", "<cp:coreProperties/>")
    return buf.getvalue()


def _read(data: bytes, part: str) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(part).decode("utf-8")


def _pic(el: str, cid: str, name: str, rid: str, descr: str | None = None) -> str:
    """One pic block: the alt-bearing element, then the blip that names its image."""
    d = f' descr="{descr}"' if descr is not None else ""
    return (f'<p:pic><p:nvPicPr><p:cNvPr id="{cid}" name="{name}"{d}/></p:nvPicPr>'
            f'<p:blipFill><a:blip r:embed="{rid}"/></p:blipFill></p:pic>')


def _slide(*pics: str, group: bool = True) -> str:
    """A slide part. `group` adds the root group-shape cNvPr every real slide carries — an
    alt-bearing element with no image of its own, which must never absorb a pic's rId."""
    root = '<p:nvGrpSpPr><p:cNvPr id="1" name=""/></p:nvGrpSpPr>' if group else ""
    return (f'<p:sld xmlns:p="p" xmlns:a="a" xmlns:r="r">{root}{"".join(pics)}</p:sld>')


def _noisy_png() -> bytes:
    """High-variance pixels, so proposals.infer_decorative does NOT claim the image is a
    near-solid block — that branch mints a NAME locator and would bypass what we are testing."""
    from PIL import Image
    rnd = random.Random(7)
    im = Image.new("RGB", (400, 300))
    im.putdata([(rnd.randrange(256), rnd.randrange(256), rnd.randrange(256))
                for _ in range(400 * 300)])
    buf = io.BytesIO()
    im.save(buf, "PNG")
    return buf.getvalue()


# ── the regression, driven by the real proposer ───────────────────────────────

@pytest.fixture()
def real_deck(tmp_path):
    """A genuine pptx with one describable image, saved by python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_noisy_png()), Inches(1), Inches(1), Inches(3), Inches(2))
    path = tmp_path / FILE
    prs.save(str(path))
    return path.read_bytes()


def test_the_vision_proposers_own_locator_reaches_its_image(real_deck, monkeypatch):
    """THE regression, and deliberately end to end: the locator is taken from the real proposer,
    never written by hand. A test that hand-writes `#Picture 1` passes while production mints
    `#rId2` and every approved description silently fails to land — which is exactly what
    happened."""
    import ai
    import remediate_office as ro
    monkeypatch.setattr(ai, "vision_is_available", lambda: True)
    monkeypatch.setattr(ai, "describe_image_structured", lambda img, **k: {
        "alt": "A dense field of coloured noise.", "model": "stub-vision",
        "grounded": False, "evidence": "vision"}, raising=False)

    proposals, _evidence = ro.alt_proposals_for_office(
        real_deck, "pptx", ai_enabled=True, context_file=FILE, scan_id=SID)
    assert proposals, "the vision proposer produced no card — the fixture no longer exercises it"
    locator = proposals[0]["locator"]
    assert re.search(r"#rId\d+$", locator), (
        f"expected the bytes-reading proposer's relationship-id locator, got {locator!r}")

    fixed, applied, unresolved = apply_alt.apply_alt_text(
        real_deck, {locator: proposals[0]["proposed_value"]})

    assert unresolved == [], "the writer could not resolve the locator the proposer minted"
    assert [a["locator"] for a in applied] == [locator]
    slide = next(n for n in zipfile.ZipFile(io.BytesIO(fixed)).namelist()
                 if n.startswith("ppt/slides/slide"))
    assert 'descr="A dense field of coloured noise."' in _read(fixed, slide)


def test_an_rid_locator_used_to_resolve_to_nothing(real_deck):
    """The bug on demand: name-only lookup, which is all this module used to do. Keeps the test
    above honest — it would pass just as well if resolution had been loosened into a guess."""
    xml = _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"))
    els = apply_alt._alt_elements(xml, "p:cNvPr")
    assert not [m for m in els if apply_alt._ATTR(m.group(1), "name").strip() == "rId2"]
    assert apply_alt.resolve_target(xml, "p:cNvPr", "rId2") is not None


# ── resolution, both fragment kinds ───────────────────────────────────────────

def test_a_relationship_id_lands_on_its_own_image_not_a_sibling():
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"),
                              _pic("p:cNvPr", "4", "Chart 4", "rId7")))
    fixed, applied, unresolved = apply_alt.apply_alt_text(data, {f"{SLIDE}#rId7": "Q3 revenue."})

    xml = _read(fixed, SLIDE)
    assert 'name="Chart 4" descr="Q3 revenue."' in xml
    assert 'name="Picture 3"/>' in xml           # the sibling is untouched
    assert unresolved == [] and len(applied) == 1


def test_name_locators_still_work_unchanged():
    """The decorative-inference branch mints these, and every existing stored row holds them."""
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2")))
    fixed, applied, unresolved = apply_alt.apply_alt_text(
        data, {f"{SLIDE}#Picture 3": "A clinician at a desk."})

    assert 'name="Picture 3" descr="A clinician at a desk."' in _read(fixed, SLIDE)
    assert unresolved == [] and applied[0]["before"] == "(no alt text)"


def test_a_shape_actually_named_like_a_relationship_id_wins():
    """Name is tried first, so a document that literally names a shape "rId2" resolves to the
    shape the reviewer was looking at rather than to whatever image holds that relationship."""
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"),
                              _pic("p:cNvPr", "4", "rId2", "rId9")))
    fixed, _applied, _unresolved = apply_alt.apply_alt_text(data, {f"{SLIDE}#rId2": "Chosen."})

    xml = _read(fixed, SLIDE)
    assert 'name="rId2" descr="Chosen."' in xml
    assert 'name="Picture 3"/>' in xml


def test_a_shape_with_no_image_never_absorbs_a_relationship():
    """A slide's root group cNvPr precedes every pic and carries no blip of its own. Scanning
    forward without stopping at the next alt-bearing element would hand it the first pic's rId
    and write the description onto the group — a description on nothing at all."""
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2")))
    fixed, applied, _unresolved = apply_alt.apply_alt_text(data, {f"{SLIDE}#rId2": "The photo."})

    xml = _read(fixed, SLIDE)
    assert 'name="Picture 3" descr="The photo."' in xml
    assert 'name=""/>' in xml                    # the group shape got nothing
    assert len(applied) == 1


def test_an_unknown_relationship_is_reported_never_guessed():
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2")))
    fixed, applied, unresolved = apply_alt.apply_alt_text(data, {f"{SLIDE}#rId99": "Nowhere."})

    assert unresolved == [f"{SLIDE}#rId99"] and applied == []
    assert fixed == data                          # untouched bytes, not a rezipped copy
    assert "Nowhere." not in _read(fixed, SLIDE)


# ── several images in one part ────────────────────────────────────────────────

def test_writes_to_one_part_do_not_displace_later_targets():
    """Each write changes the part's length, so a target resolved before the write would point
    at the wrong offset afterwards. The later image is written FIRST here, so an implementation
    that resolved everything up front lands the second description in the wrong place."""
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"),
                              _pic("p:cNvPr", "4", "Chart 4", "rId7")))
    fixed, applied, unresolved = apply_alt.apply_alt_text(data, {
        f"{SLIDE}#rId7": "The later image, written first.",
        f"{SLIDE}#Picture 3": "The earlier image, written second.",
    })

    xml = _read(fixed, SLIDE)
    assert 'name="Picture 3" descr="The earlier image, written second."' in xml
    assert 'name="Chart 4" descr="The later image, written first."' in xml
    assert unresolved == [] and len(applied) == 2


def test_an_existing_description_is_replaced_and_reported_as_before():
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2", descr="image.png")))
    fixed, applied, _unresolved = apply_alt.apply_alt_text(data, {f"{SLIDE}#rId2": "A real one."})

    assert applied[0]["before"] == "image.png"
    xml = _read(fixed, SLIDE)
    assert 'descr="A real one."' in xml and "image.png" not in xml


# ── the other two Office formats ──────────────────────────────────────────────

def test_docx_resolves_the_relationship_inside_the_drawing():
    """In wordprocessing the wp:docPr precedes the blip within the same w:drawing — the same
    adjacency, so the same rule, but pic-block spans do not exist here to lean on."""
    xml = ('<w:document xmlns:w="w" xmlns:wp="wp" xmlns:a="a" xmlns:r="r"><w:body>'
           '<w:drawing><wp:docPr id="1" name="Picture 1"/><a:blip r:embed="rId5"/></w:drawing>'
           '<w:drawing><wp:docPr id="2" name="Picture 2"/><a:blip r:embed="rId6"/></w:drawing>'
           "</w:body></w:document>")
    fixed, applied, unresolved = apply_alt.apply_alt_text(_pkg(DOC, xml), {f"{DOC}#rId6": "Fig 2."})

    assert 'name="Picture 2" descr="Fig 2."' in _read(fixed, DOC)
    assert 'name="Picture 1"/>' in _read(fixed, DOC)
    assert unresolved == [] and len(applied) == 1


def test_xlsx_drawing_resolves_the_relationship():
    xml = ('<xdr:wsDr xmlns:xdr="xdr" xmlns:a="a" xmlns:r="r">'
           '<xdr:pic><xdr:nvPicPr><xdr:cNvPr id="2" name="Chart 1"/></xdr:nvPicPr>'
           '<xdr:blipFill><a:blip r:embed="rId1"/></xdr:blipFill></xdr:pic></xdr:wsDr>')
    fixed, applied, unresolved = apply_alt.apply_alt_text(
        _pkg(DRAWING, xml), {f"{DRAWING}#rId1": "Revenue by quarter."})

    assert 'name="Chart 1" descr="Revenue by quarter."' in _read(fixed, DRAWING)
    assert unresolved == [] and len(applied) == 1


# ── the namespace axis, now owned upstream ───────────────────────────────────
# A default-namespace xlsx drawing (<pic>/<cNvPr>, openpyxl-family output) used to be
# unresolvable here for a SECOND reason: _ALT_TAG_FOR_PART was a hand-copy of the producer's
# table and kept the prefixed-only `xdr:cNvPr`. #50 removed the copy — the table is derived
# from formats.office.images.ALT_TARGETS now, and tests/test_apply_alt.py asserts the mirror
# against that source. Only the INTERSECTION is still ours: a relationship-id locator inside a
# default-namespace part exercises both axes at once, and neither test file covers it alone.

def test_a_relationship_id_resolves_in_a_default_namespace_drawing():
    xml = ('<wsDr xmlns="xdr" xmlns:a="a" xmlns:r="r">'
           '<pic><nvPicPr><cNvPr id="2" name="Chart 1"/></nvPicPr>'
           '<blipFill><a:blip r:embed="rId1"/></blipFill></pic></wsDr>')
    fixed, applied, unresolved = apply_alt.apply_alt_text(
        _pkg(DRAWING, xml), {f"{DRAWING}#rId1": "Revenue by quarter."})

    assert unresolved == [] and len(applied) == 1
    out = _read(fixed, DRAWING)
    assert '<cNvPr id="2" name="Chart 1" descr="Revenue by quarter."/>' in out
    assert "?:" not in out          # the pattern itself must never reach the document


def test_a_decorative_marking_also_reaches_an_image_by_relationship():
    """Both fragment kinds now go through resolve_target, so the decorative lane (#43) inherits
    this fix. It matters for the branches that resolve an image by its bytes and then find it
    decorative — those mint an rId, and their marking used to land nowhere either."""
    data = _pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"),
                              _pic("p:cNvPr", "4", "Chart 4", "rId7")))
    fixed, applied, unresolved = apply_alt.apply_alt_text(
        data, {}, decorative=[f"{SLIDE}#rId7"])

    xml = _read(fixed, SLIDE)
    assert unresolved == [] and len(applied) == 1
    assert "adec:decorative" in xml and 'val="1"' in xml
    assert xml.index("Chart 4") < xml.index("adec:decorative")   # marked the right image
    assert 'name="Picture 3"/>' in xml                           # sibling untouched


# ── the loop this unblocks ────────────────────────────────────────────────────

@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rid.db")
    return store_mod.Store()


class _Blob:
    def __init__(self, data): self.data, self.uploads = data, []
    def download_remediated(self, owner, sid, f): return self.data
    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data; self.uploads.append((f, mime)); return "http://b/2"


def test_an_approved_vision_draft_now_certifies_the_file(store, monkeypatch):
    """What the reviewer experiences. Before this, the description was approved, counted as
    content the document owed, and never written — so the file sat approved and unpublishable
    for good."""
    import core, handlers
    store.init_scan_run(SID, "drive", 1, "2026-07-10T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": FILE, "engine": "office", "status": "fail", "score": 40, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "drv1",
        "issues": [{"ruleId": "PPTX-ALT-001", "wcag": "1.1.1", "severity": "CRITICAL"}],
    }, "2026-07-10T00:00:00Z")
    store.record_remediation(SID, FILE, drive_write_url="http://d/1", blob_url="http://b/1")
    item = store.enqueue_proposals(SID, FILE, "1.1.1", [
        {"locator": f"{SLIDE}#rId2", "before": "(no alt text — image skipped by screen readers)",
         "proposed_value": "A dense field of coloured noise.", "rationale": "vision",
         "source": "AI vision model (stub-vision)"},
    ], rule_name="Non-text Content")
    store.update_hitl_item(item, "approved", None, None)
    store.approve_proposal_values(item, [])
    assert store.count_unapplied_approved_values(SID, FILE) == 1     # owed, not yet written

    blob = _Blob(_pkg(SLIDE, _slide(_pic("p:cNvPr", "3", "Picture 3", "rId2"))))
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    from proposals import Verification
    monkeypatch.setattr(handlers, "_verify_residual", lambda b, f: Verification(True, ()))
    handlers._apply_approved_values({"scan_id": SID, "file": FILE}, {})

    assert 'descr="A dense field of coloured noise."' in _read(blob.data, SLIDE)
    assert store.count_unapplied_approved_values(SID, FILE) == 0     # promise kept
    assert store.get_file_record(SID, FILE)["compliant"] == 1
    assert [d["action"] for d in store.list_decisions(scan_id=SID)].count("apply.unresolved") == 0

"""The 3.1.2 pptx language-of-parts lane, proved end to end (WCAG 3.1.2 Language of Parts).

The docx twin of this lane is `tests/test_remediation_verified_docx_language.py`, and the same
five properties are asserted here: the original deck trips the finding, an approval changes the
document, a REAL re-scan verifies it, unrelated content survives, and a broken engine earns no
credit. Nothing but the blob store is patched.

WHY IT IS ITS OWN FILE RATHER THAN A PARAMETER ON THE docx ONE. The criterion, the proposer and
the store getter are shared; the WRITE is not. `apply_text_values.apply_language_parts` marks
Word runs with `<w:lang w:val="…">` and PresentationML runs with `lang="…"` on `<a:rPr>` — two
different attributes in two different namespaces, reached by two different part walks. A
parametrised test would assert the shared half twice and the differing half nowhere, which is
the half that can break.

The whole chain is deterministic here as on docx: `propose_language_parts` is langdetect over
extracted text with no model, so the test starts from the deck and takes the reviewer's value
from the draft ACP actually offers rather than authoring one.

WHAT THIS CLAIMS: the runs carrying the foreign passage now declare `fr`, the deck survived, and
ACP's own criterion stops firing. NOT that the code is the right one for the passage —
langdetect's answer is a probabilistic guess a reviewer approves, which is why the lane is
`assisted` and never `auto`. A wrong-but-approved code would clear the criterion exactly as a
right one does, and no detector can tell them apart.
"""
from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

pytest.importorskip("pptx")
pytest.importorskip("langdetect")

FILE = "policy.pptx"
SID = "rv-pptx-312"
SLIDE = "ppt/slides/slide1.xml"

TITLE = "Accessibility policy"
EN = "The board reviewed the programme and approved the following statement."
FR = ("Les organismes publics doivent rendre leurs documents accessibles à toutes les personnes "
      "en situation de handicap, conformément à la réglementation européenne en vigueur.")
TAIL = "Unrelated closing line that must survive the write."


def _deck(*, mark: str | None = None) -> bytes:
    """An English deck with one French passage. `mark` pre-marks it, for the clean control."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = TITLE
    tf = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4)).text_frame
    tf.word_wrap = True
    tf.text = EN
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = FR
    if mark:
        run.font._rPr.set("lang", mark)
    tf.add_paragraph().text = TAIL

    out = Path(tempfile.mkdtemp()) / FILE
    prs.save(out)
    return out.read_bytes()


def _assess(data: bytes) -> set[str]:
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / FILE).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), FILE, detect_pii=False)
    return {sc for i in (fd or {}).get("issues", []) if (sc := _extract_sc(i.get("wcag", "")))}


def _spill(data: bytes) -> Path:
    p = Path(tempfile.mkdtemp()) / FILE
    p.write_bytes(data)
    return p


def _text(data: bytes) -> str:
    """The deck's extracted text — the words a reader sees, through ACP's own extractor."""
    import pii
    return " ".join((pii.extract_text(_spill(data)) or "").split())


def _slide_xml(data: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(SLIDE).decode("utf-8")


def _proposals(data: bytes) -> list[dict]:
    import pii
    from proposals import propose_language_parts
    return propose_language_parts(pii.extract_text(_spill(data)))


class _Blob:
    def __init__(self, data: bytes):
        self.data, self.uploads = data, []

    def download_remediated(self, owner, sid, f):
        return self.data

    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data
        self.uploads.append((f, mime))
        return "http://b/2"


@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rv.db")
    return store_mod.Store()


def _seed(store, props: list[dict], values: list[str]) -> str:
    store.init_scan_run(SID, "drive", 1, "2026-08-31T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": FILE, "engine": "office", "status": "pass", "score": 60, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "d1",
        "issues": [{"ruleId": "LANG_PARTS_UNMARKED", "wcag": "3.1.2 Language of Parts",
                    "severity": "SERIOUS"}],
    }, "2026-08-31T00:00:00Z")
    store.record_remediation(SID, FILE, drive_write_url="http://d/1", blob_url="http://b/1")
    item_id = store.enqueue_proposals(SID, FILE, "3.1.2", [
        {k: p.get(k) for k in ("locator", "before", "proposed_value", "rationale", "source")}
        for p in props], rule_name="Language of Parts")
    store.update_hitl_item(item_id, "approved", None, None)
    store.approve_proposal_values(item_id, values)
    return item_id


def _run_lane(monkeypatch, store, blob):
    import core
    import handlers
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    handlers._apply_approved_values({"scan_id": SID, "file": FILE}, {})


@pytest.fixture(scope="module")
def unmarked() -> bytes:
    return _deck()


# ── 1. the finding and its proposal ──────────────────────────────────────────

def test_a_real_assessment_reports_3_1_2_on_an_unmarked_passage(unmarked):
    assert "3.1.2" in _assess(unmarked)


def test_the_same_deck_with_the_passage_marked_is_not_flagged():
    """The control, and it also establishes independently of the writer that what clears the
    criterion is the language mark."""
    assert "3.1.2" not in _assess(_deck(mark="fr-FR"))


def test_the_proposer_offers_a_language_code_for_the_passage(unmarked):
    props = _proposals(unmarked)
    assert props, "a 3.1.2 finding with no proposal is a finding a reviewer cannot act on"
    assert props[0]["proposed_value"] == "fr"
    assert props[0]["locator"] and props[0]["locator"] in FR, (
        "the locator is not a prefix of the passage it describes, so the writer cannot find it")


def test_an_english_only_deck_is_not_proposed_for():
    """Self-gating, asserted: the proposer must not invent a foreign span."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = TITLE
    tf = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3)).text_frame
    tf.text = EN
    tf.add_paragraph().text = TAIL
    out = Path(tempfile.mkdtemp()) / FILE
    prs.save(out)
    assert _proposals(out.read_bytes()) == []


# ── 2. approval → write → re-scan → credit, through the real path ────────────

@pytest.fixture()
def applied(store, monkeypatch, unmarked):
    props = _proposals(unmarked)
    blob = _Blob(unmarked)
    _seed(store, props, [p["proposed_value"] for p in props])
    _run_lane(monkeypatch, store, blob)
    return blob, store


def test_the_passage_runs_now_declare_the_language(applied):
    """PresentationML spells this as `lang="fr"` on `<a:rPr>` — a bare attribute in the
    DrawingML namespace, not Word's `<w:lang w:val="…">` element. Asserting the pptx spelling
    is the reason this file is not a parameter on the docx one."""
    blob, _ = applied
    xml = _slide_xml(blob.data)
    assert 'lang="fr"' in xml, "no language mark reached the slide"


def test_not_one_visible_character_changed(applied):
    """Stronger than "unrelated content survives", and available on this criterion alone: a
    language mark is metadata, so the whole extracted text must be identical. Any difference
    means the writer touched prose."""
    blob, _ = applied
    assert _text(blob.data) == _text(_deck())


def test_the_deck_still_opens_and_keeps_its_structure(applied):
    from pptx import Presentation
    blob, _ = applied
    path = _spill(blob.data)
    assert zipfile.ZipFile(path).testzip() is None
    texts = [sh.text_frame.text for sh in Presentation(str(path)).slides[0].shapes
             if sh.has_text_frame]
    joined = "\n".join(texts)
    for fragment in (TITLE, EN, FR, TAIL):
        assert fragment in joined, f"{fragment!r} did not survive the write"


def test_a_second_real_assessment_no_longer_reports_3_1_2(applied):
    blob, _ = applied
    assert "3.1.2" not in _assess(blob.data)


def test_the_row_is_credited_and_the_copy_is_stored(applied):
    blob, store = applied
    assert store.count_unapplied_approved_values(SID, FILE) == 0
    assert blob.uploads


# ── 3. where the lane must NOT credit ────────────────────────────────────────

def test_a_value_that_is_not_a_language_code_is_refused(store, monkeypatch, unmarked):
    """`apply_language_parts` writes into the attribute assistive technology reads as a
    language. A reviewer typing prose there must be reported, never written — a bogus code is
    worse than no mark, because it makes a screen reader switch voice to nothing sensible.

    THE REFUSAL IS ASSERTED ON THE WRITER, not on the stored copy, and that distinction is
    what makes this test able to fail. A lane that credits nothing leaves `blob.data` as the
    deck went IN, so "no lang mark in the stored bytes" and "3.1.2 still fires on them" are
    both true of the ORIGINAL and would hold even if the writer had happily written prose into
    the attribute. Established by bite check: disabling `_LANG_CODE` in apply_text_values left
    the earlier version of this test — and its docx twin — green.
    """
    from apply_text_values import apply_language_parts
    props = _proposals(unmarked)
    locator = props[0]["locator"]

    written, applied, unresolved = apply_language_parts(unmarked, "pptx",
                                                       {locator: "French please"})
    assert applied == [], "prose was written where assistive technology reads a language"
    assert unresolved == [locator], "the bad row was dropped silently instead of reported"
    assert written == unmarked, "the deck was rewritten for a value that was refused"
    assert 'lang="French please"' not in _slide_xml(written)

    blob = _Blob(unmarked)
    _seed(store, props, ["French please"])
    _run_lane(monkeypatch, store, blob)
    assert store.count_unapplied_approved_values(SID, FILE) == 1
    assert not blob.uploads


def test_an_already_marked_deck_is_left_alone(store, monkeypatch):
    marked = _deck(mark="fr-FR")
    assert "3.1.2" not in _assess(marked)
    blob = _Blob(marked)
    _seed(store, [], [])
    _run_lane(monkeypatch, store, blob)
    assert blob.data == marked and not blob.uploads


# ── 4. a broken engine earns nothing ─────────────────────────────────────────

@pytest.mark.parametrize("name,script,timeout", [
    ("cannot be launched", None, None),
    ("exits non-zero", "#!/bin/sh\necho boom >&2\nexit 9\n", None),
    ("hangs past the timeout", "#!/bin/sh\nsleep 30\n", "2"),
])
def test_a_broken_office_analyser_never_credits_this_lane_either(monkeypatch, unmarked,
                                                                 name, script, timeout):
    """Re-asserted per lane rather than assumed to inherit: the fail-open lived in one shared
    seam. 3.1.2 comes from textchecks, which runs after the .NET call in its own try/except, so
    the residual is a real set even with no analyser at all."""
    import stat as _stat

    import scanner
    if script is None:
        monkeypatch.setattr(scanner, "DOTNET", "/nonexistent/dotnet", raising=False)
    else:
        fake = Path(tempfile.mkdtemp()) / "dotnet"
        fake.write_text(script)
        fake.chmod(fake.stat().st_mode | _stat.S_IEXEC)
        monkeypatch.setattr(scanner, "DOTNET", str(fake), raising=False)
    if timeout:
        monkeypatch.setenv("ACP_OFFICE_CLI_TIMEOUT", timeout)

    from proposals import verify_residual_scs
    residual = verify_residual_scs(unmarked, FILE)
    assert residual is not None, (
        f"an office CLI that {name} made the re-scan return None — every approved value on this "
        f"lane would be credited on a scan that never happened")
    assert "3.1.2" in residual

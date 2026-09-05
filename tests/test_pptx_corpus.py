"""The labelled .pptx ground-truth corpus — and the check that its labels are earned.

Third format after .docx and .xlsx, same rule as both: a pair is declared only when a FIRST-PARTY
detector (pure Python in api/office_structure.py, no .NET engine) was driven against the fixture
and confirmed to fire, with an adversarial counterpart confirmed to stay silent. Coverage is
counted from declarations, so an undetected fixture would raise the number #1009 reports without
raising what it measures.

pptx reaches further than xlsx did — nine of seventeen pairs — because it has the most
first-party detectors of any format.

TWO DETECTOR SUBTLETIES THESE FIXTURES ENCODE, both found by the fixture failing to detect
rather than by reading the code:

  * 1.4.3 needs an EXPLICIT shape solid fill as well as an explicit run colour. A bare textbox
    has no fill, so the detector cannot know what the text sits on and correctly says nothing.
    The first draft declared 1.4.3 and detected zero — which is precisely the failure mode these
    tests exist to catch, caught on its author.
  * 1.4.3 on pptx is judged at the WCAG LARGE-text bar (3:1), not 4.5:1: run font size is often
    inherited from the placeholder and not reliably knowable, so flagging only below the
    large-text threshold guarantees every finding is a real failure at any size.

EVERY VIOLATION HAS A PAIRED ADVERSARIAL FIXTURE THAT DIFFERS IN ONE THING. The two shape
fixtures are the same shape at two outline colours; the two link fixtures are the same link with
the underline on and off; the two focus-order decks hold the same two placeholders in opposite
document order. A corpus whose violation and control differ in several ways cannot say which
difference the detector reacted to.
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "api"))
sys.path.insert(0, str(ROOT / "scripts"))

import office_structure as osx  # noqa: E402
import ocr as _ocr  # noqa: E402
import pii as _pii  # noqa: E402
import textchecks as _tc  # noqa: E402
from engines import NO_OFFICE, OFFICE_OK  # noqa: E402

_spec = importlib.util.spec_from_file_location(
    "gen_pptx_corpus", ROOT / "scripts" / "gen_pptx_corpus.py")
gen = importlib.util.module_from_spec(_spec)
sys.modules["gen_pptx_corpus"] = gen
_spec.loader.exec_module(gen)

import corpus_expectations as ce  # noqa: E402


@pytest.fixture(scope="module")
def corpus(tmp_path_factory):
    out = tmp_path_factory.mktemp("pptx-corpus")
    manifest, problems = gen.build_all(out / "docs")
    assert not problems, f"fixtures declare verdicts the engine cannot emit: {problems}"
    return out, {row["name"]: row for row in manifest}


def _ocr_wcags(path: Path, ext: str) -> set[str]:
    """Criteria the OCR lane reports — 1.4.5, read out of the document's PIXELS rather than its
    structure. A real scan runs this pass alongside the structural one, so `_wcags` is their
    union; checking only the structural lane would make every 1.4.5 fixture invisible and every
    "this fixture is single-criterion" assertion below weaker than it reads.

    Empty when tesseract is unavailable — `ocr.is_available()` gates it, and the 1.4.5
    assertions skip rather than fail on a bare checkout (both CI pipelines install tesseract, so
    the skip is a fallback and not the normal state; see test_ocr_is_present_in_ci below)."""
    return {(f.get("wcag") or "").split()[0] for f in _ocr.images_of_text(path, ext)
            if f.get("wcag")}


def _text_wcags(path: Path, ext: str) -> set[str]:
    """Criteria the TEXT lane reports — 1.3.3 and 3.1.2, decided by the document's PROSE rather
    than its structure or its pixels. These two calls are what scanner.py makes (scanner.py:3483):
    extract the text, read the document's own language marks, then judge.

    `language_marked_spans` is passed rather than omitted because 3.1.2 asks whether a foreign
    passage's language is IDENTIFIED, and dropping the marks would make a correctly-marked
    document indistinguishable from an unmarked one — the detector would fire on both and no
    control could ever be clean."""
    text = _pii.extract_text(path) or ""
    return {(f.get("wcag") or "").split()[0]
            for f in _tc.content_findings(text, osx.language_marked_spans(path, ext))
            if f.get("wcag")}


def _wcags(path: Path) -> set[str]:
    """Every criterion a real scan of this file reports, across BOTH lanes: the first-party pptx
    structure checks and the OCR pass over its embedded images. The union is what makes the
    single-criterion assertions below mean anything."""
    structural = {(f.get("wcag") or "").split()[0] for f in osx.checks_for(path, ".pptx")
                  if f.get("wcag")}
    return structural | _ocr_wcags(path, ".pptx") | _text_wcags(path, ".pptx")


# ── the labels are legal for their lane ──────────────────────────────────────────

def test_every_declaration_is_a_verdict_the_engine_can_emit(corpus):
    _out, rows = corpus
    for name, row in rows.items():
        for sc, verdict in row["expect"].items():
            allowed = ce.possible_verdicts(sc, "pptx")
            assert verdict in allowed, (
                f"{name} expects {sc}={verdict}, but ({sc}, pptx) can only emit {sorted(allowed)}")


def test_no_review_lane_fixture_claims_pass(corpus):
    """All nine declared pptx pairs are review-lane — none can certify — so no fixture here may
    expect PASS however clean it is (ADR 0016)."""
    _out, rows = corpus
    for name, row in rows.items():
        for sc, verdict in row["expect"].items():
            if not ce.can_ever_pass(sc, "pptx"):
                assert verdict != "PASS", (
                    f"{name} expects PASS on {sc}, which .pptx cannot certify")


# ── the labels are EARNED ────────────────────────────────────────────────────────


# Criteria that ONLY the OCR lane can report — read out of a page's PIXELS, so unreachable
# without tesseract however correct the fixture is. Kept as an explicit set rather than inferred:
# inferring it would mean running the lane to find out, which is the thing being gated.
_OCR_ONLY_SC = {"1.4.5", "1.4.9"}


@pytest.mark.parametrize("name,sc", [
    ("title-empty", "2.4.6"),
    ("link-vague", "2.4.4"),
    ("link-underline-off", "1.4.1"),
    ("contrast-fail", "1.4.3"),
    ("shape-faint-outline", "1.4.11"),
    ("focus-order", "2.4.3"),
    ("embedded-control", "4.1.2"),
    ("embedded-control", "2.1.2"),
    ("picture-no-alt", "1.1.1"),
    ("image-of-text", "1.4.5"),
    ("sensory-instruction", "1.3.3"),
    ("language-parts", "3.1.2"),
])
def test_each_violation_fixture_is_actually_detected(corpus, name, sc):
    # THE SKIP THE DOCSTRING ALREADY PROMISED. `_ocr_wcags` says the 1.4.5 assertions "skip
    # rather than fail on a bare checkout"; they did not — there was no guard here, so a
    # developer without tesseract got three hard FAILURES from a complete, correct checkout.
    # That is worse than noise: a suite that is red for an environmental reason trains everyone
    # to read red as "probably the usual three", which is exactly how a real regression gets
    # waved through.
    #
    # Skipping loses no coverage, because losing it in CI is caught separately and loudly:
    # test_ocr_is_present_in_ci asserts _ocr.is_available() whenever CI/TF_BUILD is
    # set, so a pipeline that stopped installing tesseract fails there rather than quietly
    # skipping here. Both halves are needed — this one keeps a bare checkout honest, that one
    # keeps CI honest.
    if sc in _OCR_ONLY_SC and not _ocr.is_available():
        pytest.skip(f"{sc} is only reachable through the OCR lane and tesseract is unavailable")
    out, _rows = corpus
    got = _wcags(out / "docs" / f"{name}.pptx")
    assert sc in got, (
        f"{name} seeds a {sc} violation that no first-party detector catches — the manifest "
        f"would claim coverage the corpus does not have. Detected: {sorted(got) or 'nothing'}")


@pytest.mark.parametrize("name,sc", [
    ("title-ok", "2.4.6"),
    ("link-descriptive-ok", "2.4.4"),
    ("link-underlined-ok", "1.4.1"),
    ("contrast-ok", "1.4.3"),
    ("shape-strong-outline-ok", "1.4.11"),
    ("focus-order-ok", "2.4.3"),
    ("no-controls-ok", "4.1.2"),
    ("no-controls-ok", "2.1.2"),
    ("no-picture-ok", "1.1.1"),
    ("image-of-text-logo-ok", "1.4.5"),
    ("sensory-instruction-ok", "1.3.3"),
    ("language-parts-ok", "3.1.2"),
    ("language-parts-marked-ok", "3.1.2"),
])
def test_each_adversarial_fixture_stays_silent(corpus, name, sc):
    """A false positive is cheap to ship and expensive to trust — it is what teaches a reviewer
    to stop reading the findings."""
    out, _rows = corpus
    got = _wcags(out / "docs" / f"{name}.pptx")
    assert sc not in got, (
        f"{name} is a false positive on {sc} — it is the case the detector is supposed to let "
        f"through. Detected: {sorted(got)}")


# ── the pairs differ in exactly one thing ────────────────────────────────────────

def test_the_shape_pair_differs_only_in_outline_colour(corpus):
    """Same shape, two outline colours. If 1.4.11 ever reported on a shape's PRESENCE rather
    than its measured ratio, the adversarial one would fire too."""
    out, _rows = corpus
    assert "1.4.11" in _wcags(out / "docs" / "shape-faint-outline.pptx")
    assert "1.4.11" not in _wcags(out / "docs" / "shape-strong-outline-ok.pptx")


def test_the_link_pair_differs_only_in_the_underline(corpus):
    """Same link text, same destination — one has u="none" on its run. 1.4.1 is about the CUE,
    not the link, and this is what says so."""
    out, _rows = corpus
    assert "1.4.1" in _wcags(out / "docs" / "link-underline-off.pptx")
    assert "1.4.1" not in _wcags(out / "docs" / "link-underlined-ok.pptx")
    # And the descriptive-label pair must not accidentally trip 1.4.1 the other way.
    assert "2.4.4" in _wcags(out / "docs" / "link-vague.pptx")


def test_the_focus_order_pair_differs_only_in_document_order(corpus):
    """The same two placeholders, swapped. 2.4.3 here is 'is the title first', and a detector
    keying off anything else would not distinguish these."""
    out, _rows = corpus
    assert "2.4.3" in _wcags(out / "docs" / "focus-order.pptx")
    assert "2.4.3" not in _wcags(out / "docs" / "focus-order-ok.pptx")


def test_one_control_answers_for_both_its_criteria(corpus):
    """An embedded control is evidence for the accessible-name question AND the keyboard-trap
    question. Asserted together as well as separately, because a change dropping either would
    still pass the other's parametrised case and look fine."""
    out, _rows = corpus
    got = _wcags(out / "docs" / "embedded-control.pptx")
    assert {"4.1.2", "2.1.2"} <= got, f"the control fired only {sorted(got)}"


def test_contrast_needs_an_explicit_fill_and_the_fixture_supplies_one(corpus):
    """Encodes the subtlety that cost the first draft: pptx_contrast_checks needs an explicit
    shape solid fill as well as an explicit run colour, so a bare textbox declares 1.4.3 and
    detects nothing. If the fixture ever loses its fill, this fails rather than silently
    covering zero."""
    out, _rows = corpus
    import zipfile
    with zipfile.ZipFile(out / "docs" / "contrast-fail.pptx") as zf:
        xml = zf.read("ppt/slides/slide1.xml").decode()
    assert "srgbClr val=\"FFFFFF\"" in xml, "the fixture's shape lost its explicit solid fill"
    assert "1.4.3" in _wcags(out / "docs" / "contrast-fail.pptx")


# ── the corpus and the coverage report agree ─────────────────────────────────────

def test_the_coverage_report_counts_this_corpus(corpus):
    import gen_fixture_coverage as gfc
    cov = gfc.coverage()
    declared = set(gen.DECLARED) | set(gen.DECLARED_ENGINE)
    assert cov["pptx"]["has_generator"] is True, (
        "gen_fixture_coverage does not know about the pptx corpus — add it to GENERATORS")
    assert sorted(cov["pptx"]["covered"]) == sorted(declared)
    assert sorted(cov["pptx"]["engine_only"]) == sorted(gen.DECLARED_ENGINE), (
        "the report's engine-only split disagrees with the generator — the headline number "
        "would be counting pairs nobody can confirm on a bare checkout as if they were the same "
        "as the rest")
    assert gfc.BASELINE["pptx"] == len(declared), (
        f"BASELINE['pptx'] is {gfc.BASELINE['pptx']} but the corpus declares "
        f"{len(declared)} — the ratchet would have slack in it")


def test_the_declared_set_matches_what_the_fixtures_actually_declare(corpus):
    _out, rows = corpus
    from_fixtures = {sc for row in rows.values() for sc in row["expect"]}
    assert from_fixtures == set(gen.DECLARED) | set(gen.DECLARED_ENGINE)


def test_every_violation_has_a_paired_adversarial_fixture(corpus):
    """Structural, not per-case: a corpus that grows a violation without its control has stopped
    measuring false positives for that criterion, and nothing else would say so."""
    _out, rows = corpus
    violations = {sc for r in rows.values() if r["kind"] == "violation" for sc in r["expect"]}
    controls = {sc for r in rows.values() if r["kind"] == "adversarial" for sc in r["expect"]}
    assert violations == controls, (
        f"criteria with no adversarial counterpart: {sorted(violations - controls)}; "
        f"with no violation: {sorted(controls - violations)}")


# ── 1.4.5 is read out of the pixels, so it needs its own lane and its own guard ──

def test_ocr_is_present_in_ci():
    """The environment-conditional skip on the 1.4.5 assertions must be a FALLBACK, never the
    normal state. Both ci.yml and azure-pipelines.yml run scripts/install_tesseract.sh, and the
    .docx gate makes the same assertion for the same reason: a skip nobody notices is one edit
    away from being how a criterion stops being covered.

    Skipped OFF CI so a bare checkout is not failed for a dependency it never claimed to have."""
    import os
    if not (os.environ.get("CI") or os.environ.get("TF_BUILD")):
        pytest.skip("not CI — tesseract is optional on a developer checkout")
    assert _ocr.is_available(), (
        "tesseract is unavailable in CI, so 1.4.5 was NOT exercised — the corpus would report "
        "the pair as covered while proving nothing about it")


# ── the engine-verified pairs: structure here, detection in CI ───────────────────
# 2.4.2 and 3.1.1 have no first-party Python detector on ANY Office format, so their labels are
# proven where the .NET analyser is built and skipped where it is not — DECLARED_ENGINE, the same
# split the xlsx corpus introduced. Both are worth the asymmetry: they are among the seventeen
# (criterion, format) pairs in the preset that can return a PASS, so before these fixtures a clean
# scan CERTIFIED a pptx file against two criteria nothing in the suite checked.

def _title_with_text(path: Path) -> str | None:
    """SlideTitleRule's predicate, transcribed: a Title/CenteredTitle placeholder holding
    non-blank text on every slide. Returns the reason it would fire, or None."""
    from pptx import Presentation
    prs = Presentation(str(path))
    for index, slide in enumerate(prs.slides):
        title = next((s for s in slide.shapes
                      if s.is_placeholder
                      and str(s.placeholder_format.type).startswith(("TITLE", "CENTER_TITLE"))),
                     None)
        if title is None:
            return f"slide {index}: no title placeholder"
        text = "".join(r.text or "" for p in title.text_frame.paragraphs for r in p.runs).strip()
        if not text:
            return f"slide {index}: title placeholder is empty"
    return None


def _content_language(path: Path) -> str | None:
    """DocumentLanguageRule's predicate, transcribed. Returns the reason it would fire, or None.

    THE SECOND BRANCH IS THE WHOLE POINT, and it is what makes the pptx rule different from the
    xlsx one of the same name: metadata language OR any lang/altLang on an a:rPr or a:endParaRPr
    in any slide, slide master, or master's layout. The rule's own comment says reading only
    PackageProperties.Language "false-positived essentially every real deck".
    """
    import re
    import zipfile
    from pptx import Presentation
    if (Presentation(str(path)).core_properties.language or "").strip():
        return None
    with zipfile.ZipFile(str(path)) as z:
        for name in z.namelist():
            if not (name.endswith(".xml") and name.startswith(gen._LANG_BEARING_PARTS)):
                continue
            body = z.read(name).decode("utf-8", "replace")
            for element in re.finditer(r"<a:(?:rPr|endParaRPr)\b[^>]*>", body):
                if re.search(r'\b(?:alt)?[Ll]ang="[^"]+"', element.group(0)):
                    return None
    return "no metadata language and no run-level lang anywhere"


def _table_without_header(path: Path) -> str | None:
    """TableHeaderRule's predicate: a table with more than one row whose `firstRow` is not set."""
    from pptx import Presentation
    for index, slide in enumerate(Presentation(str(path)).slides):
        position = 0
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            if len(table.rows) > 1 and not table.first_row:
                return f"slide {index} table {position}: firstRow not set, {len(table.rows)} rows"
            position += 1
    return None


def _reading_order_mismatch(path: Path) -> str | None:
    """ReadingOrderRule's predicate, and the two things about it that are easy to get wrong.

    TAB ORDER IS ASSIGNED BEFORE THE POSITION FILTER. The C# does `.Select((shape, tabOrder) =>
    ...)` and only then `.Where(s => s.HasPos)`, so an unpositioned shape consumes a tab index
    without taking a visual rank. Filtering first would shift every index and change which
    fixtures fire.

    A PLACEHOLDER USUALLY HAS NO POSITION. python-pptx placeholders inherit their geometry from
    the layout and write no `a:xfrm`, so the title and body of this corpus's focus-order decks are
    invisible to this rule entirely — which is why 2.4.3 and 1.3.2 need different fixtures despite
    describing the same-sounding problem.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    for index, slide in enumerate(Presentation(str(path)).slides):
        positioned = []
        for tab, sp in enumerate(slide.shapes._spTree.findall(qn("p:sp"))):
            properties = sp.find(qn("p:spPr"))
            transform = properties.find(qn("a:xfrm")) if properties is not None else None
            offset = transform.find(qn("a:off")) if transform is not None else None
            if offset is None:
                continue
            positioned.append((tab, int(offset.get("y")), int(offset.get("x"))))
        if len(positioned) < 2:
            continue
        for rank, (tab, _y, _x) in enumerate(sorted(positioned, key=lambda r: (r[1], r[2]))):
            if abs(rank - tab) > 1:
                return f"slide {index}: shape at tab {tab} ranks {rank} visually"
    return None


# Every engine predicate, keyed by the criterion it raises. The undeclared-finding sweep walks
# this map rather than a hand-written list, so adding a fifth engine pair extends the sweep by
# construction instead of by remembering to.
ENGINE_PREDICATES = {
    "1.3.1": _table_without_header,
    "1.3.2": _reading_order_mismatch,
    "2.4.2": _title_with_text,
    "3.1.1": _content_language,
}


def test_a_placeholder_carries_no_position_so_the_focus_order_decks_cannot_trip_1_3_2(corpus):
    """WHY 2.4.3 AND 1.3.2 NEED DIFFERENT FIXTURES, measured rather than reasoned.

    The focus-order pair moves the title placeholder to the end of document order, which sounds
    like exactly what ReadingOrderRule looks for. It is invisible to it twice over: placeholders
    write no `a:xfrm` so neither shape is positioned, and even two positioned shapes swapping
    moves each by one rank while the rule needs more than one.

    If python-pptx ever starts writing explicit geometry for placeholders, this fails and the
    focus-order fixtures need re-labelling — which is the notice this test exists to give.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    out, rows = corpus
    deck = Presentation(str(out / rows["focus-order"]["file"]))
    positioned = [
        sp for sp in deck.slides[0].shapes._spTree.findall(qn("p:sp"))
        if (pr := sp.find(qn("p:spPr"))) is not None
        and (xf := pr.find(qn("a:xfrm"))) is not None and xf.find(qn("a:off")) is not None
    ]
    assert len(positioned) < 2, (
        "the focus-order deck now has two or more explicitly-positioned shapes, so "
        "ReadingOrderRule can see it — check whether it raises an undeclared 1.3.2")
    assert _reading_order_mismatch(out / rows["focus-order"]["file"]) is None


@pytest.mark.parametrize("name,fires", [
    ("table-no-header", True),
    ("table-header-ok", False),
])
def test_the_table_fixtures_carry_or_withhold_what_the_rule_reads(corpus, name, fires):
    out, rows = corpus
    reason = _table_without_header(out / rows[name]["file"])
    if fires:
        assert reason, f"{name} is the 1.3.1 fixture and its table now has a header row"
    else:
        assert reason is None, f"{name} should designate a header row and does not"


@pytest.mark.parametrize("name,fires", [
    ("reading-order", True),
    ("reading-order-ok", False),
])
def test_the_reading_order_fixtures_carry_or_withhold_what_the_rule_reads(corpus, name, fires):
    out, rows = corpus
    reason = _reading_order_mismatch(out / rows[name]["file"])
    if fires:
        assert reason, f"{name} is the 1.3.2 fixture and its shapes are now in visual order"
    else:
        assert reason is None, f"{name} should be in visual order and is not: {reason}"


def test_the_reading_order_pair_differs_only_in_document_order(corpus):
    """Same three boxes, same positions, opposite document order. If 1.3.2 ever reported on the
    PRESENCE of several boxes rather than on their order, the control would fire too."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    out, rows = corpus

    def offsets(name):
        deck = Presentation(str(out / rows[name]["file"]))
        return sorted(
            (int(off.get("y")), int(off.get("x")))
            for sp in deck.slides[0].shapes._spTree.findall(qn("p:sp"))
            if (pr := sp.find(qn("p:spPr"))) is not None
            and (xf := pr.find(qn("a:xfrm"))) is not None
            and (off := xf.find(qn("a:off"))) is not None)

    assert offsets("reading-order") == offsets("reading-order-ok")
    assert _reading_order_mismatch(out / rows["reading-order"]["file"])
    assert _reading_order_mismatch(out / rows["reading-order-ok"]["file"]) is None


@pytest.mark.parametrize("name,fires", [
    ("no-slide-title", True),
    ("title-empty", True),
    ("slide-title-ok", False),
    ("title-ok", False),
])
def test_the_title_fixtures_carry_or_withhold_what_the_rule_reads(corpus, name, fires):
    """What this file CAN prove without the .NET analyser, and the half that actually rots.

    Detection is asserted below, gated on the engine. But a fixture silently losing the property
    it was built around is a corpus defect no engine is needed to catch — and it is the likelier
    failure, because it happens whenever someone edits the base deck rather than the fixture.
    Splitting the two means a broken fixture fails everywhere and only the detection claim waits
    for CI.
    """
    out, rows = corpus
    reason = _title_with_text(out / rows[name]["file"])
    if fires:
        assert reason, f"{name} is a 2.4.2 fixture and now has a titled slide"
    else:
        assert reason is None, f"{name} should have a real title and does not: {reason}"


@pytest.mark.parametrize("name,fires", [("no-language", True), ("language-ok", False)])
def test_the_language_fixtures_carry_or_withhold_what_the_rule_reads(corpus, name, fires):
    out, rows = corpus
    reason = _content_language(out / rows[name]["file"])
    if fires:
        assert reason, f"{name} is the 3.1.1 fixture and now declares a language somewhere"
    else:
        assert reason is None, f"{name} should declare a language and does not"


def test_the_xlsx_recipe_would_not_have_worked_here(corpus):
    """THE FINDING THIS PAIR COST, kept as a test because a paragraph cannot fail.

    Two formats have a rule class called DocumentLanguageRule. The xlsx one reads
    `PackageProperties.Language` and nothing else; the pptx one ALSO scans slides, masters and
    layouts for run-level lang. So the xlsx fixture recipe — clear the core property — declares
    the pair on pptx and detects nothing, because python-pptx's default template ships
    `<a:rPr lang="en-US" .../>` in ppt/slideMasters/slideMaster1.xml.

    Built here rather than asserted, so that a python-pptx release which stops shipping those
    attributes turns this into a failure to re-read rather than a comment that has become false.
    """
    from pptx import Presentation
    out, _rows = corpus
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Q3 regional revenue"
    prs.core_properties.language = ""          # the xlsx recipe, verbatim
    transplant = out / "docs" / "_xlsx-recipe-transplanted.pptx"
    prs.save(transplant)

    assert _content_language(transplant) is None, (
        "clearing only the core property now DOES silence the pptx rule — python-pptx has "
        "stopped shipping run-level lang in its template, and gen_pptx_corpus.strip_run_languages "
        "may no longer be needed")
    assert _content_language(out / "docs" / "no-language.pptx"), (
        "the real fixture must still trip the rule the transplant cannot")
    transplant.unlink()


def test_no_fixture_carries_an_undeclared_engine_finding(corpus):
    """THE CORRECTNESS FIX THIS PAIR SURFACED, and the xlsx corpus's hardest-won test.

    An empty title placeholder is the 2.4.6 violation AND, under the analyser, the 2.4.2 one —
    SlideTitleRule flags "present but contains no text" in the same branch as a missing
    placeholder. So `title-empty` was labelled single-criterion and was simply wrong in CI, in the
    way that is hardest to notice: nothing on a bare checkout can raise 2.4.2, so nothing here
    could contradict it.

    This sweeps every fixture against both engine predicates and requires any that would fire to
    have declared it. It is what stops the next fixture reintroducing the same silent mislabel.
    """
    out, rows = corpus
    for name, row in rows.items():
        path = out / row["file"]
        for sc, predicate in ENGINE_PREDICATES.items():
            reason = predicate(path)
            if reason:
                assert sc in row["expect"], (
                    f"{name} would raise {sc} under the analyser ({reason}) and does not declare "
                    f"it — its label is wrong in CI")


def test_the_sweep_covers_every_engine_declared_pair():
    """ANTI-VACUOUS ON THE SWEEP ITSELF. A predicate map that fell behind DECLARED_ENGINE would
    let the next engine pair be added with no undeclared-finding check at all, and the sweep above
    would still pass — silently checking three criteria out of four."""
    assert set(ENGINE_PREDICATES) == set(gen.DECLARED_ENGINE), (
        "ENGINE_PREDICATES and DECLARED_ENGINE disagree; every engine-verified criterion needs a "
        "transcribed predicate or the sweep does not cover it")


@pytest.mark.skipif(not OFFICE_OK, reason=NO_OFFICE)
@pytest.mark.parametrize("name,sc,fires", [
    ("no-slide-title", "2.4.2", True),
    ("slide-title-ok", "2.4.2", False),
    ("no-language", "3.1.1", True),
    ("language-ok", "3.1.1", False),
    ("table-no-header", "1.3.1", True),
    ("table-header-ok", "1.3.1", False),
    ("reading-order", "1.3.2", True),
    ("reading-order-ok", "1.3.2", False),
])
def test_the_engine_confirms_the_declared_pairs(corpus, name, sc, fires):
    """The detection half, and the reason these two sit in DECLARED_ENGINE rather than DECLARED.

    SC ids come through `assessment_policy._extract_sc`, not a string split: the .NET analyser
    reports `wcag` in enum form ("SC_2_4_2") where the first-party checks report "2.4.2 Page
    Titled", so splitting on whitespace yields "SC_2_4_2" and matches nothing. That is what the
    first CI run of the xlsx version of this test failed on — both detectors HAD fired and the
    assertion could not see it.
    """
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    out, rows = corpus
    path = out / rows[name]["file"]
    fd, _ = analyse_and_assess(path.parent, path.name, detect_pii=False)
    # NOT a walrus in a comprehension. PEP 572 binds an assignment expression in the CONTAINING
    # scope, so `{sc for i in ... if (sc := ...)}` rebinds the parametrised `sc` to whichever
    # criterion the last issue happened to carry. Every assertion below then judged the wrong
    # criterion — and in the `fires=True` direction it asserted that a value just extracted from
    # `found` was in `found`, which is true by construction. Those rows passed vacuously in CI
    # until a `fires=False` row happened to leak a criterion that WAS present and failed with the
    # give-away message "table-header-ok is the clean control for 1.1.1".
    found = {s for i in (fd or {}).get("issues", []) if (s := _extract_sc(i.get("wcag", "")))}
    if fires:
        assert sc in found, (
            f"{name} declares {sc} but the analyser reported {sorted(found) or 'nothing'}")
    else:
        assert sc not in found, f"{name} is the clean control for {sc} but the analyser flagged it"

    # The other engine pair must stay quiet whichever way this fixture goes. `no-slide-title`
    # withholds a title but stamps a language; `no-language` keeps its title. Before the base deck
    # stamped a language explicitly, this held only because python-pptx's template happens to
    # carry run-level lang — which is an accident, not a guarantee.
    for other in set(gen.DECLARED_ENGINE) - {sc}:
        if other in rows[name]["expect"]:
            continue
        assert other not in found, (
            f"{name} also raised {other}, which it does not declare — the base deck has stopped "
            f"supplying something every fixture relied on, so the labels across this corpus are "
            f"now wrong in CI")


def _reading_order_gap(path: Path) -> int:
    """The largest |visualRank - tabOrder| on the first slide — the quantity the rule thresholds.

    Returned as a number rather than a boolean because the MARGIN is the thing worth pinning: a
    fixture that fires at exactly the threshold, or a control that is quiet at exactly it, is one
    unrelated edit away from flipping.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    slide = Presentation(str(path)).slides[0]
    positioned = []
    for tab, sp in enumerate(slide.shapes._spTree.findall(qn("p:sp"))):
        properties = sp.find(qn("p:spPr"))
        transform = properties.find(qn("a:xfrm")) if properties is not None else None
        offset = transform.find(qn("a:off")) if transform is not None else None
        if offset is not None:
            positioned.append((tab, int(offset.get("y")), int(offset.get("x"))))
    if len(positioned) < 2:
        return 0
    return max(abs(rank - tab) for rank, (tab, _y, _x)
               in enumerate(sorted(positioned, key=lambda r: (r[1], r[2]))))


def test_the_reading_order_control_has_margin_not_luck(corpus):
    """THE FRAGILITY A BITE CHECK FOUND, pinned so it cannot come back.

    ReadingOrderRule assigns tab order over every `p:sp` and only then discards the ones with no
    `a:off`, so an unpositioned placeholder consumes an index without taking a rank. python-pptx
    placeholders inherit their geometry and write no `a:xfrm` — so with the title unpositioned, a
    PERFECTLY ORDERED deck already sat at a gap of 1, which is the rule's tolerance exactly. One
    more unpositioned shape ahead of the boxes and the clean control would have fired, and the
    false positive would have been read as a detector bug.

    `_position_title` gives the title an explicit `a:xfrm`, which puts the control at 0. This
    asserts the margin rather than the outcome: a control that is quiet AT the threshold passes an
    "is it quiet?" test and is still one edit from flipping.
    """
    out, rows = corpus
    assert _reading_order_gap(out / rows["reading-order-ok"]["file"]) == 0, (
        "the ordered control no longer sits at a gap of zero — it is drifting back toward the "
        "rule's threshold, where an unrelated edit flips it into a false positive")
    assert _reading_order_gap(out / rows["reading-order"]["file"]) > 1, (
        "the violation no longer exceeds the rule's threshold")


def _graphic_frame_without_alt_text(path: Path) -> str | None:
    """AltTextRule's graphic-frame branch — the one a table falls into.

    The rule walks `Descendants<GraphicFrame>()` as well as `Descendants<Picture>()`, so a table
    with no `descr` raises 1.1.1. The first-party `office_non_text_content_checks` reads pictures
    only, which is why a bare checkout cannot see this and CI can.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    for index, slide in enumerate(Presentation(str(path)).slides):
        for frame in slide.shapes._spTree.iter(qn("p:graphicFrame")):
            properties = frame.find(qn("p:nvGraphicFramePr") + "/" + qn("p:cNvPr"))
            if properties is None or not (properties.get("descr") or "").strip():
                return f"slide {index}: graphic frame with no descr"
    return None


def test_no_fixture_carries_a_graphic_frame_without_alt_text(corpus):
    """THE FAILURE CI FOUND AND A BARE CHECKOUT STRUCTURALLY CANNOT, swept across the corpus.

    The 1.3.1 table fixtures raised an undeclared 1.1.1 because a table is a GraphicFrame and
    AltTextRule reads those. Every fixture that declares 1.1.1 may carry the finding; nothing else
    may. This is the picture-only sibling of the engine sweep above, kept separate because 1.1.1
    is a first-party DECLARED pair whose first-party detector disagrees with the analyser about
    what counts as non-text content.
    """
    out, rows = corpus
    for name, row in rows.items():
        reason = _graphic_frame_without_alt_text(out / row["file"])
        if reason and "1.1.1" not in row["expect"]:
            raise AssertionError(
                f"{name} has a graphic frame with no alt text ({reason}), so the analyser raises "
                f"1.1.1 — which it does not declare. Give the frame a descr, or declare 1.1.1.")

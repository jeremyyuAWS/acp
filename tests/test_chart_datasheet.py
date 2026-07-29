"""Native-chart datasheet proposer (1.1.1) — grounded alt text from the chart's own data.

A native (DrawingML) chart stores its series data in chart XML, so its alt text can be
EXACT rather than a vision guess at a rendered picture: title, series names, categories and
values are read straight from the file. Every number is the chart's own — the ADR-0016 gold
standard. Files with no native chart (embedded chart images) yield nothing here (handled by
the vision/OCR proposers).
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))

import proposals  # noqa: E402

pptx = pytest.importorskip("pptx")


def _chart_pptx(tmp: Path) -> Path:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Revenue ($M)", (1.2, 1.5, 1.8, 2.1))
    cd.add_series("Expenses ($M)", (0.9, 1.0, 1.1, 1.3))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
                            Inches(6), Inches(4), cd)
    gf.chart.has_title = True
    gf.chart.chart_title.text_frame.text = "Quarterly Financials 2026"
    p = tmp / "deck.pptx"
    prs.save(str(p))
    return p


def test_chart_alt_and_datasheet_are_exact():
    with tempfile.TemporaryDirectory() as d:
        out = proposals.propose_chart_datasheet(_chart_pptx(Path(d)), ".pptx")
    assert len(out) == 1
    p = out[0]
    # Concise alt names the chart and its range, from the chart's OWN title/series/categories.
    assert "Quarterly Financials 2026" in p["proposed_value"]
    assert "Revenue ($M)" in p["proposed_value"] and "Expenses ($M)" in p["proposed_value"]
    assert "Q1–Q4" in p["proposed_value"]
    # The datasheet carries every exact value paired with its category — nothing invented.
    sheet = p["rationale"]
    assert "Q1 1.2" in sheet and "Q4 2.1" in sheet          # revenue row, trimmed numbers
    assert "Q2 1" in sheet                                   # 1.0 → '1'
    assert "deterministic" in p["source"] and "chart XML" in p["source"]


def test_num_trims_cleanly():
    assert proposals._num("1.0") == "1"
    assert proposals._num("1.50") == "1.5"
    assert proposals._num("2.1") == "2.1"
    assert proposals._num("Region") == "Region"             # non-numeric untouched


def test_openpyxl_xlsx_chart_yields_an_exact_datasheet(tmp_path):
    """An openpyxl-authored workbook caches NO chart values (only a cell range) and writes the
    chart part in the default chart namespace — the `c:`-prefixed regex parser sees nothing there.
    The chart_data fallback resolves the cells, so the workbook still gets a grounded proposal."""
    opx = pytest.importorskip("openpyxl")
    from openpyxl.chart import BarChart, Reference
    wb = opx.Workbook(); ws = wb.active
    ws.append(["Region", "Revenue"])
    for c, v in zip(["North", "South", "East", "West"], [120, 70, 150, 50]):
        ws.append([c, v])
    ch = BarChart(); ch.title = "Sales by region"
    ch.add_data(Reference(ws, min_col=2, min_row=1, max_row=5), titles_from_data=True)
    ch.set_categories(Reference(ws, min_col=1, min_row=2, max_row=5))
    ws.add_chart(ch, "E2")
    p = tmp_path / "book.xlsx"
    wb.save(str(p))

    out = proposals.propose_chart_datasheet(p, ".xlsx")
    assert len(out) == 1, "a cell-referenced native chart must still produce a datasheet proposal"
    prop = out[0]
    assert "Sales by region" in prop["proposed_value"]
    assert "North–West" in prop["proposed_value"]                 # real category span
    # Every figure is the chart's own, read from the worksheet cells it points at.
    assert "North 120" in prop["rationale"] and "East 150" in prop["rationale"]
    assert "deterministic" in prop["source"]


def test_no_native_chart_yields_nothing(tmp_path):
    import zipfile
    p = tmp_path / "plain.pptx"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("ppt/slides/slide1.xml", "<p:sld/>")
    assert proposals.propose_chart_datasheet(p, ".pptx") == []


def test_wiring():
    src = (Path(__file__).resolve().parent.parent / "api" / "handlers.py").read_text()
    assert "propose_chart_datasheet" in src
    # chart datasheets ride the same 1.1.1 card as the per-image alt drafts (one enqueue).
    assert "(chart_sheets or []) + (img_props or [])" in src

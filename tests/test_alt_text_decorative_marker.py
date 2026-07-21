"""Regression guard: the vendored .NET AltTextRule (docx/pptx/xlsx) must not flag an
image that is correctly marked decorative (empty descr + the OOXML `adec:decorative`
extLst marker some authoring tools write for "Mark as decorative").

Found while cross-referencing the detectors against a customer checklist: the
remediation side (api/remediate_office.py's `_inject_descr`) already skips writing
alt text into decorative-marked images, but AltTextRule.cs in all three formats never
checked for the marker at all — so a correctly-authored decorative image was reported
as "missing alt text" on every scan. Needs the built .NET CLI to run, so it self-skips
when the analysers didn't run (CLI not built).
"""
import io
import re
import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))
import scanner  # noqa: E402

_DECORATIVE_EXT = (
    '<a:extLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">'
    '<adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative" val="1"/>'
    "</a:ext></a:extLst>"
)


def _mark_decorative(xml: str, tag: str, decorative_id: str) -> str:
    """Strip any descr= from every self-closing `<tag id="…" …/>` (so every candidate
    starts from a true empty-alt baseline), then inject the decorative extLst into the
    one whose id matches `decorative_id`, converting it to an open/close pair."""
    def repl(m):
        attrs = re.sub(r'\s*descr="[^"]*"', "", m.group(1))
        if re.search(rf'\bid="{re.escape(decorative_id)}"', attrs):
            return f"<{tag}{attrs}>{_DECORATIVE_EXT}</{tag}>"
        return f"<{tag}{attrs}/>"
    return re.sub(rf"<{tag}\b([^>]*?)/>", repl, xml)


def _scs(fdict) -> list[str]:
    out = []
    for i in fdict.get("issues", []):
        w = str(i.get("wcag", "")).strip()
        out.append(w[3:].replace("_", ".") if w.startswith("SC_") else w.split(" ", 1)[0])
    return out


def _engine_ran(fdict) -> bool:
    return any(str(i.get("wcag", "")).startswith("SC_") for i in fdict.get("issues", []))


def _png_bytes() -> bytes:
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", (50, 50), (10, 20, 30)).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture(scope="module")
def docx_scanned(tmp_path_factory):
    from docx import Document

    d = tmp_path_factory.mktemp("decorative-docx")
    doc = Document()
    doc.add_picture(io.BytesIO(_png_bytes()), width=None)
    doc.add_picture(io.BytesIO(_png_bytes()), width=None)
    name = "decorative.docx"
    path = d / name
    doc.save(str(path))

    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml").decode("utf-8")
        others = {n: z.read(n) for n in z.namelist() if n != "word/document.xml"}
    xml = _mark_decorative(xml, "wp:docPr", "1")   # image 1 → decorative, image 2 → plain
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in others.items():
            z.writestr(n, data)
        z.writestr("word/document.xml", xml)

    fdict, _ = scanner.analyse_and_assess(d, name, detect_pii=False)
    return {"scs": _scs(fdict), "engine_ran": _engine_ran(fdict)}


@pytest.fixture(scope="module")
def pptx_scanned(tmp_path_factory):
    from pptx import Presentation
    from pptx.util import Inches

    d = tmp_path_factory.mktemp("decorative-pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(0), Inches(0), Inches(1), Inches(1))
    slide.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(2), Inches(2), Inches(1), Inches(1))
    name = "decorative.pptx"
    path = d / name
    prs.save(str(path))

    with zipfile.ZipFile(path) as z:
        slide_name = next(n for n in z.namelist() if n.startswith("ppt/slides/slide"))
        xml = z.read(slide_name).decode("utf-8")
        others = {n: z.read(n) for n in z.namelist() if n != slide_name}
    xml = _mark_decorative(xml, "p:cNvPr", "2")    # id 1 is the slide's own group shape
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in others.items():
            z.writestr(n, data)
        z.writestr(slide_name, xml)

    fdict, _ = scanner.analyse_and_assess(d, name, detect_pii=False)
    return {"scs": _scs(fdict), "engine_ran": _engine_ran(fdict)}


@pytest.fixture(scope="module")
def xlsx_scanned(tmp_path_factory):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    d = tmp_path_factory.mktemp("decorative-xlsx")
    wb = Workbook()
    ws = wb.active
    ws.add_image(XLImage(io.BytesIO(_png_bytes())), "B2")
    ws.add_image(XLImage(io.BytesIO(_png_bytes())), "D4")
    name = "decorative.xlsx"
    path = d / name
    wb.save(str(path))

    with zipfile.ZipFile(path) as z:
        drawing_name = next(n for n in z.namelist() if re.match(r"xl/drawings/drawing\d+\.xml$", n))
        xml = z.read(drawing_name).decode("utf-8")
        others = {n: z.read(n) for n in z.namelist() if n != drawing_name}
    xml = _mark_decorative(xml, "cNvPr", "1")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in others.items():
            z.writestr(n, data)
        z.writestr(drawing_name, xml)

    fdict, _ = scanner.analyse_and_assess(d, name, detect_pii=False)
    return {"scs": _scs(fdict), "engine_ran": _engine_ran(fdict)}


def test_docx_decorative_image_not_flagged_but_plain_empty_alt_still_is(docx_scanned):
    if not docx_scanned["engine_ran"]:
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    assert docx_scanned["scs"].count("1.1.1") == 1, (
        "expected exactly one 1.1.1 finding (the non-decorative image) — "
        "either the decorative image was wrongly flagged, or the plain empty-alt "
        "image was wrongly skipped"
    )


def test_pptx_decorative_image_not_flagged_but_plain_empty_alt_still_is(pptx_scanned):
    if not pptx_scanned["engine_ran"]:
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    assert pptx_scanned["scs"].count("1.1.1") == 1, (
        "expected exactly one 1.1.1 finding (the non-decorative image) — "
        "either the decorative image was wrongly flagged, or the plain empty-alt "
        "image was wrongly skipped"
    )


def test_xlsx_decorative_image_not_flagged_but_plain_empty_alt_still_is(xlsx_scanned):
    if not xlsx_scanned["engine_ran"]:
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    assert xlsx_scanned["scs"].count("1.1.1") == 1, (
        "expected exactly one 1.1.1 finding (the non-decorative image) — "
        "either the decorative image was wrongly flagged, or the plain empty-alt "
        "image was wrongly skipped"
    )

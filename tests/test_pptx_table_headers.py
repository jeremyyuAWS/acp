"""Deterministic WCAG 1.3.1 remediation for PowerPoint tables — the WRITE.

Detection lives in the vendored .NET analyser (PPTX-TABLE-001, pinned by
tests/test_pptx_engine_detection.py). This module tests the fix that clears it:
`remediate_office._pptx_mark_table_headers` sets `<a:tblPr firstRow="1">` on every
multi-row DrawingML table lacking a marked header row — the direct analogue of the docx
`<w:tblHeader>` and xlsx headerRowCount writes. It is string-surgical (only the tblPr tag
is touched) and fails closed on degenerate/single-row tables.
"""
import re
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))
from remediate_office import _pptx_mark_table_headers, _remediate_pptx_slides  # noqa: E402


def _table(rows: int, *, tblpr: str | None = 'firstRow="0" bandRow="1"') -> str:
    """A DrawingML table with `rows` rows, wrapped so there is surrounding markup to
    prove byte-preservation against. `tblpr` is the attribute string of <a:tblPr>, or
    None to omit the element entirely."""
    pr = "" if tblpr is None else f'<a:tblPr {tblpr}><a:tableStyleId>{{5C22544A}}</a:tableStyleId></a:tblPr>'
    tr = "".join(f'<a:tr h="370840"><a:tc><a:txBody><a:p><a:r><a:t>r{i}</a:t></a:r></a:p>'
                 f'</a:txBody></a:tc></a:tr>' for i in range(rows))
    return (f'<a:tbl>{pr}<a:tblGrid><a:gridCol w="3048000"/></a:tblGrid>{tr}</a:tbl>')


def _slide(table_xml: str) -> str:
    return (f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<p:sld xmlns:p="ns-p" xmlns:a="ns-a"><p:cSld><p:spTree>'
            f'<p:graphicFrame><a:graphic><a:graphicData>{table_xml}'
            f'</a:graphicData></a:graphic></p:graphicFrame></p:spTree></p:cSld></p:sld>')


def test_multirow_table_gets_first_row():
    xml = _slide(_table(3))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 1
    assert '<a:tblPr firstRow="1" bandRow="1">' in out
    assert 'firstRow="0"' not in out


def test_missing_tblpr_is_inserted_as_first_child():
    xml = _slide(_table(2, tblpr=None))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 1
    # inserted as the tbl's first child, before <a:tblGrid>
    assert '<a:tbl><a:tblPr firstRow="1"/><a:tblGrid>' in out


def test_single_row_table_is_skipped():
    # Mutation-check the guard: a 1-row (degenerate/layout) table must NOT be marked.
    # Drop the `len(_A_TR.findall(block)) <= 1` guard and this flips firstRow -> "1".
    xml = _slide(_table(1))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 0
    assert out == xml
    assert 'firstRow="0"' in out


def test_zero_row_table_is_skipped():
    xml = _slide(_table(0))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 0
    assert out == xml


def test_already_marked_header_is_not_recounted():
    xml = _slide(_table(3, tblpr='firstRow="1" bandRow="1"'))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 0
    assert out == xml


def test_firstRow_true_is_treated_as_already_marked():
    xml = _slide(_table(3, tblpr='firstRow="true"'))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 0
    assert out == xml


def test_only_the_tblpr_bytes_change_on_flip():
    xml = _slide(_table(3))
    out, _ = _pptx_mark_table_headers(xml)
    # Every byte outside the flipped attribute value is preserved.
    assert out.replace('firstRow="1"', 'firstRow="0"', 1) == xml


def test_only_the_tblpr_bytes_change_on_insert():
    xml = _slide(_table(2, tblpr=None))
    out, _ = _pptx_mark_table_headers(xml)
    # The only new bytes are the inserted tblPr element.
    assert out.replace('<a:tblPr firstRow="1"/>', '', 1) == xml


def test_multiple_tables_each_marked_once():
    xml = _slide(_table(3) + _table(4))
    out, n = _pptx_mark_table_headers(xml)
    assert n == 2
    assert out.count('firstRow="1"') == 2
    assert 'firstRow="0"' not in out


def test_records_a_diff_per_table():
    diffs: list[dict] = []
    _pptx_mark_table_headers(_slide(_table(3)), diffs)
    assert [d["rule_id"] for d in diffs] == ["1.3.1"]
    assert 'firstRow="1"' in diffs[0]["after"]


def test_integration_through_remediate_pptx_slides(tmp_path):
    """End-to-end via the pptx slide remediator on a real python-pptx file: a headerless
    3-row data table comes out with firstRow set, and the applied list names 1.3.1."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    gf = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    gf.table.first_row = False                      # headerless -> <a:tblPr firstRow="0">
    for ri in range(3):
        for ci in range(2):
            gf.table.cell(ri, ci).text = f"r{ri}c{ci}"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    with zipfile.ZipFile(p) as z:
        entries = {n: z.read(n) for n in z.namelist()}
    applied = _remediate_pptx_slides(entries)

    slide_xml = next(v.decode("utf-8") for k, v in entries.items()
                     if re.fullmatch(r"ppt/slides/slide\d+\.xml", k))
    assert 'firstRow="1"' in slide_xml
    assert 'firstRow="0"' not in slide_xml
    assert any("1.3.1" in a for a in applied)

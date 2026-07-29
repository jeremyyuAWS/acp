"""Native-chart alt text (chart_data.py): the exact values, read from the file — never guessed.

Builds a real native PowerPoint chart with KNOWN data via python-pptx, then asserts the extracted
alt states those exact values (the whole point: no llava confabulation for native charts).
"""
import io
import sys
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))
import chart_data  # noqa: E402
from chart_fixtures import docx_chart_entries, rezip, strip_caches  # noqa: E402

pptx = pytest.importorskip("pptx")


def _native_chart_pptx(cats, vals, title="Q4 Revenue by Region") -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Revenue", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(1), Inches(1), Inches(6), Inches(4), data)
    gf.chart.has_title = True
    gf.chart.chart_title.text_frame.text = title
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


def test_reads_the_real_values_high_and_low():
    data = _native_chart_pptx(["North", "South", "East", "West"], [120, 70, 150, 50])
    charts = chart_data.charts_in(data, ".pptx")
    assert len(charts) == 1
    alt = chart_data.describe_chart(charts[0])
    # The exact truth: East is highest at 150, West lowest at 50 — stated because it was READ.
    assert "East" in alt and "150" in alt
    assert "West" in alt and "50" in alt
    assert "highest" in alt.lower() and "lowest" in alt.lower()
    assert "Q4 Revenue by Region" in alt


def test_chart_alts_convenience_and_type():
    data = _native_chart_pptx(["A", "B", "C"], [10, 30, 20])
    alts = chart_data.chart_alts(data, ".pptx")
    assert len(alts) == 1
    assert alts[0].lower().startswith(("bar chart", "column", "chart"))
    assert "B" in alts[0] and "30" in alts[0]      # the real max


def test_non_office_and_garbage_return_empty_never_raise():
    assert chart_data.charts_in(b"", ".pptx") == []
    assert chart_data.charts_in(b"not a zip", ".pptx") == []
    assert chart_data.charts_in(b"%PDF-1.7", ".pdf") == []
    assert chart_data.chart_alts(b"junk", ".docx") == []


def test_slide_chart_descr_injects_accurate_alt_onto_the_chart_shape():
    import io
    import zipfile
    data = _native_chart_pptx(["North", "South", "East", "West"], [120, 70, 150, 50])
    entries = {n: zipfile.ZipFile(io.BytesIO(data)).read(n) for n in zipfile.ZipFile(io.BytesIO(data)).namelist()}
    changed = chart_data.slide_chart_descr(entries)
    assert changed, "a native chart with no alt should get one"
    (new_xml, alts) = next(iter(changed.values()))
    xml = new_xml.decode("utf-8")
    # The descr landed on a <p:cNvPr> and states the REAL high/low — not a guess.
    assert 'descr="Bar chart' in xml and "East at 150" in xml and "West at 50" in xml
    assert alts and "East at 150" in alts[0]


def _native_chart_xlsx(cats, vals, title="Q4 Revenue by Region") -> bytes:
    opx = pytest.importorskip("openpyxl")
    from openpyxl.chart import BarChart, Reference
    wb = opx.Workbook(); ws = wb.active
    ws.append(["Region", "Revenue"])
    for c, v in zip(cats, vals):
        ws.append([c, v])
    ch = BarChart(); ch.title = title
    data = Reference(ws, min_col=2, min_row=1, max_row=1 + len(cats))
    catref = Reference(ws, min_col=1, min_row=2, max_row=1 + len(cats))
    ch.add_data(data, titles_from_data=True); ch.set_categories(catref)
    ws.add_chart(ch, "E2")
    buf = io.BytesIO(); wb.save(buf); return buf.getvalue()


def test_xlsx_chart_reads_the_real_cell_values():
    # A native xlsx chart references cells (empty numCache) — the resolver reads the actual VALUES
    # from the worksheet, and the category LABELS too, whether the sheet stores them as shared
    # strings (Excel) or inline strings (openpyxl writes `t="inlineStr"` with <is><t>).
    data = _native_chart_xlsx(["North", "South", "East", "West"], [120, 70, 150, 50])
    alts = chart_data.chart_alts(data, ".xlsx")
    assert alts and alts[0].lower().startswith("bar chart")
    assert "150" in alts[0] and "50" in alts[0]                 # real high/low, read from cells
    assert "East at 150" in alts[0] and "West at 50" in alts[0]  # real labels, not "item 3"/"item 4"


def test_slide_chart_descr_is_idempotent_never_overwrites():
    import io
    import zipfile
    data = _native_chart_pptx(["A", "B"], [1, 2])
    z = zipfile.ZipFile(io.BytesIO(data))
    entries = {n: z.read(n) for n in z.namelist()}
    first = chart_data.slide_chart_descr(entries)
    assert first                                             # first pass adds the alt
    for name, (new_xml, _) in first.items():
        entries[name] = new_xml
    # Second pass on the now-described chart must be a no-op — a real descr is never overwritten.
    assert chart_data.slide_chart_descr(entries) == {}


# ── xlsx drawings: BOTH namespace flavours must get the descr ─────────────────────────────────────
_SD_NS = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"


def _xlsx_entries(cats, vals, title="Sales by region") -> dict:
    import zipfile
    z = zipfile.ZipFile(io.BytesIO(_native_chart_xlsx(cats, vals, title)))
    return {n: z.read(n) for n in z.namelist()}


def _to_xdr_prefixed(xml: str) -> str:
    """Rewrite a default-namespace spreadsheetDrawing into the `xdr:`-prefixed flavour Excel
    itself writes — the SAME parts, the other legal namespace spelling. (Every unprefixed tag in
    that part is spreadsheetDrawing-namespaced; `a:`/`c:` ones already carry their own prefix.)"""
    import re as _re
    xml = xml.replace(f'xmlns="{_SD_NS}"', f'xmlns:xdr="{_SD_NS}"')
    return _re.sub(r"<(/?)([A-Za-z][\w-]*)(?=[\s/>])", r"<\1xdr:\2", xml)


def test_xlsx_default_namespace_drawing_gets_a_descr():
    # openpyxl-family generators emit <graphicFrame>/<cNvPr> in the DEFAULT spreadsheetDrawing
    # namespace. Matching only the xdr:-prefixed flavour left these charts with NO alt at all.
    entries = _xlsx_entries(["North", "South", "East", "West"], [120, 70, 150, 50])
    drawing = entries["xl/drawings/drawing1.xml"].decode()
    assert "<cNvPr " in drawing and "xdr:" not in drawing        # the fixture really is default-ns

    changed = chart_data.chart_descr_edits(entries, ".xlsx")
    assert changed, "a default-namespace xlsx chart must still get its deterministic 1.1.1 alt"
    new_xml, alts = changed["xl/drawings/drawing1.xml"]
    out = new_xml.decode()
    assert 'descr="Bar chart' in out and "East at 150" in out and "West at 50" in out
    assert alts and "Sales by region" in alts[0]
    # The emitted tag reuses the name that was MATCHED — no invented xdr: prefix, still well-formed.
    assert "<cNvPr " in out and "xdr:cNvPr" not in out
    ET.fromstring(out)


def test_xlsx_xdr_prefixed_drawing_still_gets_a_descr():
    entries = _xlsx_entries(["North", "South", "East", "West"], [120, 70, 150, 50])
    name = "xl/drawings/drawing1.xml"
    entries[name] = _to_xdr_prefixed(entries[name].decode()).encode()
    assert b"<xdr:cNvPr " in entries[name]                       # the fixture really is prefixed

    changed = chart_data.chart_descr_edits(entries, ".xlsx")
    assert changed, "the Excel-authored (xdr:-prefixed) flavour must keep working"
    out = changed[name][0].decode()
    assert 'descr="Bar chart' in out and "East at 150" in out
    # ...and here the prefix is preserved, because the replacement re-emits the matched tag name.
    assert "<xdr:cNvPr " in out and out.count("<cNvPr ") == 0
    ET.fromstring(out)


@pytest.mark.parametrize("prefixed", [False, True])
def test_xlsx_chart_descr_is_idempotent_in_both_flavours(prefixed):
    entries = _xlsx_entries(["A", "B"], [1, 2])
    name = "xl/drawings/drawing1.xml"
    if prefixed:
        entries[name] = _to_xdr_prefixed(entries[name].decode()).encode()
    first = chart_data.chart_descr_edits(entries, ".xlsx")
    assert first
    for part, (new_xml, _) in first.items():
        entries[part] = new_xml
    # A real descr is never overwritten — second pass is a no-op, same as the pptx path.
    assert chart_data.chart_descr_edits(entries, ".xlsx") == {}


# ── docx/pptx: values that live in the EMBEDDED workbook, not in the chart's caches ────────────────
# A docx/pptx chart keeps caches of its numbers and the authoritative copy in an xlsx zipped inside
# the package. Word and PowerPoint populate the caches, so the caches are normally enough. A
# generator that leaves them empty puts the values ONLY in that nested workbook — and the cell
# resolver cannot help, because it looks for `xl/…` parts that a docx/pptx package does not have.

def test_pptx_chart_with_empty_caches_reads_the_embedded_workbook():
    """The realistic case, on a genuine PowerPoint-shaped package: python-pptx writes both the
    caches and ppt/embeddings/Microsoft_Excel_Sheet1.xlsx. Strip the caches and the embedded
    workbook is the only surviving copy of the data."""
    import zipfile
    data = _native_chart_pptx(["North", "South", "East", "West"], [120, 70, 150, 50])
    z = zipfile.ZipFile(io.BytesIO(data))
    entries = {n: z.read(n) for n in z.namelist()}
    assert "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx" in entries    # the fixture really embeds one
    entries["ppt/charts/chart1.xml"] = strip_caches(entries["ppt/charts/chart1.xml"])

    charts = chart_data.charts_in(rezip(entries), ".pptx")
    assert charts, "a cache-less pptx chart must fall back to its embedded workbook"
    alt = chart_data.describe_chart(charts[0])
    assert "East at 150" in alt and "West at 50" in alt              # real labels AND real values


def test_pptx_chart_with_empty_caches_and_no_embedding_yields_nothing():
    """No caches and no embedded workbook: there is nowhere left to read the numbers from, so the
    chart is skipped. It must not raise, and must not invent an alt over `item N` placeholders."""
    import zipfile
    data = _native_chart_pptx(["North", "South"], [1, 2])
    z = zipfile.ZipFile(io.BytesIO(data))
    entries = {n: z.read(n) for n in z.namelist() if "embeddings" not in n}
    entries["ppt/charts/chart1.xml"] = strip_caches(entries["ppt/charts/chart1.xml"])
    assert chart_data.charts_in(rezip(entries), ".pptx") == []


def test_docx_chart_reads_values_from_its_embedded_workbook():
    entries = docx_chart_entries(["North", "South", "East", "West"], [120, 70, 150, 50])
    assert b"Cache/>" in entries["word/charts/chart1.xml"]    # the fixture really caches nothing

    charts = chart_data.charts_in(rezip(entries), ".docx")
    assert charts, "a cache-less docx chart must resolve through word/embeddings/*.xlsx"
    alt = chart_data.describe_chart(charts[0])
    assert "Sales by region" in alt
    # Real labels, not "item 1"/"item 3" — the inlineStr branch, reached via the embedded workbook.
    assert "East at 150" in alt and "West at 50" in alt


def test_docx_chart_descr_lands_on_wp_docPr_with_embedded_values():
    """The user-facing payoff: the resolved values reach the 1.1.1 alt written onto <wp:docPr>."""
    entries = docx_chart_entries(["North", "South", "East", "West"], [120, 70, 150, 50])
    changed = chart_data.chart_descr_edits(entries, ".docx")
    assert changed, "a docx chart with no alt should get one"
    new_xml, alts = changed["word/document.xml"]
    out = new_xml.decode()
    assert 'descr="Bar chart' in out and "East at 150" in out
    assert alts and "Sales by region" in alts[0]
    ET.fromstring(out)


def test_docx_chart_without_its_embedding_is_skipped_not_guessed():
    entries = docx_chart_entries(["North", "South"], [1, 2])
    del entries["word/embeddings/data.xlsx"]
    assert chart_data.charts_in(rezip(entries), ".docx") == []
    assert chart_data.chart_descr_edits(entries, ".docx") == {}


def test_corrupt_embedding_is_survived_not_raised():
    entries = docx_chart_entries(["North", "South"], [1, 2])
    entries["word/embeddings/data.xlsx"] = b"not a zip at all"
    assert chart_data.charts_in(rezip(entries), ".docx") == []

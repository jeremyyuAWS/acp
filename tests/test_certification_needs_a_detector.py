"""A cell may only certify a PASS if a detector in THIS codebase can emit its criterion.

The second half of the pptx 2.1.1 defect, and the half no existing guard covers.

`test_review_outcome.py` asserts the lane can't be bypassed: a 🟡/🔴 criterion never returns PASS.
That says nothing about a 🟢 cell. A 🟢 cell certifies precisely when its validator "ran and found
nothing" — and membership in `RULE_FORMATS` is the only thing that establishes "ran". Nothing
checks that a detector exists behind the membership, so a pair listed there with no implementation
resolves a clean scan to a green PASS on zero evidence, and no test complains.

That is exactly how pptx 2.1.1 shipped. `config/rule-catalog.json` declares `PPTX-ANIM-001`
(WCAG 2.1.1) whose `source` is
`digital-accessibility/DigitalA11y.Analysers.DotNet/Pptx/Rules/AnimationOrderRule.cs` — a path in
an UPSTREAM repository that this one does not vendor. `RULE_FORMATS` trusted the declaration, the
vendored analysers never implemented the rule, and `office_structure.checks_for` dispatches ten
pptx checks of which not one emits 2.1.1. The lane happened to be 🔴 human, so the fix for that
defect masks it; had the lane been 🟢 the fabricated PASS would still be shipping.

So the catalog is a CLAIM, not evidence, and this module refuses it as such. Evidence is a finding
observed coming out of the real scan path on a file built to violate the criterion — the same
"run the detectors against built files rather than reading them" rule that found nine understated
capability cells and the 1.4.3 PDF shipping damage.

Three evidence tiers, because not every engine is vendored (ADR 0012):
  FIRST_PARTY — pure-Python detector; always exercised here.
  ENGINE      — vendored .NET analyser; exercised when built, skipped when not.
  PDF_ENGINE  — the PDF analyser, loaded at runtime from ACP_PDF_ENGINE and NOT vendored, so it
                cannot be exercised locally at all. Named, never assumed.
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

import assessment_policy as ap  # noqa: E402
import office_structure  # noqa: E402
import remediation_capability as rc  # noqa: E402
import rule_registry  # noqa: E402
import scanner  # noqa: E402

FIRST_PARTY, ENGINE, PDF_ENGINE = "first-party", "engine", "pdf-engine"

# (format, SC) -> which detector tier backs its ability to certify a PASS.
# Every cell that certifies MUST appear here; test_every_certifying_cell_is_declared enforces it,
# so adding one to RULE_FORMATS forces you to say what implements it.
EVIDENCE: dict[tuple[str, str], str] = {
    # html — scanner.py, pure Python, all seventeen exercised below.
    **{("html", sc): FIRST_PARTY for sc in (
        "1.3.1", "1.3.4", "1.3.5", "1.4.1", "1.4.2", "1.4.3", "1.4.4", "1.4.10", "1.4.12",
        "2.4.1", "2.4.2", "2.4.3", "2.4.7", "2.5.3", "3.1.1", "3.3.2", "4.1.2")},
    # docx — pseudo-heading and content-control labels are first-party; contrast, title and
    # language come from the .NET analyser (there is no first-party docx contrast check).
    ("docx", "1.3.1"): FIRST_PARTY,
    ("docx", "3.3.2"): FIRST_PARTY,
    ("docx", "1.4.3"): ENGINE,
    ("docx", "2.4.2"): ENGINE,
    ("docx", "3.1.1"): ENGINE,
    # pptx — only contrast is first-party.
    ("pptx", "1.4.3"): FIRST_PARTY,
    ("pptx", "1.3.1"): ENGINE,
    ("pptx", "1.3.2"): ENGINE,
    ("pptx", "2.4.2"): ENGINE,
    ("pptx", "3.1.1"): ENGINE,
    ("xlsx", "1.4.3"): FIRST_PARTY,
    ("xlsx", "1.3.1"): ENGINE,
    ("xlsx", "1.3.2"): ENGINE,
    ("xlsx", "2.4.2"): ENGINE,
    ("xlsx", "3.1.1"): ENGINE,
    # pdf — contrast and the bookmark outline are first-party (pikepdf/pdfplumber); title and
    # language are rules of the non-vendored analyser, named in scanner._pdf_correct_title.
    ("pdf", "1.4.3"): FIRST_PARTY,
    ("pdf", "2.4.1"): FIRST_PARTY,
    ("pdf", "2.4.2"): PDF_ENGINE,   # rule id "pdf.document-title"
    ("pdf", "3.1.1"): PDF_ENGINE,   # rule id "pdf.document-language"
}


def _certifying() -> set[tuple[str, str]]:
    """Every (format, SC) that resolves a CLEAN scan to a certified PASS, read from the code."""
    rule_registry.load()
    return {(fmt, sc)
            for fmt, scs in rc.assessment_table().items()
            for sc in scs
            if ap._rule_outcome(sc, fmt, 0, 0) == "PASS"}


def _scs(issues) -> set[str]:
    out = set()
    for i in issues:
        w = str(i.get("wcag", ""))
        m = re.match(r"(\d+\.\d+\.\d+)", w) or re.match(r"SC_(\d+)_(\d+)_(\d+)$", w)
        if m:
            out.add(m.group(1) if "." in m.group(0) else ".".join(m.groups()))
    return out


# ── THE GUARD ─────────────────────────────────────────────────────────────────
def test_every_certifying_cell_is_declared():
    """No cell may certify a PASS without a declared detector tier.

    This is the guard the pptx 2.1.1 class needed. It costs nothing to run and fails the moment
    someone adds a (criterion, format) to RULE_FORMATS — or flips a lane to 🟢 — without saying
    what actually implements the check.
    """
    undeclared = sorted(_certifying() - set(EVIDENCE))
    assert not undeclared, (
        f"these cells certify a PASS with no declared detector: {undeclared}. Add an EVIDENCE "
        "entry naming the tier that implements it, and a fixture proving it fires — a "
        "rule-catalog entry is NOT evidence (see this module's docstring)."
    )


def test_no_stale_evidence_entries():
    """The converse: EVIDENCE must not claim cells that no longer certify, or it rots into a
    list nobody trusts. pptx 2.1.1 must never appear here."""
    stale = sorted(set(EVIDENCE) - _certifying())
    assert not stale, f"EVIDENCE declares cells that no longer certify: {stale}"


def test_the_catalog_is_not_evidence():
    """A rule-catalog entry whose implementation lives upstream must never certify anything.

    pptx 2.1.1 is the case in point: declared as PPTX-ANIM-001 with a `source` outside this repo,
    implemented nowhere in it. Its lane is 🔴 human so it now resolves REVIEW — but this asserts
    the CATALOG side too, so the next rule declared-but-not-vendored is caught even if its lane
    is 🟢 and the lane guard therefore stays silent.
    """
    catalog = json.loads((ACP / "config" / "rule-catalog.json").read_text())
    claimed = {(fmt, r["wcag_sc"]) for fmt in ("docx", "pptx", "xlsx", "pdf")
               for r in catalog.get(fmt, [])}
    assert ("pptx", "2.1.1") in claimed, "fixture check: the catalog should still declare it"
    # Anything the catalog claims but no tier here declares must not be certifying.
    undeclared_claims = sorted(c for c in claimed if c not in EVIDENCE)
    offenders = [c for c in undeclared_claims if c in _certifying()]
    assert not offenders, (
        f"the catalog claims a rule for {offenders} and they certify a PASS, but no detector "
        "tier is declared for them — a claim is not an implementation"
    )
    assert ap._rule_outcome("2.1.1", "pptx", 0, 0) != "PASS"


# ── TEETH: the declared first-party detectors really do fire ──────────────────
def _html_pages() -> dict[str, str]:
    """One page per certifying html criterion; several fire together, which is fine."""
    return {
        "broad": (
            '<html><head><meta name="viewport" content="width=1024, user-scalable=no">'
            '<style>*:focus{outline:none}</style></head><body>'
            '<p>See <a href="/x">this</a></p><input type="text" name="email">'
            '<input type="checkbox" tabindex="3"><button aria-label="Submit form">Send</button>'
            '<audio autoplay src="a.mp3"></audio><div>no landmark</div></body></html>'),
        "colour_only": ('<html lang="en"><head><title>t</title></head><body><main><p>See '
                        '<a href="/a" style="color:#c00;text-decoration:none">the report</a>'
                        '</p></main></body></html>'),
        "contrast": ('<html lang="en"><head><title>t</title></head><body><main>'
                     '<p style="color:#eeeeee;background:#ffffff">Faint</p></main></body></html>'),
        "orientation": ('<html lang="en"><head><title>t</title><style>'
                        '@media (orientation: portrait){ .app { display: none } }'
                        '</style></head><body><main><p>x</p></main></body></html>'),
        "line_height": ('<html lang="en"><head><title>t</title></head><body><main>'
                        '<p style="line-height:10px">Tight</p></main></body></html>'),
        "no_skip": ('<html lang="en"><head><title>t</title></head><body><nav>'
                    '<a href="/a">Home</a></nav><div>content</div></body></html>'),
        "required": ('<html lang="en"><head><title>t</title></head><body><main><form>'
                     '<input type="text" name="q" required></form></main></body></html>'),
    }


def test_first_party_html_detectors_fire(tmp_path):
    emitted: set[str] = set()
    for name, html in _html_pages().items():
        p = tmp_path / f"{name}.html"
        p.write_text(html)
        emitted |= _scs(scanner._analyse_html(p).get("issues", []))
    want = {sc for (fmt, sc), tier in EVIDENCE.items() if fmt == "html" and tier == FIRST_PARTY}
    missing = sorted(want - emitted, key=lambda s: [int(x) for x in s.split(".")])
    assert not missing, (
        f"html cells certify a PASS but no detector emitted them on a violating page: {missing}"
    )


def _office_fixtures(d: Path) -> None:
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.shared import Pt, RGBColor
    doc = Document()
    h = doc.add_paragraph().add_run("Looks Like A Heading"); h.bold = True; h.font.size = Pt(20)
    faint = doc.add_paragraph().add_run("Faint grey body text on white")   # docx 1.4.3 (ENGINE)
    faint.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE); faint.font.size = Pt(11)
    para = doc.add_paragraph("Agree?")
    sdt = OxmlElement("w:sdt"); sp = OxmlElement("w:sdtPr")
    sp.append(OxmlElement("w:checkbox")); sdt.append(sp)
    body = OxmlElement("w:sdtContent"); body.append(para._p); sdt.append(body)
    doc.element.body.append(sdt)
    doc.save(d / "a.docx")

    from pptx import Presentation
    from pptx.dml.color import RGBColor as PRGB
    from pptx.util import Inches, Pt as PPt
    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(1))
    box.fill.solid(); box.fill.fore_color.rgb = PRGB(0xFF, 0xFF, 0xFF)
    run = box.text_frame.paragraphs[0].add_run(); run.text = "Faint"
    run.font.color.rgb = PRGB(0xEE, 0xEE, 0xEE); run.font.size = PPt(24)
    # a multi-row table with no designated header row — pptx 1.3.1 (ENGINE)
    gf = s.shapes.add_table(3, 2, Inches(1), Inches(5.2), Inches(6), Inches(1.5))
    gf.table.first_row = False
    for i in range(3):
        for j in range(2):
            gf.table.cell(i, j).text = f"r{i}c{j}"
    prs.save(d / "a.pptx")

    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.worksheet.table import Table
    wb = Workbook(); ws = wb.active
    for row in [["Name", "Amount"], ["a", 1], ["b", 2]]:
        ws.append(row)
    for rr in ws.iter_rows():
        for c in rr:
            c.font = Font(color="EEEEEE")
    # a DEFINED table whose header row is not designated — xlsx 1.3.1 (ENGINE)
    tbl = Table(displayName="Data", ref="A1:B3"); tbl.headerRowCount = 0
    ws.add_table(tbl)
    wb.save(d / "a.xlsx")

    from reportlab.lib.colors import Color
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(d / "a.pdf"))
    for i in range(office_structure._MIN_PAGES_FOR_OUTLINE + 2):
        c.setFillColor(Color(0.93, 0.93, 0.93)); c.setFont("Helvetica", 12)
        c.drawString(72, 720, f"Faint page {i + 1} text"); c.showPage()
    c.save()


def test_first_party_office_and_pdf_detectors_fire(tmp_path):
    """checks_for is pure Python, so this needs no .NET build and runs everywhere."""
    _office_fixtures(tmp_path)
    for fmt, ext in (("docx", ".docx"), ("pptx", ".pptx"), ("xlsx", ".xlsx"), ("pdf", ".pdf")):
        emitted = _scs(office_structure.checks_for(tmp_path / f"a{ext}", ext))
        want = {sc for (f, sc), tier in EVIDENCE.items() if f == fmt and tier == FIRST_PARTY}
        missing = sorted(want - emitted, key=lambda s: [int(x) for x in s.split(".")])
        assert not missing, (
            f"{fmt} cells declared FIRST_PARTY certify a PASS but nothing emitted them: {missing}"
        )


def test_engine_backed_detectors_fire_when_the_engine_is_built(tmp_path):
    """The .NET-backed cells, exercised only when the analyser is actually built.

    Skipping rather than failing matches tests/test_pptx_engine_detection.py: an unbuilt CLI is a
    local-environment state, not a defect. In CI the analysers are built, so these do run there.
    """
    _office_fixtures(tmp_path)
    engine = scanner._analyse_office(tmp_path)
    if not any(v.get("succeeded") for v in engine.values()):
        pytest.skip("the .NET Office analysers did not run (CLI not built)")
    for fmt, ext in (("docx", ".docx"), ("pptx", ".pptx"), ("xlsx", ".xlsx")):
        emitted = _scs(engine.get(f"a{ext}", {}).get("issues", []))
        want = {sc for (f, sc), tier in EVIDENCE.items() if f == fmt and tier == ENGINE}
        # Only assert the ones this fixture is built to violate; a fixture gap is not a
        # detector gap, and claiming otherwise is how nine capability cells got misread.
        exercised = want & {"1.4.3", "2.4.2", "1.3.1"}
        missing = sorted(exercised - emitted, key=lambda s: [int(x) for x in s.split(".")])
        assert not missing, f"{fmt} ENGINE cells did not fire: {missing}"


def test_pdf_engine_cells_are_named_not_assumed():
    """The two cells no local run can reach must at least name a real rule in the scan path.

    ADR 0012 does not vendor the PDF analyser, so pdf 2.4.2 / 3.1.1 cannot be fired here. Two
    things ARE checkable, and they are the honest substitute:

      * each rule is named in code that runs against the real engine — `pdf.document-title` in
        scanner._pdf_correct_title, `pdf.document-language` in tests/test_scan.py, which asserts
        it stays silent on a PDF that declares a language (an engine-backed assertion that only
        means anything because the rule exists);
      * a missing engine reports an ERROR, so absence never reads as a clean pass. That is the
        property that stops "no engine" from becoming "no findings, therefore PASS".
    """
    haystack = "\n".join(p.read_text(errors="ignore")
                         for p in [ACP / "api" / "scanner.py", ACP / "tests" / "test_scan.py"])
    for rule in ("pdf.document-title", "pdf.document-language"):
        assert rule in haystack, (
            f"{rule} is no longer referenced anywhere that runs against the engine — pdf "
            "certification for its criterion has lost the only evidence there was"
        )
    # Matched loosely on purpose: pinning the exact literal would break on a reformat and read as
    # a lost safety property when nothing had changed.
    pdf_src = (ACP / "api" / "scanner.py").read_text()
    body = pdf_src[pdf_src.index("def _analyse_pdf"):]
    body = body[:body.index("\ndef ", 1)]
    assert re.search(r'["\']succeeded["\']\s*:\s*False', body), (
        "an unavailable PDF engine must report succeeded=False, never an empty clean scan that "
        "downstream reads as 'validator ran, no findings' and certifies"
    )

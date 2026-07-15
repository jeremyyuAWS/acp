"""Batch alt at assess time (the '14 images show a template' fix).

Before this, an Office file's unlabelled images only got a per-image thumbnail + AI draft at
REMEDIATION time; at assess/review time the 1.1.1 card showed a single fill-in template. Now
alt_proposals_for_office runs the exact fix-time alt logic standalone, so the review card
arrives with a per-image proposal (thumbnail + a vision draft when the model is reachable,
thumbnail-only when it isn't) for every unlabelled image — no remediation pass required.
"""
from __future__ import annotations

import io
import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))

import remediate_office as ro  # noqa: E402

pptx = pytest.importorskip("pptx")


def _deck_with_images(n: int) -> bytes:
    import os

    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(n):
        # Noisy content (not a near-uniform fill) so decorative inference does NOT fire —
        # these are "real content, needs alt" images, the case the fix targets.
        img = Image.frombytes("RGB", (80, 60), os.urandom(80 * 60 * 3))
        b = io.BytesIO()
        img.save(b, "PNG")
        b.seek(0)
        s.shapes.add_picture(b, Inches(1 + i), Inches(1), Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_vision_off_still_gives_a_thumbnail_per_image(monkeypatch):
    # AI off ⇒ no drafts, but every unlabelled image is still surfaced as evidence (a
    # thumbnail), so the card shows a per-image editor instead of one blank template box.
    import ai
    monkeypatch.setattr(ai, "vision_is_available", lambda: False)
    props, evidence = ro.alt_proposals_for_office(_deck_with_images(3), ".pptx", ai_enabled=False)
    assert len(evidence) == 3                       # one picture per image, for the reviewer to see
    assert all(e.get("thumb") and e.get("locator") for e in evidence)


def test_vision_on_pre_drafts_each_image(monkeypatch):
    # Vision reachable ⇒ every ungrounded image gets a real draft proposal with a thumbnail,
    # so the reviewer approves/edits N descriptions instead of authoring from scratch.
    import ai
    monkeypatch.setattr(ai, "vision_is_available", lambda: True)
    calls = {"n": 0}

    def fake_desc(img, **kw):
        calls["n"] += 1
        return {"alt": f"A described image number {calls['n']}", "grounded": False,
                "model": "llava:13b"}
    monkeypatch.setattr(ai, "describe_image_structured", fake_desc)

    props, evidence = ro.alt_proposals_for_office(_deck_with_images(3), ".pptx", ai_enabled=True)
    drafts = [p for p in props if (p.get("proposed_value") or "").startswith("A described image")]
    assert len(drafts) == 3                          # one AI draft per image
    assert all(p.get("thumb") for p in drafts)       # each carries its picture


def test_unsupported_ext_and_bad_bytes_degrade_to_empty():
    assert ro.alt_proposals_for_office(b"", ".pdf", ai_enabled=True) == ([], [])
    assert ro.alt_proposals_for_office(b"not a zip", ".pptx", ai_enabled=True) == ([], [])


def test_handler_wires_assess_time_alt():
    src = (Path(__file__).resolve().parent.parent / "api" / "handlers.py").read_text()
    assert "alt_proposals_for_office" in src
    assert "attach_hitl_evidence(scan_id, filename, \"1.1.1\", img_evidence)" in src


def test_consensus_verdict_attached_when_validator_configured(monkeypatch):
    # With a distinct validator model set, every ungrounded draft carries a pre-computed
    # agreement verdict so the card shows "two models agree" without a reviewer click.
    import ai
    import remediate_office as ro2
    monkeypatch.setattr(ai, "_ALT_VALIDATOR_MODEL", "llava:13b")
    monkeypatch.setattr(ai, "vision_is_available", lambda: True)
    monkeypatch.setattr(ai, "describe_image_structured",
                        lambda img, **kw: {"alt": "a described chart", "grounded": False, "model": "qwen2.5vl:7b"})
    monkeypatch.setattr(ai, "validate_alt_text",
                        lambda img, alt, **kw: {"verdict": "consistent", "second_opinion": "a chart",
                                                "validator_model": "llava:13b"})
    props, _ = ro2.alt_proposals_for_office(_deck_with_images(2), ".pptx", ai_enabled=True)
    withagree = [p for p in props if p.get("agreement")]
    assert withagree and all(p["agreement"]["verdict"] == "consistent" for p in withagree)
    assert withagree[0]["agreement"]["validator_model"] == "llava:13b"


def test_no_validator_no_agreement(monkeypatch):
    import ai
    import remediate_office as ro2
    monkeypatch.setattr(ai, "_ALT_VALIDATOR_MODEL", "")   # no distinct validator → no cross-check
    monkeypatch.setattr(ai, "vision_is_available", lambda: True)
    monkeypatch.setattr(ai, "describe_image_structured",
                        lambda img, **kw: {"alt": "x", "grounded": False, "model": "qwen2.5vl:7b"})
    props, _ = ro2.alt_proposals_for_office(_deck_with_images(1), ".pptx", ai_enabled=True)
    assert all(not p.get("agreement") for p in props)


def test_card_renders_precomputed_agreement():
    pe = (Path(__file__).resolve().parent.parent / "frontend" / "src" / "ProposalEditors.jsx").read_text()
    assert "p?.agreement" in pe and "Two models agree" in pe

"""The 1.3.3 sensory-rewrite lane on docx, pptx AND xlsx, proved end to end.

Three lanes in one module because they are one criterion, one store getter and one writer —
`apply_text_values.apply_sensory_rewrite`, parameterised internally by a `_Dialect` per format.
Parametrising the test the same way is what asserts the differing half: the three dialects
disagree on the paragraph/run/text tag names, on the text element's spelling
(`<w:t xml:space="preserve">` / `<a:t>` / `<t xml:space="preserve">`) and on which parts are
walked (one document, every slide, every worksheet plus the shared-string table). A test written
for docx alone would prove the shared machinery three times over and the dialects not at all.

Same bar as the lanes before: the original document trips the finding, an approval changes the
saved document, a REAL re-scan verifies it, unrelated content survives, and a broken engine
earns no credit. Nothing but the blob store is patched.

WHY THE APPROVED VALUE IS AUTHORED RATHER THAN PROPOSED — stated because it is the one place
this proof is narrower than the 2.4.4 and 3.1.2 ones. `propose_sensory_rewrite` returns [] when
no AI text model is available (`ai.model_is_available()` gates it), so a proposer-driven test
would skip on any host without one — and a lane that can only be proved where a model happens to
be installed is not proved. A reviewer rewriting the sentence themselves is also the real
workflow whenever the draft is wrong, which on a Low-confidence subjective judgement is often.
`handlers` reads the approved value identically either way, so the lane under test is the same;
what is NOT covered here is the proposer, and the 3.1.2 proofs cover that seam on the same
writer module.

WHAT THIS CLAIMS: the sentence a reader sees no longer identifies its target by shape or
position alone, the document survived, and ACP's own criterion stops firing. NOT that the new
sentence is good instructional prose — 1.3.3 is a judgement, which is why the lane is `assisted`
and its proposals are always Low confidence.

WHERE THE DETECTOR'S EDGE IS, and it matters for reading the controls below: `textchecks`
anchors on an instruction VERB followed within one sentence by a shape-only or position-only
reference. "Click the green button" is NOT a finding — colour alone is deliberately excluded as
too noisy — so the fixture says "on the right", and the rewrite has to remove the POSITION, not
the colour, to clear it. A control that changed only the colour word would pass a weaker test
and leave the document failing.
"""
from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

SID = "rv-office-133"

SENSORY = "Click the button on the right to submit the form."
APPROVED = "Select Submit to send the form."
STILL_SENSORY = "Press the control at the bottom of the page to submit."
TITLE = "How to submit"
TAIL = "Unrelated closing text that must survive the write."


# ── the three documents ───────────────────────────────────────────────────────

def _docx(sentence: str) -> bytes:
    from docx import Document
    d = Document()
    d.add_heading(TITLE, level=1)
    p = d.add_paragraph(sentence)
    p.add_run(" Bold tail.").bold = True
    d.add_paragraph(TAIL)
    out = Path(tempfile.mkdtemp()) / "guide.docx"
    d.save(out)
    return out.read_bytes()


def _pptx(sentence: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = TITLE
    tf = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3)).text_frame
    tf.word_wrap = True
    tf.text = sentence
    tf.add_paragraph().text = TAIL
    out = Path(tempfile.mkdtemp()) / "guide.pptx"
    prs.save(out)
    return out.read_bytes()


def _xlsx(sentence: str) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Instructions"
    ws["A1"] = TITLE
    ws["A2"] = sentence
    ws["A3"] = TAIL
    ws["B2"] = 42
    out = Path(tempfile.mkdtemp()) / "guide.xlsx"
    wb.save(out)
    return out.read_bytes()


BUILD = {"docx": _docx, "pptx": _pptx, "xlsx": _xlsx}
FILENAME = {f: f"guide.{f}" for f in BUILD}
DEP = {"docx": "docx", "pptx": "pptx", "xlsx": "openpyxl"}
# The part carrying the prose, per format — the difference the parametrisation exists to assert.
#
# xlsx names the WORKSHEET here, and that is a fact about the fixture rather than about Excel.
# SpreadsheetML has two ways to hold a cell's text: pooled in `xl/sharedStrings.xml` and
# referenced by index (`t="s"`), or written inline in the cell (`t="inlineStr"`, `<is><t>`).
# openpyxl writes inline strings and no shared-string table at all, so this fixture exercises
# only that half — `test_the_rewrite_also_reaches_a_real_shared_string_table` below covers the
# pooled form, which is what Excel itself writes and what most real workbooks are.
PROSE_PART = {"docx": "word/document.xml",
              "pptx": "ppt/slides/slide1.xml",
              "xlsx": "xl/worksheets/sheet1.xml"}


@pytest.fixture(params=("docx", "pptx", "xlsx"))
def fmt(request):
    pytest.importorskip(DEP[request.param])
    return request.param


# ── helpers ───────────────────────────────────────────────────────────────────

def _assess(data: bytes, fmt: str) -> set[str]:
    """The SCs a REAL assessment reports — the same call the production re-verification makes."""
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    name = FILENAME[fmt]
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / name).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), name, detect_pii=False)
    return {sc for i in (fd or {}).get("issues", []) if (sc := _extract_sc(i.get("wcag", "")))}


def _spill(data: bytes, fmt: str) -> Path:
    p = Path(tempfile.mkdtemp()) / FILENAME[fmt]
    p.write_bytes(data)
    return p


def _text(data: bytes, fmt: str) -> str:
    """The document's extracted text — what a reader sees, through ACP's own extractor, which
    is also what the 1.3.3 detector reads."""
    import pii
    return " ".join((pii.extract_text(_spill(data, fmt)) or "").split())


def _part(data: bytes, name: str) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(name).decode("utf-8")


class _Blob:
    def __init__(self, data: bytes):
        self.data, self.uploads = data, []

    def download_remediated(self, owner, sid, f):
        return self.data

    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data
        self.uploads.append((f, mime))
        return "http://b/2"


@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rv.db")
    return store_mod.Store()


def _seed(store, fmt: str, locator: str | None, value: str | None) -> str:
    """A scanned + remediated file with one 1.3.3 card, and `value` approved on it."""
    name = FILENAME[fmt]
    store.init_scan_run(SID, "drive", 1, "2026-08-31T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": name, "engine": "office", "status": "pass", "score": 60, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "d1",
        "issues": [{"ruleId": "SENSORY_INSTRUCTION", "wcag": "1.3.3 Sensory Characteristics",
                    "severity": "SERIOUS"}],
    }, "2026-08-31T00:00:00Z")
    store.record_remediation(SID, name, drive_write_url="http://d/1", blob_url="http://b/1")
    props = ([{"locator": locator, "before": SENSORY, "proposed_value": "",
               "rationale": "r", "source": "reviewer"}] if locator else [])
    item_id = store.enqueue_proposals(SID, name, "1.3.3", props,
                                      rule_name="Sensory Characteristics")
    store.update_hitl_item(item_id, "approved", None, None)
    store.approve_proposal_values(item_id, [value] if value else [])
    return item_id


def _run_lane(monkeypatch, store, blob, fmt: str):
    """The production handler, with the re-scan UNPATCHED."""
    import core
    import handlers
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    handlers._apply_approved_values({"scan_id": SID, "file": FILENAME[fmt]}, {})


# The locator is a prose prefix, not a part#id or an href — the third locator scheme in this
# suite, and the reason `apply_text_values` exists as its own writer.
LOCATOR = SENSORY[:60]


# ── 1. the finding ────────────────────────────────────────────────────────────

def test_a_real_assessment_reports_1_3_3_on_a_position_only_instruction(fmt):
    assert "1.3.3" in _assess(BUILD[fmt](SENSORY), fmt)


def test_the_rewritten_control_is_not_flagged(fmt):
    """The same document with the position removed. Without this, a detector that fired on any
    imperative sentence would satisfy the test above and nothing would notice."""
    assert "1.3.3" not in _assess(BUILD[fmt](APPROVED), fmt)


def test_colour_alone_is_not_a_finding(fmt):
    """The detector's edge, pinned so the controls below are read correctly. `textchecks`
    anchors on shape or POSITION, never colour — colour alone is noisier and overlaps 1.4.1 —
    so a rewrite that drops "green" and keeps "on the right" fixes nothing, and a fixture that
    relied on colour would never have failed in the first place."""
    assert "1.3.3" not in _assess(BUILD[fmt]("Click the green button to submit the form."), fmt)


# ── 2. approval → write → re-scan → credit, through the real path ─────────────

@pytest.fixture()
def applied(store, monkeypatch, fmt):
    original = BUILD[fmt](SENSORY)
    blob = _Blob(original)
    _seed(store, fmt, LOCATOR, APPROVED)
    _run_lane(monkeypatch, store, blob, fmt)
    return blob, store, fmt, original


def test_the_saved_document_carries_the_rewritten_sentence(applied):
    blob, _, fmt, _ = applied
    text = _text(blob.data, fmt)
    assert APPROVED in text
    assert "on the right" not in text, "the position reference is still in the document"


def test_the_rewrite_lands_in_the_formats_own_prose_part(applied):
    """The dialect assertion. Each format keeps its text somewhere different — Word in the
    document part, PowerPoint in the slide, Excel in the cell or in a pooled table — and a
    writer that only knew one of them would pass the extracted-text check on that format
    alone."""
    blob, _, fmt, _ = applied
    assert APPROVED in _part(blob.data, PROSE_PART[fmt])


def test_unrelated_content_survives(applied):
    blob, _, fmt, _ = applied
    text = _text(blob.data, fmt)
    assert TITLE in text and TAIL in text


def test_the_saved_file_still_opens(applied):
    """Through the library that had no part in writing the change."""
    blob, _, fmt, _ = applied
    path = _spill(blob.data, fmt)
    assert zipfile.ZipFile(path).testzip() is None
    if fmt == "docx":
        from docx import Document
        assert any(TITLE in p.text for p in Document(str(path)).paragraphs)
    elif fmt == "pptx":
        from pptx import Presentation
        assert len(Presentation(str(path)).slides) == 1
    else:
        import openpyxl
        wb = openpyxl.load_workbook(str(path))
        assert wb.sheetnames == ["Instructions"]
        assert wb["Instructions"]["B2"].value == 42, "an unrelated cell value was lost"


def test_a_second_real_assessment_no_longer_reports_1_3_3(applied):
    """THE claim: a fresh assessment of the SAVED bytes, not the writer's return value."""
    blob, _, fmt, _ = applied
    assert "1.3.3" not in _assess(blob.data, fmt)


def test_the_row_is_credited_and_the_copy_is_stored(applied):
    blob, store, fmt, _ = applied
    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 0
    assert blob.uploads


# ── 3. where the lane must NOT credit ─────────────────────────────────────────
#
# The bytes an uncredited lane wrote are DISCARDED — `_apply_one_value_kind` returns the
# original `working` and uploads nothing — so the control establishes what the write produces by
# calling the writer directly, then runs the real lane on the same approval.

def test_a_rewrite_that_is_still_sensory_is_not_credited(store, monkeypatch, fmt):
    """The reviewer swaps one position reference for another: "on the right" becomes "at the
    bottom of the page". The write succeeds, the sentence genuinely changes, and 1.3.3 still
    fails — so nothing may be credited and nothing published.

    This is the test that distinguishes a real re-scan from `residual=set()`. Under that stub it
    would pass by construction, and the reviewer's file would be certified against a criterion
    it still fails.
    """
    from apply_text_values import apply_sensory_rewrite
    original = BUILD[fmt](SENSORY)

    written, ap, _ = apply_sensory_rewrite(original, fmt, {LOCATOR: STILL_SENSORY})
    assert ap, "the writer refused the value, so this control is not about crediting"
    assert STILL_SENSORY in _text(written, fmt), "the rewrite did not reach the document"
    assert "1.3.3" in _assess(written, fmt), (
        "the replacement sentence no longer fails 1.3.3, so this control cannot distinguish a "
        "withheld credit from a cleared one")

    blob = _Blob(original)
    _seed(store, fmt, LOCATOR, STILL_SENSORY)
    _run_lane(monkeypatch, store, blob, fmt)

    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 1
    assert store.mark_file_compliant_if_reviewed(SID, FILENAME[fmt]) is False
    assert not blob.uploads
    assert blob.data == original


def test_an_approval_aimed_at_a_sentence_the_document_does_not_have_is_not_credited(
        store, monkeypatch, fmt):
    """The locator names prose that is not there. Nothing resolves, nothing is written, nothing
    is credited — never a guess at a different sentence."""
    original = BUILD[fmt](SENSORY)
    blob = _Blob(original)
    _seed(store, fmt, "A sentence this document does not contain at all", APPROVED)
    _run_lane(monkeypatch, store, blob, fmt)

    assert blob.data == original
    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 1
    assert not blob.uploads


def test_a_document_with_no_sensory_instruction_is_left_byte_identical(store, monkeypatch, fmt):
    clean = BUILD[fmt](APPROVED)
    assert "1.3.3" not in _assess(clean, fmt)
    blob = _Blob(clean)
    _seed(store, fmt, None, None)
    _run_lane(monkeypatch, store, blob, fmt)
    assert blob.data == clean and not blob.uploads


# ── 3b. the other half of SpreadsheetML's text storage ────────────────────────

def _xlsx_with_shared_strings(sentence: str) -> bytes:
    """The same workbook, but with its text POOLED in `xl/sharedStrings.xml`.

    openpyxl only ever writes inline strings, so the parametrised fixture above cannot reach
    the pooled form — and the pooled form is what Excel itself writes, which makes it the shape
    almost every real workbook arrives in. Rather than trust that `_xlsx_parts` reads both, this
    converts the package by hand and lets the lane prove it.
    """
    import re as _re
    data = _xlsx(sentence)
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        parts = {n: z.read(n) for n in z.namelist()}

    sheet = parts["xl/worksheets/sheet1.xml"].decode("utf-8")
    pool: list[str] = []

    def _pool(m: _re.Match) -> str:
        text = m.group(2)
        if text not in pool:
            pool.append(text)
        return f'{m.group(1)}t="s"><v>{pool.index(text)}</v></c>'

    sheet = _re.sub(r'(<c\b[^>]*?)t="inlineStr">\s*<is><t[^>]*>([^<]*)</t></is>\s*</c>',
                    _pool, sheet)
    assert pool, "the fixture no longer holds inline strings, so nothing was pooled"
    parts["xl/worksheets/sheet1.xml"] = sheet.encode("utf-8")

    si = "".join(f"<si><t xml:space=\"preserve\">{t}</t></si>" for t in pool)
    parts["xl/sharedStrings.xml"] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        f'count="{len(pool)}" uniqueCount="{len(pool)}">{si}</sst>').encode("utf-8")

    ct = parts["[Content_Types].xml"].decode("utf-8")
    parts["[Content_Types].xml"] = ct.replace(
        "</Types>",
        '<Override PartName="/xl/sharedStrings.xml" ContentType="application/vnd.'
        'openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"/></Types>'
    ).encode("utf-8")

    rels = parts["xl/_rels/workbook.xml.rels"].decode("utf-8")
    parts["xl/_rels/workbook.xml.rels"] = rels.replace(
        "</Relationships>",
        '<Relationship Id="rIdSst" Type="http://schemas.openxmlformats.org/officeDocument/'
        '2006/relationships/sharedStrings" Target="sharedStrings.xml"/></Relationships>'
    ).encode("utf-8")

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w") as z:
        for name, payload in parts.items():
            z.writestr(name, payload)
    return out.getvalue()


def test_the_shared_string_fixture_is_really_pooled():
    """The bite check on the fixture itself: if the conversion silently did nothing, the test
    below would be a second copy of the inline-string case wearing a different name."""
    pytest.importorskip("openpyxl")
    data = _xlsx_with_shared_strings(SENSORY)
    assert SENSORY in _part(data, "xl/sharedStrings.xml")
    assert "inlineStr" not in _part(data, "xl/worksheets/sheet1.xml")
    assert 't="s"' in _part(data, "xl/worksheets/sheet1.xml")


def test_the_rewrite_also_reaches_a_real_shared_string_table(store, monkeypatch):
    """The pooled form, end to end through the real lane."""
    pytest.importorskip("openpyxl")
    original = _xlsx_with_shared_strings(SENSORY)
    assert "1.3.3" in _assess(original, "xlsx"), "the pooled fixture does not trip the finding"

    blob = _Blob(original)
    _seed(store, "xlsx", LOCATOR, APPROVED)
    _run_lane(monkeypatch, store, blob, "xlsx")

    assert APPROVED in _part(blob.data, "xl/sharedStrings.xml")
    assert "1.3.3" not in _assess(blob.data, "xlsx")
    assert store.count_unapplied_approved_values(SID, FILENAME["xlsx"]) == 0
    assert blob.uploads

    import openpyxl
    wb = openpyxl.load_workbook(str(_spill(blob.data, "xlsx")))
    assert wb["Instructions"]["A2"].value == APPROVED
    assert wb["Instructions"]["B2"].value == 42


# ── 4. a broken engine earns nothing ──────────────────────────────────────────

@pytest.mark.parametrize("name,script,timeout", [
    ("cannot be launched", None, None),
    ("exits non-zero", "#!/bin/sh\necho boom >&2\nexit 9\n", None),
    ("hangs past the timeout", "#!/bin/sh\nsleep 30\n", "2"),
])
def test_a_broken_office_analyser_never_credits_this_lane_either(monkeypatch, fmt,
                                                                 name, script, timeout):
    """Re-asserted per lane rather than assumed to inherit: the fail-open #1058 closed lived in
    ONE shared seam, so a regression there takes every lane at once. 1.3.3 comes from
    textchecks, which runs after the .NET call in its own try/except, so the residual is a real
    set even with no analyser at all."""
    import stat as _stat

    import scanner
    if script is None:
        monkeypatch.setattr(scanner, "DOTNET", "/nonexistent/dotnet", raising=False)
    else:
        fake = Path(tempfile.mkdtemp()) / "dotnet"
        fake.write_text(script)
        fake.chmod(fake.stat().st_mode | _stat.S_IEXEC)
        monkeypatch.setattr(scanner, "DOTNET", str(fake), raising=False)
    if timeout:
        monkeypatch.setenv("ACP_OFFICE_CLI_TIMEOUT", timeout)

    from proposals import verify_residual_scs
    residual = verify_residual_scs(BUILD[fmt](SENSORY), FILENAME[fmt])
    assert residual is not None, (
        f"an office CLI that {name} made the re-scan return None — every approved value on this "
        f"lane would be credited on a scan that never happened")
    assert "1.3.3" in residual

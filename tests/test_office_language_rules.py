"""Regression tests for the owned Office analyser's WCAG 3.1.1 / 3.1.2 language rules.

These guard the fix in engine/office-analysers (ADR 0012): the DocumentLanguage rules
must read the *content* language (docx w:lang, pptx a:rPr/@lang), not only the optional
dc:language metadata, and LanguageOfParts must only flag genuine foreign-script passages
rather than every unmarked run.

The test shells out to the built .NET CLI, so it self-skips when the toolchain
(python-docx/pptx, dotnet, or the built AcpScan.Cli.dll) is not present — e.g. on a CI
runner without .NET. Where the toolchain exists it is a true end-to-end guard.
"""
from __future__ import annotations

import json
import os
import re
import subprocess
import zipfile
from pathlib import Path

import pytest

docx = pytest.importorskip("docx")
pptx = pytest.importorskip("pptx")
from docx import Document  # noqa: E402
from docx.oxml.ns import qn  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

_ACP = Path(__file__).resolve().parents[1]
_DOTNET = Path(os.environ.get("ACP_DOTNET") or os.path.expanduser("~/.dotnet/dotnet"))
_CLI = Path(os.environ.get("ACP_OFFICE_CLI")
            or _ACP / "spike/dotnet/AcpScan.Cli/bin/Release/net10.0/AcpScan.Cli.dll")

pytestmark = pytest.mark.skipif(
    not _DOTNET.exists() or not _CLI.exists(),
    reason="office analyser CLI or dotnet runtime not built/available",
)


def _strip_docx_lang(src: Path, dst: Path) -> None:
    with zipfile.ZipFile(src) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}
    for n in ("word/styles.xml", "word/document.xml"):
        if n in data:
            x = data[n].decode("utf-8")
            x = re.sub(r"<w:lang\b[^>]*/>", "", x)
            x = re.sub(r"<w:lang\b[^>]*>.*?</w:lang>", "", x)
            data[n] = x.encode("utf-8")
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])


def _strip_pptx_lang(src: Path, dst: Path) -> None:
    with zipfile.ZipFile(src) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}
    for n in names:
        if n.startswith("ppt/") and n.endswith(".xml"):
            data[n] = re.sub(rb'\s+lang="[^"]*"', b"", data[n])
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])


def _pptx_lang_to_altlang(src: Path, dst: Path) -> None:
    """Rewrite every run-level `lang="..."` attribute to `altLang="..."`, so the only
    declared language anywhere in the deck is via altLang (DrawingML's east-Asian-language
    counterpart to lang, e.g. CJK text in an otherwise Latin-script deck)."""
    with zipfile.ZipFile(src) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}
    for n in names:
        if n.startswith("ppt/") and n.endswith(".xml"):
            data[n] = re.sub(rb'\blang="', b'altLang="', data[n])
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])


def _build_fixtures(d: Path) -> None:
    # Content language via run-level w:lang, blank metadata -> must PASS.
    doc = Document()
    run = doc.add_paragraph().add_run(
        "This is an English paragraph with more than twenty words in it so that the "
        "language of parts rule has a long run to consider during analysis here today.")
    rpr = run._element.get_or_add_rPr()
    rpr.append(rpr.makeelement(qn("w:lang"), {qn("w:val"): "en-US"}))
    doc.core_properties.language = ""
    doc.save(d / "content_lang.docx")
    _strip_docx_lang(d / "content_lang.docx", d / "no_lang.docx")

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tb.text_frame.text = "English presentation text in a normal run of words."
    pres.core_properties.language = ""
    pres.save(d / "run_lang.pptx")
    _strip_pptx_lang(d / "run_lang.pptx", d / "no_lang.pptx")
    _pptx_lang_to_altlang(d / "run_lang.pptx", d / "altlang_only.pptx")

    # metadata set + long same-language run without run-lang -> must NOT false-positive.
    same = Document()
    same.add_paragraph(
        "This paragraph is a perfectly ordinary English sentence written with well over "
        "twenty words so the language of parts rule has a long run to evaluate here now.")
    same.core_properties.language = "en-US"
    same.save(d / "parts_same_lang.docx")

    # metadata set + a long CJK-script run without run-lang -> genuine 3.1.2, must FLAG.
    foreign = Document()
    foreign.add_paragraph("An English introduction establishing the Latin document script.")
    foreign.add_paragraph(" ".join(["中文"] * 25))
    foreign.core_properties.language = "en-US"
    foreign.save(d / "parts_foreign.docx")


def _run_cli(indir: Path, tmp_path: Path) -> dict[str, list[dict]]:
    out = tmp_path / "out.json"
    env = {**os.environ, "DOTNET_ROOT": os.path.expanduser("~/.dotnet"),
           "DOTNET_CLI_TELEMETRY_OPTOUT": "1", "DOTNET_NOLOGO": "1"}
    subprocess.run([str(_DOTNET), str(_CLI), str(indir), str(out)],
                   capture_output=True, text=True, env=env, check=True)
    return {item["file"]: item.get("issues", []) for item in json.loads(out.read_text())}


def _has_lang_finding(issues: list[dict]) -> bool:
    # SC_3_1_1 and SC_3_1_2 are now distinct WcagCriterion members (see ADR 0012's
    # 2026-07-23 amendment) — DocumentLanguageRule and LanguageOfPartsRule no longer
    # share a code, so this can key on the SC directly instead of a title substring.
    return any(i.get("wcag") == "SC_3_1_1" for i in issues)


def _has_parts_finding(issues: list[dict]) -> bool:
    return any(i.get("wcag") == "SC_3_1_2" for i in issues)


def test_document_language_reads_content_not_only_metadata(tmp_path):
    _build_fixtures(tmp_path)
    res = _run_cli(tmp_path, tmp_path)

    # Content language present (run w:lang / pptx a:rPr lang) but blank metadata -> PASS.
    assert not _has_lang_finding(res["content_lang.docx"])
    assert not _has_lang_finding(res["run_lang.pptx"])

    # No language anywhere -> still FLAGGED.
    assert _has_lang_finding(res["no_lang.docx"])
    assert _has_lang_finding(res["no_lang.pptx"])

    # Language declared only via altLang (DrawingML's east-Asian-language counterpart to
    # lang) -> must PASS, not be treated as undeclared.
    assert not _has_lang_finding(res["altlang_only.pptx"])


def test_language_of_parts_only_flags_foreign_script(tmp_path):
    _build_fixtures(tmp_path)
    res = _run_cli(tmp_path, tmp_path)

    # Ordinary same-language long run -> no false positive.
    assert not _has_parts_finding(res["parts_same_lang.docx"])
    # A genuine foreign-script passage with no language -> flagged.
    assert _has_parts_finding(res["parts_foreign.docx"])

"""The 2.4.4 pptx link-text lane, proved end to end (WCAG 2.4.4 Link Purpose (In Context)).

The fifth lane to meet the REMEDIATION-VERIFIED bar, asserting the same five properties as the
four before it: the original deck trips the finding, an approval changes the document, a REAL
re-scan verifies it, unrelated content survives, and a broken engine earns no credit. Nothing
but the blob store is patched — `handlers._apply_approved_values` runs the production seam
through `proposals.verify_residual_scs` to `scanner.analyse_and_assess`.

WHAT IS DIFFERENT ABOUT THIS ONE, and it is not cosmetic. On pptx the link lane must clear TWO
criteria before it credits anything, not one:

    handlers._LINK_SCS_BY_EXT = {"docx": ("2.4.4", "2.4.9"), "pptx": ("2.4.4", "2.4.9"), …}

so `scs_to_clear` is {2.4.4, 2.4.9} and `cleared` is false if EITHER still fails. That makes a
control available here that no single-criterion lane can express: an approved value which fixes
2.4.4 and breaks 2.4.9 in the same write. `test_a_value_that_trades_2_4_4_for_2_4_9_is_not_
credited` does exactly that — it approves, for the vague link, the text the OTHER hyperlink
already uses for a different destination. The vague-text finding genuinely goes away, and the
deck is not better: two links now read identically and point at different places. The lane must
withhold credit, and it can only know that because it re-scanned rather than trusting the write.

WHERE THE WRITE LANDS. `apply_link_text._apply_pptx` rebuilds the run as `rPr + <a:t>`, and in
PresentationML the hyperlink lives INSIDE `<a:rPr>` as `<a:hlinkClick r:id="…">` — so keeping
the run properties is what keeps the link a link, and the same move keeps the author's bold.
Both are asserted below rather than assumed, because a writer that dropped rPr would still make
the visible text correct and would silently destroy the hyperlink it was asked to relabel.

WHAT THIS CLAIMS: ACP's own criteria stop firing on a deck ACP changed, the destination is
unmoved and the rest of the deck survived. NOT that the new wording is the best a human would
choose — 2.4.4 is "in context", so a detector reading link text alone approximates it in both
directions, and this file proves the write and the verification, not the prose.
"""
from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

pytest.importorskip("pptx")

FILE = "q3-review.pptx"
SID = "rv-pptx-244"
SLIDE1 = "ppt/slides/slide1.xml"
SLIDE2 = "ppt/slides/slide2.xml"
RELS2 = "ppt/slides/_rels/slide2.xml.rels"

VAGUE_HREF = "https://example.org/q3-accessibility-report-2026"
GOOD_HREF = "https://example.org/procurement-policy"
VAGUE_TEXT = "click here"
GOOD_TEXT = "supplier accessibility requirements"
TITLE1 = "Q3 accessibility review"
TITLE2 = "Where to read more"
# The replacement a reviewer settles on. Deliberately not the proposer's draft: the point of the
# approval step is that a human may edit it, and the lane must write what they approved.
APPROVED_TEXT = "Q3 accessibility report"


def _deck(second_text: str) -> bytes:
    """Two slides, two hyperlinks: one already descriptive, one under test.

    The second link is not decoration. A writer that relabelled every hyperlink would pass a
    single-link deck, and this is what makes `test_the_other_hyperlink_is_untouched` able to
    fail. Putting them on DIFFERENT slides also exercises the per-slide rels lookup — each
    slide has its own `_rels/slideN.xml.rels`, and a writer that resolved rIds against the
    wrong one would rewrite the wrong link.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = TITLE1
    r1 = s1.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1)).text_frame \
        .paragraphs[0].add_run()
    r1.text = GOOD_TEXT
    r1.hyperlink.address = GOOD_HREF

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = TITLE2
    r2 = s2.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1)).text_frame \
        .paragraphs[0].add_run()
    r2.text = second_text
    r2.font.bold = True                    # authored formatting that must survive the rewrite
    r2.hyperlink.address = VAGUE_HREF

    out = Path(tempfile.mkdtemp()) / FILE
    prs.save(out)
    return out.read_bytes()


def _assess(data: bytes) -> set[str]:
    """The SCs a REAL assessment reports — the same call the production re-verification makes."""
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / FILE).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), FILE, detect_pii=False)
    return {sc for i in (fd or {}).get("issues", []) if (sc := _extract_sc(i.get("wcag", "")))}


def _spill(data: bytes) -> str:
    p = Path(tempfile.mkdtemp()) / FILE
    p.write_bytes(data)
    return str(p)


def _part(data: bytes, name: str) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(name).decode("utf-8")


def _proposals(data: bytes) -> list[dict]:
    """What ACP actually offers a reviewer for this deck. `derive_link_text` reads the URL, so
    this is deterministic with no model — the whole chain can start from the document."""
    from proposals import propose_link_texts
    return propose_link_texts(_spill(data), "pptx", ai_enabled=False)


class _Blob:
    """The only thing patched in this module. Stores bytes verbatim; decides nothing."""

    def __init__(self, data: bytes):
        self.data, self.uploads = data, []

    def download_remediated(self, owner, sid, f):
        return self.data

    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data
        self.uploads.append((f, mime))
        return "http://b/2"


@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rv.db")
    return store_mod.Store()


def _seed(store, props: list[dict], values: list[str]) -> int:
    """A scanned + remediated deck with one 2.4.4 card, and `values` approved on it."""
    store.init_scan_run(SID, "drive", 1, "2026-08-31T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": FILE, "engine": "office", "status": "pass", "score": 60, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "d1",
        "issues": [{"ruleId": "PPTX_LINK_PURPOSE_VAGUE", "wcag": "2.4.4 Link Purpose (In Context)",
                    "severity": "MODERATE"}],
    }, "2026-08-31T00:00:00Z")
    store.record_remediation(SID, FILE, drive_write_url="http://d/1", blob_url="http://b/1")
    item_id = store.enqueue_proposals(SID, FILE, "2.4.4", [
        {k: p.get(k) for k in ("locator", "before", "proposed_value", "rationale", "source")}
        for p in props], rule_name="Link Purpose (In Context)")
    store.update_hitl_item(item_id, "approved", None, None)
    store.approve_proposal_values(item_id, values)
    return item_id


def _run_lane(monkeypatch, store, blob):
    """The production handler, with the re-scan UNPATCHED."""
    import core
    import handlers
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    handlers._apply_approved_values({"scan_id": SID, "file": FILE}, {})


@pytest.fixture(scope="module")
def original() -> bytes:
    return _deck(VAGUE_TEXT)


# ── 1. the finding and its proposal ──────────────────────────────────────────

def test_a_real_assessment_reports_2_4_4_on_the_original(original):
    assert "2.4.4" in _assess(original), (
        "the fixture is supposed to fail 2.4.4 before remediation; without that every "
        "'the fix cleared it' assertion below is vacuous")


def test_the_descriptive_control_deck_is_not_flagged():
    """The same deck with descriptive text on both links. Without this, a detector that fired
    on every hyperlink would satisfy the test above and nothing would notice."""
    assert "2.4.4" not in _assess(_deck(APPROVED_TEXT))


def test_the_original_does_not_already_fail_2_4_9(original):
    """The two link texts differ, so the deck starts clean on the duplicate-text criterion.

    Load-bearing for `test_a_value_that_trades_2_4_4_for_2_4_9_is_not_credited`: if 2.4.9 were
    already failing, that control could not tell a criterion the write BROKE from one it never
    fixed.
    """
    assert "2.4.9" not in _assess(original)


def test_the_proposer_offers_a_value_aimed_at_the_vague_link(original):
    """A finding a reviewer cannot act on is not remediation. The proposal must name the
    destination it belongs to and carry a concrete replacement."""
    mine = [p for p in _proposals(original) if p.get("locator") == VAGUE_HREF]
    assert mine, f"no proposal for {VAGUE_HREF}"
    assert mine[0]["before"] == VAGUE_TEXT
    assert mine[0]["proposed_value"].strip(), "a proposal with no value is a finding with no fix"
    assert mine[0].get("sc") == "2.4.4"


def test_no_proposal_is_raised_for_the_already_descriptive_link(original):
    """Self-gating, asserted. The cheapest place to get this wrong is offering reviewers work
    that is not needed."""
    assert GOOD_HREF not in {p.get("locator") for p in _proposals(original)}


# ── 2. approval → write → re-scan → credit, through the real path ────────────

@pytest.fixture()
def applied(store, monkeypatch, original):
    """The whole chain from the deck: propose, approve an EDITED value, run the lane."""
    props = [p for p in _proposals(original) if p.get("locator") == VAGUE_HREF]
    blob = _Blob(original)
    _seed(store, props, [APPROVED_TEXT])
    _run_lane(monkeypatch, store, blob)
    return blob, store


def test_the_saved_deck_carries_the_approved_text(applied):
    blob, _ = applied
    xml = _part(blob.data, SLIDE2)
    assert APPROVED_TEXT in xml
    assert VAGUE_TEXT not in xml, "the vague text is still in the deck a user will open"


def test_the_hyperlink_target_is_unchanged(applied):
    """Relabelling must not move the destination. A 'fix' that silently repointed a link would
    read as a success everywhere else in this file."""
    blob, _ = applied
    assert VAGUE_HREF in _part(blob.data, RELS2)


def test_the_run_is_still_a_hyperlink(applied):
    """In PresentationML the link is `<a:hlinkClick>` INSIDE `<a:rPr>`, and the writer rebuilds
    the run as rPr + a new `<a:t>`. A writer that dropped the run properties would produce
    perfectly correct visible text and leave the reader with no link at all — which no other
    assertion here would catch."""
    blob, _ = applied
    assert "hlinkClick" in _part(blob.data, SLIDE2)


def test_the_authors_formatting_survives(applied):
    blob, _ = applied
    assert 'b="1"' in _part(blob.data, SLIDE2), "the bold run property was dropped"


def test_the_other_hyperlink_and_both_titles_are_untouched(applied):
    blob, _ = applied
    assert GOOD_TEXT in _part(blob.data, SLIDE1)
    assert GOOD_HREF in _part(blob.data, "ppt/slides/_rels/slide1.xml.rels")
    assert TITLE1 in _part(blob.data, SLIDE1)
    assert TITLE2 in _part(blob.data, SLIDE2)


def test_the_saved_deck_still_opens(applied):
    """Through python-pptx, which had no part in writing the change."""
    from pptx import Presentation
    blob, _ = applied
    path = _spill(blob.data)
    assert zipfile.ZipFile(path).testzip() is None
    prs = Presentation(path)
    assert len(prs.slides) == 2
    assert any(sh.has_text_frame and APPROVED_TEXT in sh.text_frame.text
               for sh in prs.slides[1].shapes)


def test_a_second_real_assessment_no_longer_reports_2_4_4(applied):
    """THE claim: a fresh assessment of the SAVED bytes, not the writer's return value."""
    blob, _ = applied
    assert "2.4.4" not in _assess(blob.data)


def test_the_row_is_credited_and_the_copy_is_stored(applied):
    """Downstream of the re-scan, not of the approval: `_apply_one_value_kind` credits only
    when the residual no longer holds the criteria."""
    blob, store = applied
    assert store.count_unapplied_approved_values(SID, FILE) == 0
    assert blob.uploads, "the corrected copy was never stored"


# ── 3. where the lane must NOT credit ────────────────────────────────────────

# THE BYTES AN UNCREDITED LANE WROTE ARE DISCARDED, which shapes both controls below. When the
# re-scan still reports the criterion, `_apply_one_value_kind` returns the ORIGINAL `working`
# bytes and nothing is uploaded — by design: an unverified copy must never become the corrected
# one. So `blob.data` after a withheld credit is the deck as it went in, and asserting the
# criterion on it would be asserting something about the original.
#
# Each control therefore establishes what the write PRODUCES by calling the writer directly
# first — that half is a fact about the document — and then runs the real lane on the same
# approval to establish what it does with it. Both halves are needed: without the first, "no
# credit" could mean the writer simply failed.

def test_an_approved_value_that_is_itself_vague_is_not_credited(store, monkeypatch, original):
    """A reviewer approves text that still says nothing. The write succeeds; the re-scan still
    reports 2.4.4; nothing may be credited and nothing may be published.

    This is the test that distinguishes a real re-scan from `residual=set()`, which every other
    apply test in this repo supplies. Under that stub it would pass by construction.
    """
    from apply_link_text import apply_link_text
    written, applied, _ = apply_link_text(original, "pptx", {VAGUE_HREF: "read more"})
    assert applied, "the writer refused the value, so this control is not about crediting"
    assert "2.4.4" in _assess(written), (
        "the written text no longer fails 2.4.4 even to the first-party detector, so this "
        "control cannot distinguish a withheld credit from a cleared one")

    props = [p for p in _proposals(original) if p.get("locator") == VAGUE_HREF]
    blob = _Blob(original)
    _seed(store, props, ["read more"])
    _run_lane(monkeypatch, store, blob)

    assert store.count_unapplied_approved_values(SID, FILE) == 1
    assert store.mark_file_compliant_if_reviewed(SID, FILE) is False
    assert not blob.uploads, "an uncleared write was published as the corrected copy"
    assert blob.data == original, "an unverified copy replaced the stored one"


def test_a_value_that_trades_2_4_4_for_2_4_9_is_not_credited(store, monkeypatch, original):
    """The control this lane exists to make possible, and no single-criterion lane can express.

    The reviewer approves, for the vague link, the exact text the OTHER hyperlink already uses
    for a DIFFERENT destination. 2.4.4 genuinely clears — "supplier accessibility requirements"
    describes a destination — and the deck is worse: two links now read identically and go to
    different places, which is 2.4.9 Link Purpose (Link Only).

    Because pptx's link lane declares scs_to_clear = {2.4.4, 2.4.9}, `cleared` is false and the
    value stays uncredited. A lane that verified only the criterion the card came from would
    certify this deck.
    """
    from apply_link_text import apply_link_text
    written, applied, _ = apply_link_text(original, "pptx", {VAGUE_HREF: GOOD_TEXT})
    assert applied, "the writer refused the value, so this control is not about crediting"
    after = _assess(written)
    assert "2.4.4" not in after, (
        "the substituted text still fails 2.4.4, so this control is testing the vague-text "
        "path again rather than the criterion the write broke")
    assert "2.4.9" in after, "the duplicate text did not trip 2.4.9; the control is vacuous"

    props = [p for p in _proposals(original) if p.get("locator") == VAGUE_HREF]
    blob = _Blob(original)
    _seed(store, props, [GOOD_TEXT])
    _run_lane(monkeypatch, store, blob)

    assert store.count_unapplied_approved_values(SID, FILE) == 1, (
        "a write that fixed one criterion and broke another was credited")
    assert not blob.uploads
    assert blob.data == original


def test_an_already_descriptive_deck_is_left_byte_identical(store, monkeypatch):
    """Nothing is proposed for it, so the lane has nothing to write. 'We did not make it worse'
    is a claim a customer is entitled to."""
    clean = _deck(APPROVED_TEXT)
    assert [p for p in _proposals(clean) if p.get("locator") == VAGUE_HREF] == []
    blob = _Blob(clean)
    _seed(store, [], [])
    _run_lane(monkeypatch, store, blob)
    assert blob.data == clean and not blob.uploads


# ── 4. a broken engine earns nothing ─────────────────────────────────────────

@pytest.mark.parametrize("name,script,timeout", [
    ("cannot be launched", None, None),
    ("exits non-zero", "#!/bin/sh\necho boom >&2\nexit 9\n", None),
    ("writes unparseable output", '#!/bin/sh\nprintf "not json" > "$3"\nexit 3\n', None),
    ("hangs past the timeout", "#!/bin/sh\nsleep 30\n", "2"),
])
def test_a_broken_office_analyser_never_credits_this_lane_either(monkeypatch, original,
                                                                 name, script, timeout):
    """Re-asserted per lane rather than assumed to inherit: the fail-open #1058 closed lived in
    ONE shared seam (`scanner._analyse_office` → `verify_residual_scs` → `cleared = residual is
    None or …`), so a regression there takes every lane at once.

    The re-scan must return a real answer and never None, which `_apply_one_value_kind` reads
    as "cleared". It can, because the first-party pptx link check runs after the .NET call in
    its own try/except.
    """
    import stat as _stat

    import scanner
    if script is None:
        monkeypatch.setattr(scanner, "DOTNET", "/nonexistent/dotnet", raising=False)
    else:
        fake = Path(tempfile.mkdtemp()) / "dotnet"
        fake.write_text(script)
        fake.chmod(fake.stat().st_mode | _stat.S_IEXEC)
        monkeypatch.setattr(scanner, "DOTNET", str(fake), raising=False)
    if timeout:
        monkeypatch.setenv("ACP_OFFICE_CLI_TIMEOUT", timeout)

    from proposals import verify_residual_scs
    residual = verify_residual_scs(original, FILE)
    assert residual is not None, (
        f"an office CLI that {name} made the re-scan return None — every approved value on this "
        f"lane would be credited on a scan that never happened")
    assert "2.4.4" in residual

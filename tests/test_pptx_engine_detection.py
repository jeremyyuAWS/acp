"""Regression guard: PowerPoint tables with no header row are flagged for WCAG 1.3.1.

docx and xlsx both flag a data table that has no designated header row (Info and
Relationships). PowerPoint has parity: the vendored .NET analyser ships
`Pptx/Rules/TableHeaderRule.cs` (rule id PPTX-TABLE-001), registered in
ServiceCollectionExtensions, which flags an `<a:tbl>` with more than one row whose
`<a:tblPr firstRow="1">` is absent. This mirrors the xlsx/docx TableHeaderRule and is
conservative: a single-row table (the common layout-table shape) is skipped.

There is deliberately NO first-party detector in office_structure.py for this — adding
one would double-count. This test pins the .NET behaviour so the coverage can't silently
regress. It needs the built .NET CLI to run, so it self-skips when the analysers did not
run (CLI not built), exactly like tests/test_xlsx_engine_detection.py.
"""
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "api"))
import scanner  # noqa: E402


def _table_header_findings(fdict) -> list[dict]:
    """The 1.3.1 table-header findings emitted for this file."""
    return [i for i in fdict.get("issues", [])
            if i.get("ruleId") == "PPTX-TABLE-001"
            or str(i.get("wcag", "")) in ("SC_1_3_1", "1.3.1 Info and Relationships")]


def _engine_ran(fdict) -> bool:
    return any(str(i.get("wcag", "")).startswith("SC_") for i in fdict.get("issues", []))


def _build(d: Path, name: str, *, first_row: bool, rows: int) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank body, title placeholder
    gf = slide.shapes.add_table(rows, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    table = gf.table
    table.first_row = first_row                          # sets <a:tblPr firstRow="0|1">
    for ri in range(rows):
        for ci in range(2):
            table.cell(ri, ci).text = f"r{ri}c{ci}"
    prs.save(str(d / name))
    return name


@pytest.fixture(scope="module")
def scanned(tmp_path_factory):
    d = tmp_path_factory.mktemp("pptx-engine")
    out = {}
    for key, name, first_row, rows in [
        ("headerless", "no_header_data.pptx", False, 3),   # data table, no header -> fires
        ("with_header", "with_header.pptx", True, 3),       # header designated -> silent
        ("one_row", "one_row_layout.pptx", False, 1),       # 1-row layout table -> silent
    ]:
        _build(d, name, first_row=first_row, rows=rows)
        fdict, _ = scanner.analyse_and_assess(d, name, detect_pii=False)
        out[key] = fdict
    return out


def test_headerless_data_table_is_flagged(scanned):
    fdict = scanned["headerless"]
    if not _engine_ran(fdict):
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    findings = _table_header_findings(fdict)
    assert findings, "headerless pptx data table not flagged (Pptx TableHeaderRule regressed)"
    assert all(f.get("severity") == "SERIOUS" for f in findings)


def test_designated_header_row_is_not_flagged(scanned):
    fdict = scanned["with_header"]
    if not _engine_ran(fdict):
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    assert not _table_header_findings(fdict), "table with firstRow set should not be flagged"


def test_single_row_layout_table_is_not_flagged(scanned):
    fdict = scanned["one_row"]
    if not _engine_ran(fdict):
        pytest.skip("the .NET office analysers did not run (CLI not built)")
    assert not _table_header_findings(fdict), "1-row layout table should not be flagged (no false positive)"

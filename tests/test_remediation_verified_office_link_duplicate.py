"""The 2.4.9 link lane on docx AND pptx, proved end to end (Link Purpose (Link Only)).

Two lanes in one module because they are one writer, one store getter and one criterion — the
only difference is which OOXML part the hyperlink lives in, and parametrising over the format is
how that difference gets asserted rather than assumed. Same bar as the lanes before: the
original document trips the finding, an approval changes the saved document, a REAL re-scan
verifies it, unrelated content survives, and a broken engine earns no credit. Nothing but the
blob store is patched.

WHAT 2.4.9 IS, and why the fixture looks like it does. "Link Purpose (Link Only)" asks that a
link's own text identify its destination WITHOUT the surrounding sentence. So the failure here
is not vague text — both links below read perfectly well — it is the SAME text pointing at two
DIFFERENT destinations. A reader tabbing through the links hears "annual report" twice and
cannot tell which is which.

That makes this the mirror image of the 2.4.4 proofs, and the mirror is the point. Because the
Office link lane declares `scs_to_clear = ("2.4.4", "2.4.9")` for both formats, a value that
clears one criterion by breaking the other must not be credited — and it can be constructed from
either side. `test_remediation_verified_pptx_link.py` builds the 2.4.4 → 2.4.9 direction (a
vague link relabelled with text another link already uses).
`test_a_value_that_trades_2_4_9_for_2_4_4_is_not_credited` below builds the 2.4.9 → 2.4.4 one: a
duplicate relabelled to "click here" genuinely stops being a duplicate, and is worse. Neither
direction is credited, and neither could be detected without a real re-scan.

WHAT THIS CLAIMS: ACP's own criteria stop firing on a document ACP changed, the destinations are
unmoved, and the rest survived. NOT that the new label is the one a human would choose.
"""
from __future__ import annotations

import io
import sys
import tempfile
import zipfile
from pathlib import Path

import pytest

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))
sys.path.insert(0, str(ACP / "scripts"))

SID = "rv-office-249"

Q3_HREF = "https://example.org/q3-report"
Q4_HREF = "https://example.org/q4-report"
DUPLICATE = "annual report"          # descriptive, and used for BOTH destinations — the failure
APPROVED = "Q4 annual report"        # what the reviewer settles on for the Q4 link
HEADING = "Quarterly reports"
BODY = "Unrelated body copy that must survive the write."


# ── the two documents ─────────────────────────────────────────────────────────

def _docx(q4_text: str) -> bytes:
    from docx import Document
    from gen_demo_fixtures import _add_hyperlink

    d = Document()
    d.add_heading(HEADING, level=1)
    p1 = d.add_paragraph("Third quarter: ")
    _add_hyperlink(p1, Q3_HREF, DUPLICATE)
    p2 = d.add_paragraph("Fourth quarter: ")
    _add_hyperlink(p2, Q4_HREF, q4_text)
    d.add_paragraph(BODY)

    out = Path(tempfile.mkdtemp()) / "reports.docx"
    d.save(out)
    return out.read_bytes()


def _pptx(q4_text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for title, href, text in (("Third quarter", Q3_HREF, DUPLICATE),
                              ("Fourth quarter", Q4_HREF, q4_text)):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
        r = s.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1)).text_frame \
            .paragraphs[0].add_run()
        r.text = text
        r.hyperlink.address = href
    box = prs.slides[1].shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(1))
    box.text_frame.text = BODY

    out = Path(tempfile.mkdtemp()) / "reports.pptx"
    prs.save(out)
    return out.read_bytes()


BUILD = {"docx": _docx, "pptx": _pptx}
FILENAME = {"docx": "reports.docx", "pptx": "reports.pptx"}
# The part carrying the SECOND link's text, per format — the only structural difference between
# the two lanes, named here rather than branched on in every assertion.
LINK_PART = {"docx": "word/document.xml", "pptx": "ppt/slides/slide2.xml"}
RELS_PART = {"docx": "word/_rels/document.xml.rels", "pptx": "ppt/slides/_rels/slide2.xml.rels"}
# Where the OTHER link lives. On docx both hyperlinks share one part; on pptx they are on
# different slides, which is the structural difference this module parametrises over.
SIBLING_PART = {"docx": "word/document.xml", "pptx": "ppt/slides/slide1.xml"}
SIBLING_RELS = {"docx": "word/_rels/document.xml.rels",
                "pptx": "ppt/slides/_rels/slide1.xml.rels"}
# Content the write must not touch, and the part it is in. The heading is a docx-only
# construct — the deck's equivalent is its slide titles — so each format names its own.
SURVIVES = {
    "docx": ((HEADING, "word/document.xml"), (BODY, "word/document.xml")),
    "pptx": (("Third quarter", "ppt/slides/slide1.xml"),
             ("Fourth quarter", "ppt/slides/slide2.xml"),
             (BODY, "ppt/slides/slide2.xml")),
}


@pytest.fixture(params=("docx", "pptx"))
def fmt(request):
    if request.param == "docx":
        pytest.importorskip("docx")
    else:
        pytest.importorskip("pptx")
    return request.param


# ── helpers ───────────────────────────────────────────────────────────────────

def _assess(data: bytes, fmt: str) -> set[str]:
    """The SCs a REAL assessment reports — the same call the production re-verification makes."""
    from assessment_policy import _extract_sc
    from scanner import analyse_and_assess
    name = FILENAME[fmt]
    with tempfile.TemporaryDirectory() as d:
        (Path(d) / name).write_bytes(data)
        fd, _ = analyse_and_assess(Path(d), name, detect_pii=False)
    return {sc for i in (fd or {}).get("issues", []) if (sc := _extract_sc(i.get("wcag", "")))}


def _spill(data: bytes, fmt: str) -> str:
    p = Path(tempfile.mkdtemp()) / FILENAME[fmt]
    p.write_bytes(data)
    return str(p)


def _part(data: bytes, name: str) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        return z.read(name).decode("utf-8")


def _proposals(data: bytes, fmt: str) -> list[dict]:
    from proposals import propose_link_texts
    return propose_link_texts(_spill(data, fmt), fmt, ai_enabled=False)


class _Blob:
    def __init__(self, data: bytes):
        self.data, self.uploads = data, []

    def download_remediated(self, owner, sid, f):
        return self.data

    def upload_remediated(self, owner, sid, f, data, mime):
        self.data = data
        self.uploads.append((f, mime))
        return "http://b/2"


@pytest.fixture()
def store(monkeypatch):
    import store as store_mod
    monkeypatch.setattr(store_mod, "_SQLITE_PATH", Path(tempfile.mkdtemp()) / "rv.db")
    return store_mod.Store()


def _seed(store, fmt: str, props: list[dict], values: list[str]) -> int:
    """A scanned + remediated file with one 2.4.9 card, and `values` approved on it."""
    name = FILENAME[fmt]
    store.init_scan_run(SID, "drive", 1, "2026-08-31T00:00:00Z", "rubric", "hash")
    store.save_file_result(SID, {
        "file": name, "engine": "office", "status": "pass", "score": 60, "compliant": 0,
        "skipped_rules": 0, "drive_file_id": "d1",
        "issues": [{"ruleId": f"{fmt.upper()}_LINK_PURPOSE_AMBIGUOUS",
                    "wcag": "2.4.9 Link Purpose (Link Only)", "severity": "MODERATE"}],
    }, "2026-08-31T00:00:00Z")
    store.record_remediation(SID, name, drive_write_url="http://d/1", blob_url="http://b/1")
    item_id = store.enqueue_proposals(SID, name, "2.4.9", [
        {k: p.get(k) for k in ("locator", "before", "proposed_value", "rationale", "source")}
        for p in props], rule_name="Link Purpose (Link Only)")
    store.update_hitl_item(item_id, "approved", None, None)
    store.approve_proposal_values(item_id, values)
    return item_id


def _run_lane(monkeypatch, store, blob, fmt: str):
    """The production handler, with the re-scan UNPATCHED."""
    import core
    import handlers
    monkeypatch.setattr(core, "store", store)
    monkeypatch.setitem(sys.modules, "blob", blob)
    handlers._apply_approved_values({"scan_id": SID, "file": FILENAME[fmt]}, {})


# ── 1. the finding and its proposal ──────────────────────────────────────────

def test_a_real_assessment_reports_2_4_9_on_reused_link_text(fmt):
    assert "2.4.9" in _assess(BUILD[fmt](DUPLICATE), fmt)


def test_the_original_does_not_fail_2_4_4(fmt):
    """Both labels are descriptive; only their REUSE is the failure.

    Load-bearing for the whole module. If the fixture also failed 2.4.4, no assertion below
    could tell which criterion the write moved, and
    `test_a_value_that_trades_2_4_9_for_2_4_4_is_not_credited` could not exist at all.
    """
    assert "2.4.4" not in _assess(BUILD[fmt](DUPLICATE), fmt)


def test_the_distinct_text_control_is_not_flagged(fmt):
    """The same document with the second link relabelled. Without this, a detector that fired on
    every pair of hyperlinks would satisfy the test above and nothing would notice."""
    assert "2.4.9" not in _assess(BUILD[fmt](APPROVED), fmt)


def test_the_proposer_offers_a_value_for_each_side_of_the_collision(fmt):
    """2.4.9 is a property of a PAIR, so a proposal naming only one of them leaves the reviewer
    guessing which to change. Both destinations must be offered, each carrying sc 2.4.9."""
    props = _proposals(BUILD[fmt](DUPLICATE), fmt)
    by_href = {p["locator"]: p for p in props}
    assert set(by_href) == {Q3_HREF, Q4_HREF}, sorted(by_href)
    for p in by_href.values():
        assert p["before"] == DUPLICATE
        assert p["proposed_value"].strip()
        assert p.get("sc") == "2.4.9"


# ── 2. approval → write → re-scan → credit, through the real path ────────────

@pytest.fixture()
def applied(store, monkeypatch, fmt):
    original = BUILD[fmt](DUPLICATE)
    props = [p for p in _proposals(original, fmt) if p["locator"] == Q4_HREF]
    blob = _Blob(original)
    _seed(store, fmt, props, [APPROVED])
    _run_lane(monkeypatch, store, blob, fmt)
    return blob, store, fmt, original


def test_relabelling_one_side_is_enough_and_the_document_says_so(applied):
    """Only ONE of the two links needs to change: the criterion is about the collision, not
    about either label. The other link keeps its original text, asserted so a writer that
    rewrote both would fail rather than pass twice as hard."""
    blob, _, fmt, _ = applied
    assert APPROVED in _part(blob.data, LINK_PART[fmt])
    assert DUPLICATE in _part(blob.data, SIBLING_PART[fmt])


def test_both_destinations_are_unchanged(applied):
    blob, _, fmt, _ = applied
    assert Q4_HREF in _part(blob.data, RELS_PART[fmt])
    assert Q3_HREF in _part(blob.data, SIBLING_RELS[fmt])


def test_unrelated_content_survives(applied):
    blob, _, fmt, _ = applied
    for text, part in SURVIVES[fmt]:
        assert text in _part(blob.data, part), f"{text!r} did not survive the write"


def test_the_saved_file_still_opens(applied):
    """Through the library that had no part in writing the change."""
    blob, _, fmt, _ = applied
    path = _spill(blob.data, fmt)
    assert zipfile.ZipFile(path).testzip() is None
    if fmt == "docx":
        from docx import Document
        assert any(HEADING in p.text for p in Document(path).paragraphs)
    else:
        from pptx import Presentation
        assert len(Presentation(path).slides) == 2


def test_a_second_real_assessment_no_longer_reports_2_4_9(applied):
    """THE claim: a fresh assessment of the SAVED bytes, not the writer's return value."""
    blob, _, fmt, _ = applied
    assert "2.4.9" not in _assess(blob.data, fmt)


def test_the_row_is_credited_and_the_copy_is_stored(applied):
    blob, store, fmt, _ = applied
    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 0
    assert blob.uploads


# ── 3. where the lane must NOT credit ────────────────────────────────────────
#
# As in the 2.4.4 proofs: the bytes an uncredited lane wrote are DISCARDED
# (`_apply_one_value_kind` returns the original `working` and uploads nothing), so each control
# establishes what the write produces by calling the writer directly, then runs the real lane on
# the same approval. Without the first half, "no credit" could mean the writer simply failed.

def test_a_value_that_trades_2_4_9_for_2_4_4_is_not_credited(store, monkeypatch, fmt):
    """The mirror of the pptx 2.4.4 control, built from the other side.

    The reviewer relabels the duplicate to "click here". 2.4.9 genuinely clears — the two links
    no longer read alike — and the document is worse: one of them now says nothing about where
    it goes. Because the Office link lane declares scs_to_clear = {2.4.4, 2.4.9}, `cleared` is
    false and nothing is credited.

    A lane verifying only the criterion its card came from would certify this document, and no
    amount of checking the WRITE could tell: the write did exactly what was approved.
    """
    from apply_link_text import apply_link_text
    original = BUILD[fmt](DUPLICATE)

    written, ap, _ = apply_link_text(original, fmt, {Q4_HREF: "click here"})
    assert ap, "the writer refused the value, so this control is not about crediting"
    after = _assess(written, fmt)
    assert "2.4.9" not in after, "the relabelled link is still a duplicate; the control is vacuous"
    assert "2.4.4" in after, "\"click here\" did not trip 2.4.4; the control is vacuous"

    props = [p for p in _proposals(original, fmt) if p["locator"] == Q4_HREF]
    blob = _Blob(original)
    _seed(store, fmt, props, ["click here"])
    _run_lane(monkeypatch, store, blob, fmt)

    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 1, (
        "a write that fixed one criterion and broke another was credited")
    assert store.mark_file_compliant_if_reviewed(SID, FILENAME[fmt]) is False
    assert not blob.uploads
    assert blob.data == original


def test_relabelling_one_link_to_the_other_s_text_leaves_the_collision(store, monkeypatch, fmt):
    """The degenerate approval: the reviewer 'fixes' the Q4 link by typing the text it already
    shares with Q3. The write succeeds and changes nothing that matters, so 2.4.9 still fails
    and nothing is credited."""
    from apply_link_text import apply_link_text
    original = BUILD[fmt](DUPLICATE)

    written, ap, _ = apply_link_text(original, fmt, {Q4_HREF: DUPLICATE.upper()})
    assert ap
    assert "2.4.9" in _assess(written, fmt), (
        "case alone separated the two labels, so this control is not about the collision")

    props = [p for p in _proposals(original, fmt) if p["locator"] == Q4_HREF]
    blob = _Blob(original)
    _seed(store, fmt, props, [DUPLICATE.upper()])
    _run_lane(monkeypatch, store, blob, fmt)

    assert store.count_unapplied_approved_values(SID, FILENAME[fmt]) == 1
    assert not blob.uploads


def test_a_document_with_distinct_link_text_is_left_byte_identical(store, monkeypatch, fmt):
    clean = BUILD[fmt](APPROVED)
    assert _proposals(clean, fmt) == []
    blob = _Blob(clean)
    _seed(store, fmt, [], [])
    _run_lane(monkeypatch, store, blob, fmt)
    assert blob.data == clean and not blob.uploads


# ── 4. a broken engine earns nothing ─────────────────────────────────────────

@pytest.mark.parametrize("name,script,timeout", [
    ("cannot be launched", None, None),
    ("exits non-zero", "#!/bin/sh\necho boom >&2\nexit 9\n", None),
    ("hangs past the timeout", "#!/bin/sh\nsleep 30\n", "2"),
])
def test_a_broken_office_analyser_never_credits_this_lane_either(monkeypatch, fmt,
                                                                 name, script, timeout):
    """Re-asserted per lane rather than assumed to inherit: the fail-open #1058 closed lived in
    ONE shared seam, so a regression there takes every lane at once."""
    import stat as _stat

    import scanner
    if script is None:
        monkeypatch.setattr(scanner, "DOTNET", "/nonexistent/dotnet", raising=False)
    else:
        fake = Path(tempfile.mkdtemp()) / "dotnet"
        fake.write_text(script)
        fake.chmod(fake.stat().st_mode | _stat.S_IEXEC)
        monkeypatch.setattr(scanner, "DOTNET", str(fake), raising=False)
    if timeout:
        monkeypatch.setenv("ACP_OFFICE_CLI_TIMEOUT", timeout)

    from proposals import verify_residual_scs
    residual = verify_residual_scs(BUILD[fmt](DUPLICATE), FILENAME[fmt])
    assert residual is not None, (
        f"an office CLI that {name} made the re-scan return None — every approved value on this "
        f"lane would be credited on a scan that never happened")
    assert "2.4.9" in residual

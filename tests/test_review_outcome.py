"""Review-aware per-criterion outcome (ADR 0023, Option A): store._rule_outcome +
REVIEW_FORMATS + _split_sc_counts.

A review-lane (criterion, format) resolves to REVIEW when its detector fires an advisory
(severity REVIEW) finding and to NOT_EVALUATED when it does not — NEVER to PASS, because a
review detector doesn't certify conformance (ADR 0016). A blocking FAIL always outranks an
advisory REVIEW. Pass/fail-lane criteria are unaffected unless an advisory finding rides them.
"""
from __future__ import annotations

import sys
from pathlib import Path

ACP = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ACP / "api"))

import store  # noqa: E402


# ── review-lane resolution (4.1.2 is a catalog rule; office is its review lane) ──
def test_review_lane_no_signal_is_not_evaluated_never_pass():
    # A control-free office doc: no advisory finding → genuine N/A, NOT a fabricated pass.
    assert store._rule_outcome("4.1.2", "docx", 0, 0) == store.NOT_EVALUATED
    assert store._rule_outcome("2.1.2", "pptx", 0, 0) == store.NOT_EVALUATED


def test_review_lane_with_signal_is_review():
    assert store._rule_outcome("4.1.2", "docx", 0, 1) == store.REVIEW
    assert store._rule_outcome("2.1.2", "xlsx", 0, 3) == store.REVIEW


def test_review_lane_blocking_finding_outranks_review():
    # If a definite FAIL also landed on a review-lane criterion, FAIL wins.
    assert store._rule_outcome("4.1.2", "docx", 1, 2) == "FAIL"


def test_review_lane_is_format_scoped():
    # 4.1.2's REVIEW_FORMATS lane is office-only, and pdf is not in it. pdf reaches REVIEW by a
    # different route: it is registered in the capability registry with PARTIAL coverage (the
    # AcroForm technique in formats/pdf/detectors/name_role_value.py), and partial coverage
    # cannot certify a pass. Same token, different reason — which is the point of the coverage
    # axis. It read NOT_EVALUATED before the registry, despite the detector having shipped.
    assert store._rule_outcome("4.1.2", "pdf", 0, 0) == store.REVIEW
    # The two routes to REVIEW differ in when they fire, and that difference is real. A
    # REVIEW_FORMATS lane needs a signal — its detector surfaces evidence or says nothing — so
    # docx with no findings is still "we did not look". A registry lane with PARTIAL coverage
    # reports REVIEW on a CLEAN scan, because the technique ran and covered part of the
    # criterion. Same token; one is "found something to look at", the other "looked partially".
    assert store._rule_outcome("4.1.2", "docx", 0, 0) == store.NOT_EVALUATED
    assert store._rule_outcome("4.1.2", "docx", 0, 1) == store.REVIEW
    # A format with neither a review lane nor a registry entry still reads "we did not look".
    assert store._rule_outcome("2.4.3", "xlsx", 0, 0) == store.NOT_EVALUATED
    # html keeps its real pass/fail lane for 4.1.2.
    assert store._rule_outcome("4.1.2", "html", 0, 0) == "PASS"
    assert store._rule_outcome("4.1.2", "html", 1, 0) == "FAIL"


# ── pass/fail lane still behaves for 🟢 auto-assess criteria ────────────────────
def test_pass_fail_lane_unchanged_for_auto_assess():
    # 3.1.1 Language is 🟢 auto-assess (deterministic present/absent) → a clean scan is a real PASS.
    assert store._rule_outcome("3.1.1", "pdf", 0, 0) == "PASS"
    assert store._rule_outcome("3.1.1", "pdf", 2, 0) == "FAIL"
    assert store._rule_outcome("1.3.4", "docx", 0, 0) == store.NOT_EVALUATED  # html-only rule


def test_review_lane_criterion_clean_is_review_not_a_certified_pass():
    # 1.1.1 is a 🟡 review-lane criterion (alt adequacy is a judgement) — a scan that finds no
    # missing alt is NOT a certified pass; it stays REVIEW ("verify"), never green (audit #174).
    assert store._rule_outcome("1.1.1", "pdf", 0, 0) == store.REVIEW
    assert store._rule_outcome("1.3.3", "docx", 0, 0) == store.REVIEW
    # …but a real FAIL still wins over the review lane.
    assert store._rule_outcome("1.1.1", "pdf", 2, 0) == "FAIL"


def test_docx_2_4_6_clean_is_review_not_a_certified_pass():
    """2.4.6 on docx used to be the one 🟢 outlier and resolved a clean scan to PASS.

    Its only detector, DOCX_HEADING_SKIP, judges heading LEVELS — so a document with a flawless
    H1→H2→H3 outline whose headings read "Section 1" / "Untitled" / "asdf" scanned clean and was
    certified. That is a pass on evidence that never bore on the criterion (ADR 0023 audit,
    Correction 2). All four document formats now agree; a real FAIL still outranks the lane.
    """
    for fmt in ("docx", "pptx", "xlsx", "pdf"):
        assert store._rule_outcome("2.4.6", fmt, 0, 0) == store.REVIEW, fmt
        assert store._rule_outcome("2.4.6", fmt, 1, 0) == "FAIL", fmt


def test_html_2_4_6_clean_is_review_not_a_certified_pass():
    """html carried the identical defect docx did, one detector over — fixed here.

    scanner.HTML_HEADING_SKIP judges heading LEVELS; remediate._fix_heading_skip clamps the gaps;
    the clean re-scan was read as certifying 2.4.6. Run through the real pipeline
    (_analyse_html → filter_issues_to_target → _rule_outcome), a page with a perfect h1→h2→h3
    outline headed "Section 1" / "Untitled" / "asdf" resolved PASS, and so did one starting at h3
    with no h1 or h2 (the `prev_level > 0` guard never judges the first heading). Neither says
    anything about whether a heading describes its topic. All five formats now agree.
    """
    assert store._rule_outcome("2.4.6", "html", 0, 0) == store.REVIEW
    assert store._rule_outcome("2.4.6", "html", 1, 0) == "FAIL"   # a real FAIL still outranks it


def test_html_2_4_6_meaningless_headings_scan_clean_but_do_not_certify():
    """The same claim through the REAL pipeline, not just the lane table.

    Both of these pages scan clean — the detector is a level check and has nothing to say about
    either — and both plainly fail 2.4.6. The outcome they resolve to must be REVIEW, so the
    honest statement ("we looked at levels; a human must judge the wording") is what reaches the
    report, instead of a green certified PASS.
    """
    import tempfile

    import scanner

    PAGES = {
        # A flawless h1→h2→h3 outline whose headings describe nothing at all.
        "meaningless": "<html><body><h1>Section 1</h1><h2>Untitled</h2><h3>asdf</h3></body></html>",
        # No h1 or h2 anywhere: the `prev_level > 0` guard means the first heading is never judged.
        "starts_at_h3": "<html><body><h3>Deep</h3><h4>Deeper</h4></body></html>",
    }
    for name, html in PAGES.items():
        tmp = Path(tempfile.mkdtemp()) / f"{name}.html"
        tmp.write_text(html)
        issues = store.filter_issues_to_target(scanner._analyse_html(tmp)["issues"], "AA")
        fails = [i for i in issues if i["wcag"].startswith("2.4.6")]
        assert not fails, f"{name} should scan clean for 2.4.6, got {fails}"
        assert store._rule_outcome("2.4.6", "html", len(fails), 0) == store.REVIEW, name


def test_no_non_certifying_lane_can_ever_return_pass():
    """The CLASS guard: only a 🟢 auto lane may certify a PASS. Every other lane, on every format.

    docx 2.4.6, html 2.4.6 and pptx 2.1.1 were the same defect found three times, one cell at a
    time, because each fix pinned its own instance. This asserts the invariant over the whole
    assessment table instead, so a new 🟡/🔴 cell cannot regress into a green PASS unnoticed —
    including a lane token that does not exist yet, since the rule is `!= auto` rather than a
    list of the non-certifying lanes (the omission that let 🔴 human through in the first place).
    """
    import remediation_capability as rc

    offenders = [
        (fmt, sc, lane)
        for fmt, scs in rc.assessment_table().items()
        for sc, lane in scs.items()
        if lane != rc.A_AUTO and store._rule_outcome(sc, fmt, 0, 0) == "PASS"
    ]
    assert not offenders, (
        "these (format, criterion) pairs certify a PASS their assessment lane says ACP "
        f"cannot judge: {offenders}"
    )


def test_auto_lanes_still_certify_so_the_guard_is_not_vacuous():
    """The guard above is a negative; this proves it did not simply delete certification.

    A blanket "never PASS" would satisfy the invariant and be worthless. 🟢 auto cells must still
    resolve a clean scan to a real PASS.
    """
    import remediation_capability as rc

    certifying = [(fmt, sc) for fmt, scs in rc.assessment_table().items() for sc, lane in scs.items()
                  if lane == rc.A_AUTO and store._rule_outcome(sc, fmt, 0, 0) == "PASS"]
    assert len(certifying) > 20, f"expected the 🟢 lane to still certify broadly, got {certifying}"
    # The named control from the bug report: pptx 1.4.3 is deterministic and must stay green.
    assert store._rule_outcome("1.4.3", "pptx", 0, 0) == "PASS"


def test_pptx_2_1_1_clean_deck_is_review_not_a_certified_pass():
    """pptx 2.1.1 Keyboard — the 🔴 human cell that resolved to a certified PASS.

    Third instance of the docx/html 2.4.6 shape, and the worst of the three: those two at least
    ran a detector whose narrow signal was over-read. There is NO 2.1.1 detector on any pptx path
    — office_structure.checks_for dispatches ten pptx checks and not one emits 2.1.1 — so bare
    RULE_FORMATS membership alone was being read as "the validator ran". A clean deck therefore
    scored green on keyboard operability, which is a property of the runtime presenting the deck
    and not of the file, and which the lane table in the same breath said only a person can judge.

    Run through the real pipeline (checks_for → filter_issues_to_target → _rule_outcome).
    """
    import tempfile

    import office_structure
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Accessibility Review"
    slide.placeholders[1].text = "Findings and next steps for the platform team."
    deck = Path(tempfile.mkdtemp()) / "clean.pptx"
    prs.save(str(deck))

    issues = store.filter_issues_to_target(office_structure.checks_for(deck, ".pptx"), "AA")
    fails = [i for i in issues if str(i.get("wcag", "")).startswith("2.1.1")]
    assert not fails, f"no pptx detector emits 2.1.1, so a clean deck must scan clean: {fails}"
    assert store._rule_outcome("2.1.1", "pptx", len(fails), 0) == store.REVIEW
    # A recorded blocking finding still outranks the lane — this is not a mute button.
    assert store._rule_outcome("2.1.1", "pptx", 1, 0) == "FAIL"
    # The lane table and the outcome function now agree, which is what the bug was about.
    import remediation_capability as rc
    assert rc.assessment_table()["pptx"]["2.1.1"] == rc.A_HUMAN


def test_advisory_review_on_a_pass_fail_criterion_surfaces_as_review():
    # An advisory finding on an in-scope pass/fail criterion (no blocking finding) → REVIEW.
    assert store._rule_outcome("1.1.1", "pdf", 0, 1) == store.REVIEW
    # …but a real FAIL still wins.
    assert store._rule_outcome("1.1.1", "pdf", 1, 1) == "FAIL"


def test_three_arg_call_is_back_compatible():
    # Existing callers pass 3 args; review_count defaults to 0.
    assert store._rule_outcome("3.1.1", "pdf", 0) == "PASS"           # 🟢 auto → pass
    assert store._rule_outcome("1.1.1", "pdf", 0) == store.REVIEW     # 🟡 review-lane → verify
    assert store._rule_outcome("4.1.2", "docx", 0) == store.NOT_EVALUATED


# ── count split by advisory severity ────────────────────────────────────────────
def test_split_sc_counts_separates_advisory_from_blocking():
    issues = [
        {"ruleId": "X", "wcag": "1.1.1 Non-text Content", "severity": "CRITICAL"},
        {"ruleId": "Y", "wcag": "1.1.1 Non-text Content", "severity": "SERIOUS"},
        {"ruleId": "OFFICE_INTERACTIVE_CONTROL_NAME_ROLE", "wcag": "4.1.2 Name, Role, Value", "severity": "REVIEW"},
        {"ruleId": "OFFICE_INTERACTIVE_CONTROL_KEYBOARD", "wcag": "2.1.2 No Keyboard Trap", "severity": "REVIEW"},
    ]
    fail, review = store._split_sc_counts(issues)
    assert fail == {"1.1.1": 2}
    assert review == {"4.1.2": 1, "2.1.2": 1}


def test_split_sc_counts_ignores_findings_with_no_sc():
    fail, review = store._split_sc_counts([{"ruleId": "Z", "wcag": "not-a-criterion", "severity": "MINOR"}])
    assert fail == {} and review == {}


def test_review_is_a_distinct_outcome_token():
    assert store.REVIEW == "REVIEW"
    assert store.REVIEW not in (store.NOT_EVALUATED, "PASS", "FAIL")

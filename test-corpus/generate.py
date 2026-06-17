#!/usr/bin/env python3
"""
Generate a synthetic WCAG-2.1 smoke-test corpus for the acp scanner.

Files are crafted with controlled pass/fail signals matched to the devSEAL
analyser rules (verified against their source), plus robustness edge cases
(corrupt / wrong-type / encrypted). Emits a ground-truth manifest so a scan
can be checked against known expected findings — the corpus doubles as an oracle.

Run:
  uv run --with python-docx --with python-pptx --with openpyxl \
         --with reportlab --with pikepdf --with Pillow python generate.py
Outputs to ./files/ and ./manifest.json (+ base64 sidecars in ./files/*.b64).
"""
from __future__ import annotations
import base64, json, os, struct, zlib
from pathlib import Path

OUT = Path(__file__).parent / "files"
OUT.mkdir(exist_ok=True)
manifest = []

def record(name, fmt, quality, outcome, rules, notes):
    manifest.append({"file": name, "format": fmt, "quality": quality,
                     "expected_outcome": outcome, "expected_rule_ids": rules, "notes": notes})

def tiny_png(path, color=(40, 60, 120)):
    from PIL import Image
    Image.new("RGB", (8, 8), color).save(path)

IMG = OUT / "_img.png"; tiny_png(IMG)

# ─────────────────────────── DOCX ───────────────────────────
from docx import Document
from docx.shared import Inches
from docx.oxml.shared import OxmlElement, qn

def add_hyperlink(p, text, url):
    part = p.part
    rid = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    h = OxmlElement("w:hyperlink"); h.set(qn("r:id"), rid)
    r = OxmlElement("w:r"); r.append(OxmlElement("w:rPr"))
    t = OxmlElement("w:t"); t.text = text; r.append(t); h.append(r); p._p.append(h)

def mark_header_row(table):
    trPr = table.rows[0]._tr.get_or_add_trPr()
    trPr.append(OxmlElement("w:tblHeader"))

def docx_compliant():
    d = Document()
    d.core_properties.title = "Q3 2024 Accessibility Report"
    d.core_properties.language = "en-US"
    d.add_heading("Q3 2024 Accessibility Report", 0)
    d.add_heading("Summary", 1)
    d.add_paragraph("Revenue rose 12% quarter over quarter.")
    pic = d.add_paragraph().add_run().add_picture(str(IMG), width=Inches(1))
    pic._inline.docPr.set("descr", "Bar chart showing Q3 revenue rising 12 percent.")
    p = d.add_paragraph("See the ")
    add_hyperlink(p, "full Q3 revenue methodology (PDF)", "https://example.com/q3")
    t = d.add_table(rows=2, cols=2); mark_header_row(t)
    t.cell(0,0).text="Region"; t.cell(0,1).text="Revenue"; t.cell(1,0).text="West"; t.cell(1,1).text="100"
    d.save(OUT/"docx-compliant.docx")
    record("docx-compliant.docx","docx","compliant","analysed",[],"title+lang+alt+header row+descriptive link; expect ~0 issues, high score")

def docx_moderate():
    d = Document()
    d.core_properties.title = "Team Offsite Notes"
    d.core_properties.language = "en-US"
    d.add_heading("Team Offsite Notes", 0)
    d.add_paragraph().add_run().add_picture(str(IMG), width=Inches(1))  # no alt
    p = d.add_paragraph("For the agenda ")
    add_hyperlink(p, "click here", "https://example.com/agenda")        # generic
    d.save(OUT/"docx-moderate.docx")
    record("docx-moderate.docx","docx","moderate","analysed",
           ["DOCX-ALT-001","DOCX-LINK-001"],"title+lang OK; missing alt + generic link text")

def docx_noncompliant():
    d = Document()  # no title, no language
    d.add_paragraph("Quarterly figures and notes.")
    d.add_paragraph().add_run().add_picture(str(IMG), width=Inches(1))  # no alt
    p = d.add_paragraph("More info ")
    add_hyperlink(p, "here", "https://example.com")                    # generic
    t = d.add_table(rows=3, cols=2)                                    # no header row
    for i,(a,b) in enumerate([("R","V"),("West","100"),("East","90")]):
        t.cell(i,0).text=a; t.cell(i,1).text=b
    d.save(OUT/"docx-noncompliant.docx")
    record("docx-noncompliant.docx","docx","non-compliant","analysed",
           ["DOCX-TITLE-001","DOCX-LANG-001","DOCX-ALT-001","DOCX-LINK-001","DOCX-TABLE-001"],
           "no title/lang, missing alt, generic link, table w/o header — many SERIOUS/CRITICAL")

def docx_empty_clean():
    d = Document()
    d.core_properties.title = "Blank But Tagged"; d.core_properties.language = "en-US"
    d.save(OUT/"docx-empty.docx")
    record("docx-empty.docx","docx","edge-empty","analysed",[],"empty body but title+lang set; tests empty-content handling, expect ~0 issues")

# ─────────────────────────── PPTX ───────────────────────────
from pptx import Presentation
from pptx.util import Inches as PInches

def pptx_compliant():
    prs = Presentation(); prs.core_properties.language = "en-US"; prs.core_properties.title="Onboarding"
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s.shapes.title.text = "Onboarding Overview"
    pic = s.shapes.add_picture(str(IMG), PInches(1), PInches(2), PInches(1), PInches(1))
    pic._element.nvPicPr.cNvPr.set("descr", "Diagram of the four onboarding stages.")
    prs.save(OUT/"pptx-compliant.pptx")
    record("pptx-compliant.pptx","pptx","compliant","analysed",[],"slide title + image alt + language; expect ~0 issues")

def pptx_noncompliant():
    prs = Presentation()  # no language
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank — no title placeholder
    tb = s.shapes.add_textbox(PInches(1), PInches(1), PInches(4), PInches(1)); tb.text_frame.text="Big News"
    pic = s.shapes.add_picture(str(IMG), PInches(1), PInches(3), PInches(1), PInches(1))
    pic._element.nvPicPr.cNvPr.attrib.pop("descr", None)  # python-pptx auto-fills descr=filename; drop it for a true 'no alt'
    prs.save(OUT/"pptx-noncompliant.pptx")
    record("pptx-noncompliant.pptx","pptx","non-compliant","analysed",
           ["PPTX-TITLE-001","PPTX-ALT-001","PPTX-LANG-001"],"no slide title, image w/o alt, no language (rule ids approximate)")

# ─────────────────────────── XLSX ───────────────────────────
from openpyxl import Workbook

def xlsx_compliant():
    wb = Workbook(); wb.properties.title="Q3 Revenue"; wb.properties.language="en-US"
    ws = wb.active; ws.title="Q3 Revenue"
    ws.append(["Region","Revenue"]); ws.append(["West",100]); ws.append(["East",90])
    wb.save(OUT/"xlsx-compliant.xlsx")
    record("xlsx-compliant.xlsx","xlsx","compliant","analysed",[],"named sheet + title + lang + header row; expect ~0 issues")

def xlsx_noncompliant():
    wb = Workbook()  # default sheet name "Sheet", no title/lang
    ws = wb.active   # name stays "Sheet"
    ws.append(["Q3 Financials","",""]); ws.merge_cells("A1:C1")       # merged cells
    ws.append(["West",100,"x"])
    hidden = wb.create_sheet("Hidden"); hidden.sheet_state="hidden"   # hidden content
    wb.save(OUT/"xlsx-noncompliant.xlsx")
    record("xlsx-noncompliant.xlsx","xlsx","non-compliant","analysed",
           ["XLSX-TITLE-001","XLSX-LANG-001","XLSX-SHEETNAME-001","XLSX-MERGED-001","XLSX-HIDDEN-001"],
           "no title/lang, generic 'Sheet' name, merged cells, hidden sheet (rule ids approximate)")

# ─────────────────────────── PDF ───────────────────────────
from reportlab.pdfgen import canvas
import pikepdf

def pdf_untagged():
    p = OUT/"pdf-untagged.pdf"
    c = canvas.Canvas(str(p)); c.drawString(72,720,"Quarterly report — untagged, no title, no language."); c.showPage(); c.save()
    record("pdf-untagged.pdf","pdf","non-compliant","analysed",
           ["PDF-TAGGED-001","PDF-TITLE-001","PDF-LANG-001"],"untagged + no title + no /Lang (rule ids approximate)")

def pdf_titled_lang():
    p = OUT/"pdf-titled-lang.pdf"
    c = canvas.Canvas(str(p)); c.setTitle("Q3 Accessibility Report"); c.drawString(72,720,"Has title + language, still untagged."); c.showPage(); c.save()
    pdf = pikepdf.open(str(p), allow_overwriting_input=True)
    pdf.Root.Lang = pikepdf.String("en-US"); pdf.save(str(p))
    record("pdf-titled-lang.pdf","pdf","moderate","analysed",
           ["PDF-TAGGED-001"],"title + /Lang set, but still untagged — title/lang pass, tagging fails")

def pdf_encrypted():
    p = OUT/"pdf-encrypted.pdf"; tmp = OUT/"_plain.pdf"
    c = canvas.Canvas(str(tmp)); c.drawString(72,720,"secret"); c.showPage(); c.save()
    pdf = pikepdf.open(str(tmp))
    pdf.save(str(p), encryption=pikepdf.Encryption(owner="ownerpw", user="userpw"))
    tmp.unlink()
    record("pdf-encrypted.pdf","pdf","edge-encrypted","error",[],"password-protected; analyser open should FAIL gracefully (Succeeded=false, error captured)")

# ─────────────────── robustness / malformed ───────────────────
def edge_corrupt_docx():
    (OUT/"edge-corrupt.docx").write_bytes(b"PK\x03\x04" + b"\x00"*40 + b"truncated-not-a-real-zip")
    record("edge-corrupt.docx","docx","edge-corrupt","error",[],"truncated/invalid OOXML zip; Open() should fail gracefully")

def edge_plaintext_docx():
    (OUT/"edge-plaintext.docx").write_bytes(b"This is plain text pretending to be a Word document.\n")
    record("edge-plaintext.docx","docx","edge-wrongtype","error",[],"plain text with .docx extension; should fail gracefully")

def edge_pdf_as_docx():
    c = canvas.Canvas(str(OUT/"_x.pdf")); c.drawString(72,720,"pdf"); c.showPage(); c.save()
    (OUT/"edge-pdf-as.docx").write_bytes((OUT/"_x.pdf").read_bytes()); (OUT/"_x.pdf").unlink()
    record("edge-pdf-as.docx","docx","edge-wrongtype","error",[],"real PDF bytes with .docx extension; should fail gracefully")

for fn in [docx_compliant, docx_moderate, docx_noncompliant, docx_empty_clean,
           pptx_compliant, pptx_noncompliant, xlsx_compliant, xlsx_noncompliant,
           pdf_untagged, pdf_titled_lang, pdf_encrypted,
           edge_corrupt_docx, edge_plaintext_docx, edge_pdf_as_docx]:
    try: fn()
    except Exception as e: print(f"!! {fn.__name__} failed: {e}")

if IMG.exists(): IMG.unlink()

# base64 sidecars + manifest + size report
print(f"{'file':32} {'size':>8}")
for m in manifest:
    f = OUT/m["file"]
    if not f.exists(): print(f"MISSING {m['file']}"); continue
    b = f.read_bytes(); (OUT/(m['file']+'.b64')).write_text(base64.b64encode(b).decode())
    m["bytes"] = len(b); m["mime"] = {
        "docx":"application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx":"application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx":"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pdf":"application/pdf"}[m["format"]]
    print(f"{m['file']:32} {len(b):>8}")
(Path(__file__).parent/"manifest.json").write_text(json.dumps(manifest, indent=2))
print(f"\n{len(manifest)} files; manifest.json written")

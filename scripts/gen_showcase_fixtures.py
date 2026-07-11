#!/usr/bin/env python3
"""Generate IMAGE-RICH accessibility showcase files for a live end-to-end demo.

Unlike gen_demo_fixtures.py (solid-colour placeholder images, tuned to trip every detector),
these carry real charts/diagrams drawn with PIL so the GPU vision model (llava:13b) produces
MEANINGFUL alt text — the point being to show the "AI proposes, human confirms" lane with
content worth describing. Each file also seeds a broad spread of in-scope findings so a live
scan → assess → remediate exercises the whole pipeline.

Usage:  python scripts/gen_showcase_fixtures.py [OUT_DIR]   # default: ~/Downloads/acp-showcase
Then upload the files to the Google Drive you sign in with, and run Discover → Assess → Remediate.
"""
from __future__ import annotations

import io
import math
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

BLUE = (31, 95, 168); ORANGE = (214, 110, 40); GREEN = (52, 140, 90)
PURPLE = (140, 70, 160); RED = (200, 60, 70); GREY = (90, 90, 96)
PALETTE = [BLUE, ORANGE, GREEN, PURPLE, RED]

FRENCH = ("Le rapport trimestriel montre une croissance significative des revenus dans "
          "toutes les régions européennes cette année.")
SENSORY = ("To approve the request, click the large round green button in the bottom-right "
           "corner of the screen.")
DENSE = ("Notwithstanding the aforementioned contractual stipulations, the counterparties "
         "hereby irrevocably acknowledge that the indemnification obligations enumerated "
         "throughout the preceding paragraphs shall remain operative in perpetuity, insofar "
         "as such obligations are not subsequently extinguished by a superseding instrument "
         "executed with commensurate formality and corresponding authority. ") * 3


def _font(sz, bold=False):
    for p in ([
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf", "/System/Library/Fonts/Helvetica.ttc",
    ]):
        try:
            return ImageFont.truetype(p, sz)
        except Exception:
            continue
    return ImageFont.load_default()


def _center(d, x, y, text, font, fill):
    w = d.textlength(text, font=font)
    d.text((x - w / 2, y), text, font=font, fill=fill)


def _png(im) -> bytes:
    b = io.BytesIO(); im.save(b, format="PNG"); return b.getvalue()


# ── chart / diagram toolkit ───────────────────────────────────────────────────

def bar_chart(title, data):
    W, H = 780, 470; im = Image.new("RGB", (W, H), "white"); d = ImageDraw.Draw(im)
    _center(d, W / 2, 20, title, _font(26, True), (25, 25, 35))
    x0, y0, x1, y1 = 90, 95, W - 45, H - 70
    mx = max(v for _, v in data) * 1.18
    n = len(data)
    for i, (lab, val) in enumerate(data):
        cx = x0 + (i + 0.5) * (x1 - x0) / n
        bw = (x1 - x0) / (n * 1.7)
        bh = (val / mx) * (y1 - y0)
        d.rectangle([cx - bw / 2, y1 - bh, cx + bw / 2, y1], fill=PALETTE[i % len(PALETTE)])
        _center(d, cx, y1 - bh - 24, str(val), _font(18, True), (40, 40, 40))
        _center(d, cx, y1 + 8, lab, _font(17), (40, 40, 40))
    d.line([x0, y1, x1, y1], fill=(80, 80, 80), width=2)
    d.line([x0, y0, x0, y1], fill=(80, 80, 80), width=2)
    return _png(im)


def line_chart(title, labels, series):
    W, H = 780, 470; im = Image.new("RGB", (W, H), "white"); d = ImageDraw.Draw(im)
    _center(d, W / 2, 20, title, _font(26, True), (25, 25, 35))
    x0, y0, x1, y1 = 90, 95, W - 45, H - 70
    allv = [v for _, s in series for v in s]
    mx = max(allv) * 1.15
    for si, (name, s) in enumerate(series):
        pts = []
        for i, v in enumerate(s):
            px = x0 + i * (x1 - x0) / (len(s) - 1)
            py = y1 - (v / mx) * (y1 - y0)
            pts.append((px, py))
        d.line(pts, fill=PALETTE[si % len(PALETTE)], width=4, joint="curve")
        for p in pts:
            d.ellipse([p[0] - 4, p[1] - 4, p[0] + 4, p[1] + 4], fill=PALETTE[si % len(PALETTE)])
        d.text((x1 - 120, y0 + si * 22), "— " + name, font=_font(16), fill=PALETTE[si % len(PALETTE)])
    for i, lab in enumerate(labels):
        px = x0 + i * (x1 - x0) / (len(labels) - 1)
        _center(d, px, y1 + 8, lab, _font(16), (40, 40, 40))
    d.line([x0, y1, x1, y1], fill=(80, 80, 80), width=2)
    d.line([x0, y0, x0, y1], fill=(80, 80, 80), width=2)
    return _png(im)


def pie_chart(title, slices):
    W, H = 700, 470; im = Image.new("RGB", (W, H), "white"); d = ImageDraw.Draw(im)
    _center(d, W / 2, 20, title, _font(26, True), (25, 25, 35))
    cx, cy, r = 260, 260, 150
    total = sum(v for _, v in slices); ang = -90
    for i, (lab, val) in enumerate(slices):
        sweep = 360 * val / total
        d.pieslice([cx - r, cy - r, cx + r, cy + r], ang, ang + sweep, fill=PALETTE[i % len(PALETTE)])
        mid = math.radians(ang + sweep / 2)
        d.text((520, 130 + i * 34), f"■ {lab}  {round(100*val/total)}%", font=_font(18),
               fill=PALETTE[i % len(PALETTE)])
        ang += sweep
    return _png(im)


def flow_diagram(title, steps):
    W, H = 860, 300; im = Image.new("RGB", (W, H), "white"); d = ImageDraw.Draw(im)
    _center(d, W / 2, 18, title, _font(24, True), (25, 25, 35))
    n = len(steps); bw, bh, gap = 150, 90, (W - 60 - n * 150) / max(1, n - 1)
    y = 130
    for i, s in enumerate(steps):
        x = 30 + i * (bw + gap)
        d.rounded_rectangle([x, y, x + bw, y + bh], radius=12, fill=(238, 244, 252),
                            outline=BLUE, width=3)
        _center(d, x + bw / 2, y + bh / 2 - 12, s, _font(17, True), (25, 55, 110))
        if i < n - 1:
            ax = x + bw; ay = y + bh / 2
            d.line([ax, ay, ax + gap, ay], fill=GREY, width=4)
            d.polygon([(ax + gap, ay), (ax + gap - 12, ay - 8), (ax + gap - 12, ay + 8)], fill=GREY)
    return _png(im)


def image_of_text():
    im = Image.new("RGB", (900, 230), "white"); d = ImageDraw.Draw(im)
    f = _font(34, True); f2 = _font(26)
    d.text((30, 24), "Q4 Revenue Report — 2026", font=f, fill=(20, 20, 40))
    d.text((30, 90), "Total revenue up 14% year over year", font=f2, fill=(40, 40, 40))
    d.text((30, 140), "across every regional business unit", font=f2, fill=(40, 40, 40))
    return _png(im)


# ── documents ─────────────────────────────────────────────────────────────────

def build_pptx(path: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    # 1 — Title slide, EMPTY title (2.4.6); subtitle French (3.1.2) + sensory (1.3.3).
    s = prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text = ""
    if len(s.placeholders) > 1:
        s.placeholders[1].text = FRENCH + "  " + SENSORY
    # 2 — two charts, no alt (1.1.1 x2 → GPU alt text).
    s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "Regional Performance"
    s.shapes.add_picture(io.BytesIO(bar_chart("Q4 Revenue by Region ($M)",
        [("North", 120), ("South", 70), ("East", 150), ("West", 50)])), Inches(0.3), Inches(1.4), Inches(4.6))
    s.shapes.add_picture(io.BytesIO(line_chart("Monthly Active Users (k)", ["Jan", "Feb", "Mar", "Apr", "May"],
        [("2025", [20, 28, 33, 41, 55]), ("2024", [15, 18, 22, 25, 30])])), Inches(5.1), Inches(1.4), Inches(4.6))
    # 3 — flow diagram (1.1.1) + two ambiguous "click here" links (2.4.9).
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(io.BytesIO(flow_diagram("Remediation Workflow",
        ["Scan", "Assess", "Remediate", "Publish"])), Inches(0.5), Inches(0.6), Inches(9))
    tf = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1.5)).text_frame
    tf.text = "For the methodology "
    for url in ("https://example.com/methodology", "https://example.com/scores"):
        r = tf.paragraphs[0].add_run(); r.text = "click here"; r.hyperlink.address = url
        tf.paragraphs[0].add_run().text = "  and  "
    # 4 — pie chart (1.1.1) + image-of-text (1.4.5) + dense text (3.1.5).
    s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "Budget & Notes"
    s.shapes.add_picture(io.BytesIO(pie_chart("Spend by Category",
        [("Engineering", 45), ("Sales", 25), ("Ops", 18), ("Other", 12)])), Inches(0.3), Inches(1.4), Inches(4.4))
    s.shapes.add_picture(io.BytesIO(image_of_text()), Inches(5.0), Inches(1.5), Inches(4.6))
    tf = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2)).text_frame
    tf.word_wrap = True; tf.text = DENSE
    prs.save(str(path))


def build_docx(path: Path):
    import re, zipfile
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor
    doc = Document()
    doc.core_properties.title = ""; doc.core_properties.language = ""
    doc.add_heading("Annual Accessibility Report — Executive Summary", level=1)
    doc.add_heading("Scoring Methodology", level=3)                      # 2.4.6 skip
    doc.add_paragraph("The following figures summarize performance across every business unit.")
    p = doc.add_paragraph(); r = p.add_run("Confidential — internal only."); r.font.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)  # 1.4.3
    doc.add_picture(io.BytesIO(bar_chart("Compliance Score by Department",
        [("Legal", 82), ("HR", 61), ("Finance", 74), ("Research", 55)])), width=Inches(5.2))  # 1.1.1 (GPU)
    doc.add_picture(io.BytesIO(flow_diagram("Certification Process",
        ["Discover", "Score", "Fix", "Certify"])), width=Inches(6.0))                          # 1.1.1 (GPU)
    doc.add_paragraph(FRENCH)                                            # 3.1.2
    doc.add_paragraph(SENSORY)                                           # 1.3.3
    for _ in range(3):
        pp = doc.add_paragraph("Each criterion was weighted by its impact on end users, then "
            "normalized so no single dimension dominates the composite score reported to the "
            "steering committee for quarterly review."); pp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # 1.4.8
    doc.add_paragraph(DENSE)                                             # 3.1.5
    t = doc.add_table(rows=4, cols=3)                                    # 1.3.1 no header
    for row, vals in zip(t.rows, [["North", "1,240", "up"], ["South", "980", "flat"],
                                   ["East", "1,510", "up"], ["West", "760", "down"]]):
        for c, v in zip(row.cells, vals):
            c.text = v
    doc.save(str(path))
    # strip default language so 3.1.1 fires
    with zipfile.ZipFile(path) as z:
        parts = {n: z.read(n) for n in z.namelist()}
    for n in ("word/styles.xml", "word/settings.xml"):
        if n in parts:
            parts[n] = re.sub(rb"<w:lang\b[^>]*/>", b"", parts[n])
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n, data in parts.items():
            z.writestr(n, data)


def build_pdf(path: Path):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    c = canvas.Canvas(str(path), pagesize=letter)                        # no title(2.4.2)/lang(3.1.1)
    heads = ["Executive Summary", "Regional Analysis", "Findings", "Budget",
             "Recommendations", "Appendix"]
    for i, h in enumerate(heads):                                        # 2.4.1 headings, no outline
        c.setFillColorRGB(0, 0, 0); c.setFont("Helvetica-Bold", 18); c.drawString(72, 720, h)
        c.setFont("Helvetica", 11)
        for j in range(3):
            c.drawString(72, 690 - j * 15, "Ordinary body paragraph text below the heading.")
        if i == 0:
            c.setFillColorRGB(0.8, 0.8, 0.8); c.drawString(72, 630, "Confidential — internal only.")  # 1.4.3
            c.setFillColorRGB(0, 0, 0); c.drawString(72, 610, FRENCH[:95]); c.drawString(72, 594, SENSORY[:95])
        if i == 1:
            c.drawImage(ImageReader(io.BytesIO(bar_chart("Revenue by Region ($M)",
                [("North", 120), ("South", 70), ("East", 150), ("West", 50)]))), 72, 300, width=340, height=205)  # 1.1.1
        if i == 3:
            c.drawImage(ImageReader(io.BytesIO(pie_chart("Spend by Category",
                [("Eng", 45), ("Sales", 25), ("Ops", 18), ("Other", 12)]))), 72, 300, width=300, height=200)      # 1.1.1
            c.drawImage(ImageReader(io.BytesIO(image_of_text())), 72, 150, width=360, height=92)                   # 1.4.5
        c.showPage()
    c.save()


def main():
    out = Path(sys.argv[1]).expanduser() if len(sys.argv) > 1 else Path.home() / "Downloads" / "acp-showcase"
    out.mkdir(parents=True, exist_ok=True)
    build_pptx(out / "showcase-deck.pptx")
    build_docx(out / "showcase-report.docx")
    build_pdf(out / "showcase-brief.pdf")
    for f in sorted(out.iterdir()):
        print(f"wrote {f}")


if __name__ == "__main__":
    main()

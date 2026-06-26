"""Generate a large, varied accessibility test corpus (~200 files) to upload to Drive.

Realistic names (Dept-DocType-NN.ext), a controlled spread of WCAG profiles
(clean / minor / serious / critical / unanalysable) across HTML/PDF/DOCX/PPTX/XLSX,
and synthetic PII sprinkled into a fraction of files to exercise the PII detector.

All PII is fake: example.com emails, reserved 555-01xx phones, standard card-network
TEST numbers, fake-but-valid-format SSNs. No real data.

  python3 scripts/gen_bulk_corpus.py --count 200 --out /tmp/acp-bulk-corpus
"""
from __future__ import annotations
import argparse, random, textwrap
from pathlib import Path

DEPTS = ["HR", "Finance", "Legal", "Marketing", "Operations", "IT", "Research",
         "Facilities", "Procurement", "Clinical", "Compliance", "Sales", "Support",
         "Product", "Security"]
DOCTYPES = ["Policy", "Handbook", "Report", "Invoice", "Memo", "Onboarding-Form",
            "Proposal", "Audit", "Newsletter", "Contract", "Spec", "FAQ",
            "Guidelines", "Minutes", "Brief", "Roadmap", "Datasheet", "Notice",
            "Summary", "Plan", "Checklist", "Review"]
BODIES = [
    "This document outlines the responsibilities, procedures, and standards that "
    "apply across the organisation. All staff are expected to review it annually.",
    "The following summary captures the key outcomes from the quarter, including "
    "progress against objectives, open risks, and recommended next steps.",
    "Please read the sections below carefully. They describe the process, the "
    "required approvals, and the escalation path for exceptions.",
    "Our commitment is to clear, accessible communication. This material is part "
    "of an ongoing effort to keep records accurate, current, and easy to follow.",
    "This report details findings and observations gathered during the review "
    "period, along with their severity and the owners assigned to each item.",
]

# ── synthetic PII pools (all fake) ──
SSNS = ["123-45-6789", "078-05-1120", "219-09-9998", "457-55-5462", "212-09-7788"]
CARDS = ["4111 1111 1111 1111", "5555 5555 5555 4444", "3782 822463 10005",
         "6011 0009 9013 9424", "5105 1051 0510 5100"]
EMAILS = ["jane.doe@example.com", "john.smith@example.org", "a.patel@example.net",
          "m.garcia@example.com", "k.brown@example.co", "r.lee@example.org"]
PHONES = ["(415) 555-0142", "212-555-0173", "650-555-0118", "(312) 555-0190",
          "206-555-0164", "+1 408-555-0127"]


def pii_lines(level: str, rnd: random.Random) -> list[str]:
    if level == "none":
        return []
    if level == "contact":
        return [f"Contact: {rnd.choice(EMAILS)} · {rnd.choice(PHONES)}"]
    # sensitive
    return [
        f"Employee SSN: {rnd.choice(SSNS)}",
        f"Card on file: {rnd.choice(CARDS)}",
        f"Email: {rnd.choice(EMAILS)} · Phone: {rnd.choice(PHONES)}",
    ]


# ── format builders ───────────────────────────────────────────────────────────
def build_html(path, prof, title, body, pii, rnd):
    lang = ' lang="en"' if prof["lang"] else ""
    head = f"<title>{title}</title>" if prof["title"] else ""
    img = f'<img src="banner.png" alt="{title} banner">' if prof["alt"] else '<img src="banner.png">'
    heading = f"<h1>{title}</h1>" if prof["headings"] else f'<p style="font-size:26px;font-weight:700">{title}</p>'
    link = '<a href="/policy">Read the full policy document</a>' if prof["links"] else '<a href="/policy">click here</a>'
    color = "#1a1a1a" if prof["contrast"] else "#cfcfcf"
    field = '<label for="e">Email</label><input id="e" name="e">' if prof["labels"] else '<input name="e" placeholder="Email">'
    pii_html = "".join(f"<p>{l}</p>" for l in pii_lines(pii, rnd))
    html = (f'<!DOCTYPE html><html{lang}><head><meta charset="utf-8">{head}</head><body>'
            f'{heading}{img}<p style="color:{color}">{body}</p><p>{link}</p>'
            f'<form>{field}</form>{pii_html}</body></html>')
    path.write_text(html)


def build_pdf(path, prof, title, body, pii, rnd):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    c = canvas.Canvas(str(path), pagesize=letter)
    if prof["title"]:
        c.setTitle(title)
    y = 730
    c.drawString(72, y, title); y -= 28
    if not prof["contrast"]:
        c.setFillColorRGB(0.9, 0.9, 0.9)
    for line in textwrap.wrap(body, 86):
        c.drawString(72, y, line); y -= 15
    c.setFillColorRGB(0, 0, 0); y -= 8
    c.drawString(72, y, "Read the full policy document" if prof["links"] else "click here"); y -= 22
    for line in pii_lines(pii, rnd):
        c.drawString(72, y, line); y -= 15
    c.save()


def build_docx(path, prof, title, body, pii, rnd):
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    if prof["title"]:
        doc.core_properties.title = title
    if prof["headings"]:
        doc.add_heading(title, 0)
    else:
        p = doc.add_paragraph(); r = p.add_run(title); r.bold = True; r.font.size = Pt(16)
    p = doc.add_paragraph(body)
    if not prof["contrast"] and p.runs:
        p.runs[0].font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    for line in pii_lines(pii, rnd):
        doc.add_paragraph(line)
    doc.save(str(path))


def build_pptx(path, prof, title, body, pii, rnd):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title if prof["title"] else ""
    body_ph = slide.placeholders[1]
    tf = body_ph.text_frame
    tf.text = body
    for line in pii_lines(pii, rnd):
        tf.add_paragraph().text = line
    prs.save(str(path))


def build_xlsx(path, prof, title, body, pii, rnd):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "Records"
    if prof["title"]:
        wb.properties.title = title
    ws["A1"] = title
    ws["A2"] = body
    for i, line in enumerate(pii_lines(pii, rnd)):
        ws[f"A{4 + i}"] = line
    wb.save(str(path))


def build_unanalysable(path, rnd):
    # Corrupt/empty file the engine cannot parse → "unanalysable".
    kind = rnd.choice(["truncated", "garbage", "empty"])
    if kind == "empty":
        path.write_bytes(b"")
    elif kind == "garbage":
        path.write_bytes(bytes(rnd.getrandbits(8) for _ in range(512)))
    else:  # truncated header
        path.write_bytes(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n" + bytes(rnd.getrandbits(8) for _ in range(64)))


BUILDERS = {"html": build_html, "pdf": build_pdf, "docx": build_docx,
            "pptx": build_pptx, "xlsx": build_xlsx}

# weighted profile generator
def make_profile(kind: str) -> dict:
    full = dict(title=True, lang=True, alt=True, headings=True, links=True, contrast=True, labels=True)
    if kind == "clean":
        return full
    if kind == "minor":
        p = dict(full); p[random.choice(["alt", "labels", "links"])] = False; return p
    if kind == "serious":
        p = dict(full)
        for k in random.sample(["alt", "headings", "links", "contrast"], 2):
            p[k] = False
        return p
    # critical
    return dict(title=False, lang=False, alt=False, headings=False, links=False, contrast=False, labels=False)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--count", type=int, default=200)
    ap.add_argument("--out", default="/tmp/acp-bulk-corpus")
    ap.add_argument("--seed", type=int, default=7)
    args = ap.parse_args()
    rnd = random.Random(args.seed)
    random.seed(args.seed)
    out = Path(args.out); out.mkdir(parents=True, exist_ok=True)
    for old in out.glob("*"):  # clean slate
        old.unlink()

    FORMAT_W = [("html", 32), ("pdf", 26), ("docx", 24), ("pptx", 11), ("xlsx", 7)]
    PROFILE_W = [("clean", 34), ("minor", 24), ("serious", 22), ("critical", 14), ("unanalysable", 6)]
    PII_W = [("none", 80), ("contact", 14), ("sensitive", 6)]
    fmts = [f for f, w in FORMAT_W for _ in range(w)]
    profs = [p for p, w in PROFILE_W for _ in range(w)]
    piis = [p for p, w in PII_W for _ in range(w)]

    counts = {"format": {}, "profile": {}, "pii": {}}
    used = set()
    n = 0
    while n < args.count:
        fmt = rnd.choice(fmts)
        kind = rnd.choice(profs)
        pii = rnd.choice(piis)
        dept, dtype = rnd.choice(DEPTS), rnd.choice(DOCTYPES)
        idx = rnd.randint(1, 99)
        name = f"{dept}-{dtype}-{idx:02d}"
        ext = "pdf" if kind == "unanalysable" and fmt not in ("html",) else fmt
        if kind == "unanalysable":
            ext = rnd.choice(["pdf", "docx", "xlsx"])
        fname = f"{name}.{ext}"
        if fname in used:
            continue
        used.add(fname)
        path = out / fname
        title = f"{dept} {dtype.replace('-', ' ')}"
        body = rnd.choice(BODIES)
        try:
            if kind == "unanalysable":
                build_unanalysable(path, rnd)
            else:
                BUILDERS[fmt](path, make_profile(kind), title, body, pii, rnd)
        except Exception as e:
            print(f"  ! skip {fname}: {e}"); used.discard(fname); continue
        counts["format"][ext] = counts["format"].get(ext, 0) + 1
        counts["profile"][kind] = counts["profile"].get(kind, 0) + 1
        counts["pii"][pii] = counts["pii"].get(pii, 0) + 1
        n += 1

    print(f"Generated {n} files in {out}\n")
    for dim in ("format", "profile", "pii"):
        print(f"  by {dim}:", ", ".join(f"{k}={v}" for k, v in sorted(counts[dim].items())))


if __name__ == "__main__":
    main()

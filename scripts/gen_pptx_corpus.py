#!/usr/bin/env python3
"""A LABELLED .pptx corpus — the third format to get ground truth, after .docx and .xlsx.

WHY THIS EXISTS. `scripts/gen_fixture_coverage.py` reports coverage per (criterion, format) pair;
pptx sat at 0 of 17 because no labelled corpus existed. This declares SIXTEEN of them.

SAME RULE AS THE XLSX CORPUS: a pair is only declared when a FIRST-PARTY detector — pure Python
in `api/office_structure.py`, no .NET engine — was driven against the fixture and confirmed to
fire, with an adversarial counterpart confirmed to stay silent. Coverage is counted from
declarations, so a fixture whose seeded violation nobody can confirm is caught would raise the
number without raising what the number measures.

    1.1.1   a picture with no descr                office_non_text_content_checks
    1.4.1   a hyperlink run with u="none"          office_color_only_checks
    1.4.3   run colour on an explicit shape fill   pptx_contrast_checks
    1.4.11  a faint shape outline on its fill      pptx_nontext_contrast_checks
    2.1.2   an embedded control                    office_control_review_checks
    2.4.3   title placeholder not first in order   pptx_focus_order_checks
    2.4.4   a vague hyperlink label                pptx_checks
    2.4.6   an empty title placeholder             pptx_checks
    1.3.3   a sensory-only instruction in a box    textchecks.detect_sensory
    4.1.2   an embedded control (same fixture)     office_control_review_checks

    1.3.1   a table with no header row designated  .NET TableHeaderRule      (DECLARED_ENGINE)
    1.3.2   three boxes out of visual order         .NET ReadingOrderRule     (DECLARED_ENGINE)
    2.4.2   a slide with no title placeholder      .NET SlideTitleRule       (DECLARED_ENGINE)
    3.1.1   no language anywhere in the package     .NET DocumentLanguageRule (DECLARED_ENGINE)

The last two are the exception to the rule above and are kept in a SEPARATE tuple for it: no
first-party Python detector exists for either on ANY Office format, so their labels are proven
where the .NET analyser is built and skipped where it is not. They earn the asymmetry by being
certifying pairs — 2.4.2 and 3.1.1 are among the seventeen (criterion, format) pairs in the whole
preset that can return a PASS, so a false clean result on them is a certification rather than an
advisory, and before these fixtures nothing in the suite checked either.

The one not here is 2.1.1, which is human-only on pptx by registration and cannot certify a pass
either way. 1.3.3 was on that list until it
was checked: it is a TEXT predicate with no engine, and is declared below. A criterion's
neighbours are not evidence of its reachability.

TWO DETECTOR SUBTLETIES WORTH KNOWING, because a fixture that ignores them silently covers
nothing:

  * 1.4.3 needs an EXPLICIT shape solid fill as well as an explicit run colour. A bare textbox
    has no fill, so the detector cannot know what the text sits on and correctly says nothing —
    the first draft of that fixture declared 1.4.3 and detected zero.
  * 1.4.3 is judged at the WCAG LARGE-text bar (3:1), not 4.5:1, because run font size is often
    inherited from the placeholder and not reliably knowable. Flagging only below the large-text
    threshold guarantees every finding is a real failure at any size.

Run:
    python scripts/gen_pptx_corpus.py --out ~/Downloads/acp-pptx-eval/sc-corpus
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "api"))
sys.path.insert(0, str(ROOT / "scripts"))

import corpus_expectations as ce  # noqa: E402

FMT = "pptx"

INK = (0x1A, 0x1F, 0x26)      # ~15:1 on white
FAINT = (0xDD, 0xDD, 0xDD)    # ~1.6:1 on white — under even the 3:1 large-text bar
PAPER = (0xFF, 0xFF, 0xFF)

# Stamped on the base deck so that only the fixture that deliberately withholds it can raise
# 3.1.1 under the .NET analyser. The xlsx corpus learned this the expensive way: openpyxl leaves
# the language unset, so before its base workbook set one EVERY fixture there also carried a
# 3.1.1 finding, and the single-criterion labels were true only because a bare container has no
# .NET. python-pptx is quieter about it — its default template's masters carry lang attributes,
# which happens to keep the rule silent — but "happens to" is not a property to rely on, and an
# explicit metadata language is the same guarantee without the accident.
DOC_LANG = "en-GB"

# A 1x1 PNG, inline — the fixture needs an image to embed, not a picture of anything, and a
# committed binary is provenance a reviewer has to take on trust.
_PNG_1PX = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000d49444154789c6360f8cf0000010101001b0b8b5c0000000049454e44ae426082")


def _deck(title: str = "Q3 regional revenue"):
    """A one-slide deck that is clean on everything a fixture does not deliberately break: a
    real title (so 2.4.6 stays quiet) and title-first placeholder order (so 2.4.3 does)."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # title only
    slide.shapes.title.text = title
    prs.core_properties.language = DOC_LANG
    return prs, slide


def _textbox(slide, text: str, *, colour=INK, fill=PAPER, underline=None, link: str | None = None):
    """A textbox with an EXPLICIT solid fill — 1.4.3 needs one to know what the text sits on."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(1))
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(*fill)
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(*colour)
    if underline is not None:
        run.font.underline = underline
    if link:
        run.hyperlink.address = link
    return tb


def _box(slide, text: str, top_in: float, left_in: float = 1.0):
    """A textbox at an EXPLICIT position — which is what makes it visible to ReadingOrderRule.

    Distinct from `_textbox` above, which exists for the contrast fixtures and carries an explicit
    fill and run colour. This one carries neither: a 1.3.2 fixture that also set colours would be
    seeding two criteria and could not be labelled single-criterion.
    """
    from pptx.util import Inches
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(4), Inches(0.8))
    tb.text_frame.paragraphs[0].add_run().text = text
    return tb


def _position_title(slide):
    """Give the title placeholder an EXPLICIT position, so it takes a visual rank as well as a
    tab index.

    THE CONTROL HAD NO MARGIN WITHOUT THIS, and a bite check is what found it. ReadingOrderRule
    assigns tab order over every `p:sp` and then discards the ones with no `a:off`, so an
    unpositioned placeholder consumes an index without taking a rank. A python-pptx placeholder
    inherits its geometry from the layout and writes no `a:xfrm` — so with the title unpositioned,
    a PERFECTLY ORDERED deck already sat at |visualRank - tabOrder| == 1, which is the rule's
    tolerance exactly. One more unpositioned shape ahead of the boxes and the clean control would
    have fired, and the false positive would have been blamed on the detector.

    Positioning the title puts it at rank 0 and tab 0, so the ordered fixture sits at a gap of
    zero — two clear of the threshold — while the violation still exceeds it.
    tests/test_pptx_corpus.py::test_the_reading_order_control_has_margin_not_luck pins that gap
    so it cannot silently erode again.
    """
    from pptx.util import Inches
    title = slide.shapes.title
    title.left, title.top = Inches(0.5), Inches(0.4)
    title.width, title.height = Inches(9), Inches(1.2)
    return title


def _table(slide, *, first_row: bool):
    """A three-row table, with the header row designated or not.

    Three rows rather than two because the rule skips a table of one row entirely, and a
    two-row table would leave no margin if that guard ever moved to two.
    """
    from pptx.util import Inches
    frame = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2))
    # ALT TEXT ON A TABLE, WHICH LOOKS UNNECESSARY AND IS NOT. Pptx/Rules/AltTextRule.cs walks
    # `Descendants<GraphicFrame>()` as well as `Descendants<Picture>()`, and a table IS a graphic
    # frame — so a table with no `descr` raises 1.1.1. The first-party
    # office_non_text_content_checks only reads pictures, so a bare checkout sees nothing and the
    # fixture looked single-criterion right up until CI ran the analyser.
    frame._element.nvGraphicFramePr.cNvPr.set("descr", "Revenue by region, two columns")
    table = frame.table
    table.first_row = first_row
    for r, row in enumerate((("Region", "Revenue"), ("North", "412"), ("South", "388"))):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    return frame


def _replace_parts(path: Path, parts: dict[str, str]) -> None:
    """Rewrite the saved deck, REPLACING any part named in `parts` and adding the rest.

    Replacing rather than appending matters: a zip may legally hold two entries with the same
    name, and readers differ on which one wins — a fixture built that way would be testing the
    zip library, not the detector.
    """
    import shutil
    import tempfile
    tmp = Path(tempfile.mkdtemp()) / path.name
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename in parts:
                continue
            zout.writestr(item, zin.read(item.filename))
        for name, body in parts.items():
            zout.writestr(name, body)
    shutil.move(str(tmp), str(path))


# ── 2.4.6 Headings and Labels ────────────────────────────────────────────────────

def f_title_empty(prs, slide):
    """AN EMPTY TITLE PLACEHOLDER IS TWO VIOLATIONS, and this fixture has to say so.

    It is the 2.4.6 case (a heading with no label) and, under the .NET analyser, also the 2.4.2
    case — SlideTitleRule flags a title placeholder "present but contains no text" in the same
    branch as a missing one. Before 2.4.2 was declared, this fixture was labelled single-criterion
    and was simply wrong in CI, in the way that is hardest to notice: nothing on a bare checkout
    can raise 2.4.2, so nothing here could contradict it.

    The dedicated 2.4.2 fixture is `no-slide-title`, which removes the placeholder entirely and so
    isolates the criterion. This one cannot: the two coincide on the same shape.
    """
    slide.shapes.title.text = ""
    _textbox(slide, "Revenue grew across every region this quarter.")
    return ({"2.4.6": "REVIEW", "2.4.2": "FAIL"},
            "a title placeholder present but left empty — 2.4.6 by the first-party detector, and "
            "2.4.2 under the analyser, which treats an empty placeholder as an absent title")


def f_title_ok(prs, slide):
    _textbox(slide, "Revenue grew across every region this quarter.")
    return {"2.4.6": "REVIEW"}, "a real title — must not be flagged (adversarial)"


# ── 2.4.4 Link Purpose ───────────────────────────────────────────────────────────

def f_link_vague(prs, slide):
    _textbox(slide, "click here", link="https://example.com/fy26-travel-policy.pdf")
    return {"2.4.4": "REVIEW"}, "a 'click here' label — the destination is unknowable"


def f_link_descriptive_ok(prs, slide):
    _textbox(slide, "Read the FY26 travel policy",
             link="https://example.com/fy26-travel-policy.pdf")
    return {"2.4.4": "REVIEW"}, "a descriptive label — must not be flagged (adversarial)"


# ── 1.4.1 Use of Color ───────────────────────────────────────────────────────────

def f_link_underline_off(prs, slide):
    # u="none" on the run's rPr: the link is set apart from surrounding text by COLOUR ALONE.
    _textbox(slide, "Read the FY26 travel policy", underline=False,
             link="https://example.com/fy26-travel-policy.pdf")
    return {"1.4.1": "REVIEW"}, "a hyperlink with its underline suppressed — colour is the only cue"


def f_link_underlined_ok(prs, slide):
    _textbox(slide, "Read the FY26 travel policy", underline=True,
             link="https://example.com/fy26-travel-policy.pdf")
    return {"1.4.1": "REVIEW"}, "the same link, underlined — must not be flagged (adversarial)"


# ── 1.4.3 Contrast (Minimum) ─────────────────────────────────────────────────────

def f_contrast_fail(prs, slide):
    _textbox(slide, "Revenue grew across every region this quarter.", colour=FAINT)
    return {"1.4.3": "REVIEW"}, "#DDDDDD on an explicit white fill ≈ 1.6:1, under even the 3:1 bar"


def f_contrast_ok(prs, slide):
    _textbox(slide, "Revenue grew across every region this quarter.")
    return {"1.4.3": "REVIEW"}, "the same shape, legible ink — must not be flagged (adversarial)"


# ── 1.4.11 Non-text Contrast ─────────────────────────────────────────────────────

def _shape(slide, outline):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(1), Inches(2), Inches(2.5), Inches(1))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(*PAPER)
    sh.line.color.rgb = RGBColor(*outline)
    return sh


def f_shape_faint_outline(prs, slide):
    _shape(slide, (0xC8, 0xC8, 0xC8))
    return {"1.4.11": "REVIEW"}, "a shape outline #C8C8C8 on white ≈ 1.8:1, under 3:1"


def f_shape_strong_outline_ok(prs, slide):
    # THE SAME SHAPE, drawn visibly. If 1.4.11 ever reported on a shape's PRESENCE rather than
    # its measured ratio, this would fire too — and that measurement is the whole criterion.
    _shape(slide, INK)
    return {"1.4.11": "REVIEW"}, "the same shape at ≈15:1 — must not be flagged (adversarial)"


# ── 2.4.3 Focus Order ────────────────────────────────────────────────────────────

def f_title_not_first(prs, slide):
    """Rebuild on a layout with two placeholders and move the title after the body."""
    return None   # handled by `post`, below — see FIXTURES


def f_focus_order(prs, slide):
    from pptx import Presentation
    # A fresh deck on the title+content layout: two placeholders, so document order is meaningful.
    prs2 = Presentation()
    s = prs2.slides.add_slide(prs2.slide_layouts[1])
    s.placeholders[1].text = "Revenue grew across every region this quarter."
    s.shapes.title.text = "Q3 regional revenue"
    tree = s.shapes._spTree
    el = s.shapes.title._element
    tree.remove(el)
    tree.append(el)                       # title now LAST in document order
    return prs2, ({"2.4.3": "REVIEW"},
                  "the title placeholder is not first in document order")


def f_focus_order_ok(prs, slide):
    from pptx import Presentation
    prs2 = Presentation()
    s = prs2.slides.add_slide(prs2.slide_layouts[1])
    s.shapes.title.text = "Q3 regional revenue"
    s.placeholders[1].text = "Revenue grew across every region this quarter."
    return prs2, ({"2.4.3": "REVIEW"},
                  "title first, body second — must not be flagged (adversarial)")


# ── 4.1.2 Name, Role, Value / 2.1.2 No Keyboard Trap ─────────────────────────────

def f_embedded_control(prs, slide):
    _textbox(slide, "Complete the survey below.")
    # ONE fixture, TWO criteria: an embedded control is evidence for the accessible-name question
    # AND the keyboard-trap question, and the detector reports both. Neither is settleable from
    # the file alone, which is why both are review-lane.
    return ({"4.1.2": "REVIEW", "2.1.2": "REVIEW"},
            "an embedded VBA project — name/role and keyboard behaviour both need a human",
            {"ppt/vbaProject.bin": "stand-in for a real macro project"})


def f_no_controls_ok(prs, slide):
    _textbox(slide, "Complete the survey below.")
    return ({"4.1.2": "REVIEW", "2.1.2": "REVIEW"},
            "a static deck with no controls — must not be flagged (adversarial)")


# ── 1.1.1 Non-text Content ───────────────────────────────────────────────────────

def f_picture_no_alt(prs, slide):
    import tempfile
    from pptx.util import Inches
    png = Path(tempfile.mkdtemp()) / "px.png"
    png.write_bytes(_PNG_1PX)
    slide.shapes.add_picture(str(png), Inches(1), Inches(2), Inches(2), Inches(2))
    return {"1.1.1": "REVIEW"}, "a picture with no descr — nothing for a screen reader"


def f_no_picture_ok(prs, slide):
    _textbox(slide, "North 412  ·  South 388  ·  East 501")
    return {"1.1.1": "REVIEW"}, "no non-text content at all — must not be flagged (adversarial)"


# ── 1.3.3 Sensory Characteristics — decided by the prose, not the deck ──────────
# The one criterion here that reads no slide structure at all: textchecks.detect_sensory reads the
# EXTRACTED TEXT for an instruction identifying a control only by shape, colour or position. The
# words are the fixture and the deck is the container — which is why the wording is identical to
# the .xlsx and .pdf corpora. A detector change then shows up as the same result in three places
# instead of three arguments about three different sentences.
SENSORY_BAD = ("To continue, click the round green button on the right. "
               "See the box below for the payment terms.")
SENSORY_OK = ("To continue, choose Submit under Payment options. "
              "The payment terms are in the Payment terms section.")


def f_sensory_instruction(prs, slide):
    _textbox(slide, SENSORY_BAD)
    return ({"1.3.3": "FAIL"},
            "an instruction identifying a control only by shape, colour and position")


def f_sensory_instruction_ok(prs, slide):
    _textbox(slide, SENSORY_OK)
    return ({"1.3.3": "REVIEW"},
            "the same instruction naming the control and the section (adversarial)")



# ── 1.4.5 Images of Text — decided by OCR, not by the container ─────────────────
# Like 1.3.3, this criterion reads the DOCUMENT'S PIXELS rather than its structure: ocr.py runs
# tesseract over each embedded raster and flags any carrying >= ocr._MIN_WORDS (10) real words.
# The image is the fixture and the deck is the container, so the wording is identical across
# the .xlsx, .pptx and .pdf corpora — see tests/test_images_of_text_corpus.py, which asserts that
# sameness so one detector change reads as one result in three places.
#
# THE ALT IS CORRECT ON PURPOSE, and that is what keeps the fixture single-criterion. 1.4.5 is
# alt-agnostic — `images_of_text` never looks at a descr — so an image of prose fails 1.4.5
# whether or not it is described. Leaving the alt off would fail 1.1.1 as well and the fixture
# would measure two things at once. (The .docx corpus's equivalent DOES declare both, which is
# correct there: gen_sc_corpus is scored by score_assessment rather than by the single-criterion
# sweep these three corpora use.)
#
# THE WORD FLOOR IS THE TRAP. _MIN_WORDS counts OCR'd TOKENS, not what was drawn — dates and
# phone numbers are not words. The .docx corpus hit this: three short lines recovered six words,
# so 1.4.9 (floor 3) fired and 1.4.5 (floor 10) did not, and the fixture read as an engine bug.
# Hence prose, and enough of it: this image OCRs to 29 words, with room for the odd merge
# ("begins on" comes back as "beginson").
IMAGE_OF_TEXT_PROSE = ("Open enrollment for the coming plan year",
                       "begins on the first of March and closes",
                       "on the last day of the same month. Call",
                       "the benefits office to change your plan.")
IMAGE_OF_TEXT_ALT = ("Open enrollment runs from 1 March to 31 March; call the benefits office "
                     "to change plan.")
# WCAG 1.4.5 exempts logotypes, and ocr._MIN_WORDS is what encodes that exemption here: two words
# is under the floor, so the control is clean for the same reason a real logo would be.
LOGO_TEXT = ("UT Health",)


def _text_png(lines, size=(900, 360)) -> bytes:
    """A PNG with real, OCR-legible text baked into it.

    Scaled up AFTER drawing because PIL's default bitmap font is too small for tesseract to read
    reliably at native size. An unreadable image-of-text fixture silently becomes a no-text
    fixture, and 1.4.5 then stops firing for a reason that has nothing to do with the detector —
    which is the failure mode a ground-truth corpus exists to make impossible."""
    import io
    from PIL import Image, ImageDraw
    im = Image.new("RGB", size, "white")
    d = ImageDraw.Draw(im)
    y = 20
    for line in lines:
        d.text((20, y), line, fill="black")
        y += 70
    im = im.resize((size[0] * 2, size[1] * 2), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, format="PNG")
    return buf.getvalue()


def _text_png_file(lines, size=(900, 360)):
    import tempfile
    png = Path(tempfile.mkdtemp()) / "image-of-text.png"
    png.write_bytes(_text_png(lines, size))
    return png


def f_image_of_text(prs, slide):
    from pptx.util import Inches
    _textbox(slide, "Enrollment notice")
    pic = slide.shapes.add_picture(str(_text_png_file(IMAGE_OF_TEXT_PROSE)),
                                   Inches(1), Inches(2), Inches(6), Inches(2.4))
    pic._element._nvXxPr.cNvPr.set("descr", IMAGE_OF_TEXT_ALT)
    return ({"1.4.5": "FAIL"},
            "a screenshot of prose pasted onto a slide, correctly described. The alt makes it "
            "reachable but not RESIZABLE — enlarging it pixelates rather than reflows, which is "
            "what 1.4.5 is about")


def f_image_of_text_logo_ok(prs, slide):
    from pptx.util import Inches
    _textbox(slide, "Enrollment notice")
    pic = slide.shapes.add_picture(str(_text_png_file(LOGO_TEXT, size=(300, 100))),
                                   Inches(1), Inches(2), Inches(2), Inches(0.7))
    pic._element._nvXxPr.cNvPr.set("descr", "UT Health")
    return ({"1.4.5": "REVIEW"},
            "a logotype with its text as the alt (adversarial). WCAG 1.4.5 exempts logos, so a "
            "finding here is a false positive")



# ── 2.4.2 Page Titled and 3.1.1 Language of Page — the ENGINE-VERIFIED pairs ────
# No first-party Python detector exists for either on any Office format, so these are confirmed
# by the .NET analyser (PptxRuleIds.SlideTitle, PptxRuleIds.DocumentLanguage) and only where it
# is built. That is what DECLARED_ENGINE keeps countable without diluting what DECLARED means —
# the same split the xlsx corpus introduced, and both are worth the asymmetry because 2.4.2 and
# 3.1.1 are among the seventeen (criterion, format) pairs in the preset that can return a PASS.
# A false clean result here is a certification, not an advisory.
#
# BUILT FROM THE VENDORED RULE SOURCE, AND THE PPTX RULE IS NOT THE XLSX ONE. Both formats have a
# class called DocumentLanguageRule and they read different things:
#
#   Xlsx/Rules/DocumentLanguageRule.cs    document.PackageProperties.Language, and nothing else
#   Pptx/Rules/DocumentLanguageRule.cs    that OR any lang/altLang on an a:rPr or a:endParaRPr in
#                                         any slide, slide master, or master's layout
#
# The xlsx recipe — clear `core_properties.language` — therefore does NOT trip the pptx rule.
# Measured, not reasoned: a deck built that way still answers "has content language" because
# python-pptx's default template ships `<a:rPr lang="en-US" smtClean="0"/>` in
# ppt/slideMasters/slideMaster1.xml. The fixture below strips those attributes from every slide,
# master and layout, which is the only way the rule's second branch goes quiet.
#
# The comment in that rule says why the branch exists: reading only PackageProperties.Language
# "false-positived essentially every real deck". So this is not an obscure edge — it is the
# normal case, and a fixture transplanted from the xlsx corpus would have declared the pair and
# detected nothing.

_LANG_BEARING_PARTS = ("ppt/slides/slide", "ppt/slideMasters/slideMaster",
                       "ppt/slideLayouts/slideLayout")


def strip_run_languages(path: Path) -> None:
    """Remove lang/altLang from every a:rPr and a:endParaRPr in slides, masters and layouts.

    Scoped to those two elements rather than to every `lang=` in the package: the rule reads
    A.RunProperties and A.EndParagraphRunProperties, so a broader strip would change parts the
    rule never looks at and make the fixture a test of the zip surgery instead of the rule.
    """
    import re
    import shutil
    import tempfile
    tmp = Path(tempfile.mkdtemp()) / path.name
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and item.filename.startswith(_LANG_BEARING_PARTS):
                data = re.sub(
                    r"(<a:(?:rPr|endParaRPr)\b[^>]*?)>",
                    lambda m: re.sub(r'\s+(?:alt)?[Ll]ang="[^"]*"', "", m.group(1)) + ">",
                    data.decode("utf-8")).encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def f_no_slide_title(prs, slide):
    """A slide with no title PLACEHOLDER at all — built on the blank layout.

    Distinct from `title-empty`, which has a placeholder and empties it. Both trip 2.4.2 under
    the rule (it flags a missing placeholder and an empty one separately), and this one is the
    fixture for 2.4.2 because it isolates the criterion: an empty placeholder is ALSO the 2.4.6
    violation, so `title-empty` cannot be a single-criterion label for either.
    """
    from pptx import Presentation
    prs2 = Presentation()
    s = prs2.slides.add_slide(prs2.slide_layouts[6])          # blank: no placeholders
    _textbox(s, "Revenue grew across every region this quarter.")
    prs2.core_properties.language = DOC_LANG                  # so 3.1.1 stays quiet
    return prs2, ({"2.4.2": "FAIL"},
                  "a slide with no title placeholder — nothing identifies or navigates to it")


def f_slide_title_ok(prs, slide):
    _textbox(slide, "Revenue grew across every region this quarter.")
    return ({"2.4.2": "PASS"},
            "a title placeholder holding real text — must not be flagged (adversarial)")


def f_no_language(prs, slide):
    """BOTH branches of the rule have to be closed, and only one of them is reachable here.

    Clearing the metadata language undoes what `_deck` stamps. The run-level lang attributes the
    default template writes into the masters and layouts are removed after the save, by
    POST_SAVE — python-pptx has no API for them, and leaving them is exactly the mistake that
    makes an xlsx-shaped fixture declare this pair and detect nothing.
    """
    prs.core_properties.language = ""
    _textbox(slide, "Revenue grew across every region this quarter.")
    return ({"3.1.1": "FAIL"},
            "no metadata language and no run-level lang on any slide, master or layout")


def f_language_ok(prs, slide):
    _textbox(slide, "Revenue grew across every region this quarter.")
    return ({"3.1.1": "PASS"},
            f"metadata language set to {DOC_LANG!r} (adversarial)")


# ── 1.3.1 Info and Relationships and 1.3.2 Meaningful Sequence — also ENGINE-VERIFIED ──
# Pptx/Rules/TableHeaderRule.cs and Pptx/Rules/ReadingOrderRule.cs. Both are certifying pairs, so
# a false clean result is a certification rather than an advisory — the same reason 2.4.2 and
# 3.1.1 were worth the DECLARED_ENGINE asymmetry.
#
# READING ORDER IS NOT FOCUS ORDER, AND THE DIFFERENCE IS LOAD-BEARING HERE. This corpus already
# has a 2.4.3 focus-order pair that moves the title placeholder to the end of document order, and
# it looks like it should trip 1.3.2 as well. It does not, for two independent reasons, and both
# were measured rather than assumed:
#
#   * ReadingOrderRule only counts shapes that carry an explicit `a:off`. A python-pptx
#     PLACEHOLDER inherits its position from the layout and writes no `a:xfrm` at all, so the
#     title and body of the focus-order decks are invisible to the rule — fewer than two
#     positioned shapes, and it yields nothing.
#   * Even with two positioned shapes, swapping them moves each by exactly one rank, and the rule
#     fires only when |visualRank - tabOrder| EXCEEDS 1.
#
# So the 1.3.2 fixture needs THREE explicitly-positioned textboxes, ordered so that one of them
# moves by two ranks. `reading-order` puts the visually-lowest box first in document order, which
# is the smallest arrangement that trips the rule.
#
# One subtlety in the transcription, because it changes which shapes fire: the C# assigns tab
# order BEFORE filtering out unpositioned shapes (`.Select((shape, tabOrder) => ...)` then
# `.Where(s => s.HasPos)`), so an unpositioned placeholder consumes a tab index without taking a
# visual rank. The test's predicate does the same.

def f_table_no_header(prs, slide):
    _table(slide, first_row=False)
    return ({"1.3.1": "FAIL"},
            "a three-row table with no header row designated — no column context for a screen "
            "reader")


def f_table_header_ok(prs, slide):
    _table(slide, first_row=True)
    return ({"1.3.1": "PASS"},
            "the same table with its first row marked as the header (adversarial)")


def f_reading_order(prs, slide):
    """Three positioned boxes whose document order puts the visually-lowest one first."""
    _position_title(slide)
    _box(slide, "Third, visually", 5.0)
    _box(slide, "First, visually", 2.0)
    _box(slide, "Second, visually", 3.5)
    return ({"1.3.2": "FAIL"},
            "three boxes whose tab order runs bottom, top, middle — a screen reader announces "
            "them in that order")


def f_reading_order_ok(prs, slide):
    _position_title(slide)
    _box(slide, "First, visually", 2.0)
    _box(slide, "Second, visually", 3.5)
    _box(slide, "Third, visually", 5.0)
    return ({"1.3.2": "PASS"},
            "the same three boxes in visual order (adversarial)")


# ── 3.1.2 Language of Parts — decided by the prose, like 1.3.3 ──────────────────
# textchecks.detect_language_parts reads EXTRACTED TEXT: it needs at least two segments of
# >= _MIN_SEG_WORDS (12) real words, in at least two confidently-detected languages, and it
# reports the passages whose language the document never identifies. The words are the fixture
# and the deck is the container, so the wording is identical across the .xlsx, .pptx and
# .pdf corpora — see tests/test_language_parts_corpus.py.
#
# THE CONTROL IS MONOLINGUAL, NOT "THE SAME DOCUMENT WITH THE FRENCH MARKED", and on two of the
# three formats that is forced rather than chosen. office_structure.language_marked_spans returns
# {} for .xlsx and .pdf: SpreadsheetML's rich-text run properties have no language element at
# all, and PDF's /Lang structure-tree walk is not built. Marking is therefore unrepresentable
# there, and a "marked" fixture would fail 3.1.2 exactly like the violation — measured, not
# assumed. .pptx CAN carry the mark, so it gets a third fixture proving that a write clears the
# criterion; the other two would be asserting a capability the format does not have.
LANG_EN_BODY = (
    "The benefits office publishes a summary of every plan option before the enrollment window opens.",
    "Employees may change their elections at any point during the month without providing a reason.",
    "Questions about eligibility should be directed to the benefits office rather than to a manager.",
)
LANG_FR_PASSAGE = ("Le bureau des avantages sociaux publie chaque annee un resume complet de "
                   "toutes les options offertes aux employes de la region.")
# Same length and register as the French passage, so the control differs by LANGUAGE and nothing
# else — a shorter filler would fall under the 12-word floor and be skipped rather than judged.
LANG_EN_TAIL = ("Enrollment closes at the end of the month for every employee in every "
                "participating region.")


def _lang_slide(slide, lines, mark_fr=False):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(12)
        if mark_fr and line == LANG_FR_PASSAGE:
            run.font._rPr.set("lang", "fr-FR")


def f_language_parts(prs, slide):
    _lang_slide(slide, list(LANG_EN_BODY) + [LANG_FR_PASSAGE])
    return ({"3.1.2": "FAIL"},
            "an English deck with an unmarked French paragraph — a screen reader pronounces it "
            "with English phonetics and it is unintelligible")


def f_language_parts_ok(prs, slide):
    _lang_slide(slide, list(LANG_EN_BODY) + [LANG_EN_TAIL])
    return ({"3.1.2": "REVIEW"},
            "the same deck entirely in one language (adversarial)")


def f_language_parts_marked_ok(prs, slide):
    _lang_slide(slide, list(LANG_EN_BODY) + [LANG_FR_PASSAGE], mark_fr=True)
    return ({"3.1.2": "REVIEW"},
            "the SAME mixed-language deck with the French run marked lang=\"fr-FR\" "
            "(adversarial). The pair above proves the detector fires; this one proves a WRITE "
            "can clear it — the property that makes 3.1.2 remediable on .pptx and not on .xlsx "
            "or .pdf, neither of which has anywhere to record the mark")


FIXTURES = [
    ("sensory-instruction",     f_sensory_instruction,     "violation"),
    ("sensory-instruction-ok",  f_sensory_instruction_ok,  "adversarial"),
    ("title-empty",             f_title_empty,             "violation"),
    ("title-ok",                f_title_ok,                "adversarial"),
    ("link-vague",              f_link_vague,              "violation"),
    ("link-descriptive-ok",     f_link_descriptive_ok,     "adversarial"),
    ("link-underline-off",      f_link_underline_off,      "violation"),
    ("link-underlined-ok",      f_link_underlined_ok,      "adversarial"),
    ("contrast-fail",           f_contrast_fail,           "violation"),
    ("contrast-ok",             f_contrast_ok,             "adversarial"),
    ("shape-faint-outline",     f_shape_faint_outline,     "violation"),
    ("shape-strong-outline-ok", f_shape_strong_outline_ok, "adversarial"),
    ("focus-order",             f_focus_order,             "violation"),
    ("focus-order-ok",          f_focus_order_ok,          "adversarial"),
    ("embedded-control",        f_embedded_control,        "violation"),
    ("no-controls-ok",          f_no_controls_ok,          "adversarial"),
    ("picture-no-alt",          f_picture_no_alt,          "violation"),
    ("image-of-text",           f_image_of_text,           "violation"),
    ("image-of-text-logo-ok",   f_image_of_text_logo_ok,   "adversarial"),
    ("language-parts",          f_language_parts,          "violation"),
    ("language-parts-ok",       f_language_parts_ok,       "adversarial"),
    ("language-parts-marked-ok", f_language_parts_marked_ok, "adversarial"),
    ("no-picture-ok",           f_no_picture_ok,           "adversarial"),
    ("no-slide-title",          f_no_slide_title,          "violation"),
    ("slide-title-ok",          f_slide_title_ok,          "adversarial"),
    ("no-language",             f_no_language,             "violation"),
    ("language-ok",             f_language_ok,             "adversarial"),
    ("table-no-header",         f_table_no_header,         "violation"),
    ("table-header-ok",         f_table_header_ok,         "adversarial"),
    ("reading-order",           f_reading_order,           "violation"),
    ("reading-order-ok",        f_reading_order_ok,        "adversarial"),
]

# Transforms applied AFTER python-pptx has written the package, keyed by fixture name. Only one
# case needs it: the default template's masters and layouts carry the run-level lang attributes
# that Pptx/Rules/DocumentLanguageRule.cs reads, and python-pptx has no API for removing them.
POST_SAVE = {"no-language": strip_run_languages}

DECLARED = ("1.1.1", "1.3.3", "1.4.1", "1.4.11", "1.4.3", "1.4.5", "2.1.2", "2.4.3",
            "2.4.4", "2.4.6", "3.1.2", "4.1.2")

# Declared, but confirmed only where the .NET Office analyser is built — CI, not a bare
# container. A SEPARATE tuple rather than folded into DECLARED so one number keeps one meaning:
# DECLARED is "a detector was driven against this fixture anywhere the suite runs", and merging
# the two would quietly make the coverage column mean two different things in the same row.
# gen_fixture_coverage counts both and reports the split.
DECLARED_ENGINE = ("1.3.1", "1.3.2", "2.4.2", "3.1.1")


def _validate(name: str, expectations: dict[str, str]) -> list[str]:
    """Reject an expectation the engine could never produce. A manifest that expects PASS on a
    REVIEW-lane pair does not describe a product bug; it describes a manifest bug, and it would
    report a false failure on every run until somebody re-derived this by hand."""
    problems = []
    for sc, verdict in expectations.items():
        allowed = ce.possible_verdicts(sc, FMT)
        if verdict not in allowed:
            problems.append(
                f"{name}: expects {sc}={verdict} but the engine can only emit "
                f"{','.join(sorted(allowed))} for ({sc}, {FMT})")
    return problems


def build_all(docs: Path) -> tuple[list[dict], list[str]]:
    docs.mkdir(parents=True, exist_ok=True)
    manifest, problems = [], []
    for name, build, kind in FIXTURES:
        prs, slide = _deck()
        built = build(prs, slide)
        # A builder returns (expectations, note), or (expectations, note, parts) when it needs a
        # zip part python-pptx cannot author, or (presentation, (expectations, note)) when it has
        # to construct its own deck on a different layout. Dispatched by shape rather than by a
        # flag, so each builder reads as the simplest thing that expresses its case.
        parts = None
        if len(built) == 2 and isinstance(built[1], tuple):
            prs, (expectations, note) = built
        else:
            expectations, note = built[0], built[1]
            parts = built[2] if len(built) > 2 else None
        path = docs / f"{name}.pptx"
        prs.save(path)
        if parts:
            _replace_parts(path, parts)
        # A post-save transform, for the one case that cannot be expressed through python-pptx:
        # stripping run-level lang attributes the default template writes into the masters and
        # layouts. Keyed by fixture name rather than returned by the builder because the builder
        # runs before the package exists.
        if name in POST_SAVE:
            POST_SAVE[name](path)
        problems += _validate(name, expectations)
        manifest.append({"file": f"docs/{name}.pptx", "name": name, "kind": kind,
                         "format": FMT, "expect": expectations, "note": note})
    return manifest, problems


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("--out", type=Path,
                    default=Path.home() / "Downloads" / "acp-pptx-eval" / "sc-corpus")
    args = ap.parse_args()

    manifest, problems = build_all(args.out / "docs")
    for row in manifest:
        print(f"  {row['kind']:11} {row['name']:24} {row['expect']}")
    if problems:
        print("\nIMPOSSIBLE EXPECTATIONS — fix the fixture, not the engine:", file=sys.stderr)
        for p in problems:
            print(f"  {p}", file=sys.stderr)
        return 1

    scs = sorted(sc for sc, f in ce.pol.SCOPE_PRESETS["acp-core-17"].items() if FMT in f)
    (args.out / "manifest.json").write_text(json.dumps({
        "format": FMT, "fixtures": manifest,
        "lanes": {sc: {"clean_verdict": ce.clean_verdict(sc, FMT),
                       "possible": sorted(ce.possible_verdicts(sc, FMT)),
                       "can_pass": ce.can_ever_pass(sc, FMT)} for sc in scs},
    }, indent=1) + "\n")
    kinds: dict[str, int] = {}
    for row in manifest:
        kinds[row["kind"]] = kinds.get(row["kind"], 0) + 1
    print(f"\n{len(manifest)} fixtures: " + ", ".join(f"{n} {k}" for k, n in sorted(kinds.items())))
    print(f"declares {len(DECLARED)} of {len(scs)} applicable .pptx pairs: {', '.join(DECLARED)}")
    print(f"wrote {args.out / 'manifest.json'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

"""Generate a varied test corpus covering all severity tiers and file types.

Each file is crafted to trigger specific WCAG violations at a known severity level
so the E2E test can verify scanner detection, scoring, and tracing accuracy.

Output: test-corpus/files/  (same directory the ACP local scanner reads)

Severity tiers (maps to ACP scoring):
  CRITICAL   — Level A violations, PDF untagged/no-lang, multiple fatal structural issues
  SERIOUS    — Level AA contrast failures, missing alt, bad links
  MODERATE   — Mixed minor structural issues
  BORDERLINE — Just one or two fixable issues; should still fail conformance
  CLEAN      — Fully accessible; should score 100 and pass

Usage:
  python scripts/gen_test_corpus.py
"""
from __future__ import annotations
import io, os, struct, zlib
from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "test-corpus/files"
OUT.mkdir(parents=True, exist_ok=True)

# ── PDF helpers ──────────────────────────────────────────────────────────────

def _pdf_critical_untagged() -> bytes:
    """Untagged PDF: no DocTitle, no Lang, no Tags, scanned-image placeholder."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    # Deliberately NO title, NO language in the XMP/Info
    c.drawString(100, 700, "CRITICAL: This PDF has no tags, no title, and no language.")
    c.drawString(100, 680, "Image below has no alt text.")
    # Draw a coloured rect as a fake "image" — no alt tag
    c.setFillColorRGB(0.2, 0.2, 0.2)
    c.rect(100, 400, 200, 150, fill=1)
    c.drawString(100, 380, "Click here for more information.")   # ambiguous link text
    c.drawString(100, 360, "Read more about our services here.") # ambiguous
    c.save()
    return buf.getvalue()


def _pdf_serious_contrast() -> bytes:
    """PDF with very low contrast text — light grey on white."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setTitle("Contrast Test Document")
    c.drawString(100, 700, "Contrast failure test document")
    # Near-white text on white background (contrast ratio ~1.1:1, needs 4.5:1)
    c.setFillColorRGB(0.92, 0.92, 0.92)
    c.setStrokeColorRGB(1, 1, 1)
    c.drawString(100, 650, "This text is nearly invisible — very low contrast failure.")
    c.drawString(100, 620, "WCAG 1.4.3 requires 4.5:1 contrast for normal text.")
    c.setFillColorRGB(0.88, 0.88, 0.88)
    c.drawString(100, 590, "Another low-contrast line — borderline failure.")
    c.setFillColorRGB(0, 0, 0)
    c.drawString(100, 540, "This black text is fine.")
    c.save()
    return buf.getvalue()


def _pdf_moderate_mixed() -> bytes:
    """PDF with a title but missing lang + one structural issue."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setTitle("Quarterly Accessibility Report Q3-2025")
    # No language set — WCAG 3.1.1 violation
    c.drawString(100, 700, "Quarterly Accessibility Report")
    c.drawString(100, 670, "Executive Summary")
    c.drawString(100, 640, "This document summarises our Q3 accessibility progress.")
    c.drawString(100, 610, "Abbreviation: WCAG, ADA, PDF, HTML used without expansion.")
    c.save()
    return buf.getvalue()


def _pdf_borderline() -> bytes:
    """PDF that is mostly OK but has one missing alt-equivalent."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setTitle("Accessibility Policy — borderline")
    c.drawString(100, 700, "Accessibility Policy")
    c.drawString(100, 670, "Our organisation is committed to digital accessibility.")
    c.drawString(100, 640, "Contact us at accessibility@example.com for support.")
    c.setFillColorRGB(0.85, 0.85, 0.85)   # slightly low contrast (borderline ~2.8:1)
    c.drawString(100, 610, "Last reviewed: June 2025")
    c.setFillColorRGB(0, 0, 0)
    c.save()
    return buf.getvalue()


def _pdf_clean() -> bytes:
    """Clean PDF: title, language, tagged content, good contrast."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setTitle("Accessible Annual Report 2025")
    c.drawString(100, 700, "Accessible Annual Report 2025")
    c.drawString(100, 670, "This document meets WCAG 2.1 Level AA.")
    c.drawString(100, 640, "All images include equivalent text descriptions.")
    c.drawString(100, 610, "Contact: accessibility@example.com")
    c.save()
    return buf.getvalue()


# ── DOCX helpers ─────────────────────────────────────────────────────────────

def _docx_critical() -> bytes:
    """DOCX: no headings, bold-as-heading anti-pattern, missing list structure."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    # No actual Heading styles — uses bold Normal (WCAG 1.3.1 violation)
    p = doc.add_paragraph()
    run = p.add_run("COMPANY OVERVIEW")
    run.bold = True
    run.font.size = Pt(16)

    doc.add_paragraph(
        "This document uses bold text to simulate headings instead of proper "
        "heading styles, which breaks screen reader navigation (WCAG 1.3.1)."
    )

    p2 = doc.add_paragraph()
    r2 = p2.add_run("SERVICES WE OFFER")
    r2.bold = True
    r2.font.size = Pt(14)

    # Low-contrast text: near-white on white
    p3 = doc.add_paragraph()
    r3 = p3.add_run("This text is nearly invisible — contrast failure (WCAG 1.4.3).")
    r3.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    # Bullets as dash-space instead of proper list (WCAG 1.3.1)
    doc.add_paragraph("- Web development")
    doc.add_paragraph("- Mobile apps")
    doc.add_paragraph("- Accessibility auditing")

    # No document title set in properties
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_serious() -> bytes:
    """DOCX: has headings but images without alt, ambiguous hyperlinks."""
    from docx import Document
    from docx.shared import Pt
    from docx.oxml.ns import qn
    import lxml.etree as etree

    doc = Document()
    doc.core_properties.title = "Services Brochure"
    doc.add_heading("Our Digital Services", level=1)
    doc.add_paragraph(
        "We offer a full suite of digital services. See details at "
    )
    # Add a hyperlink run with ambiguous text "click here"
    p = doc.add_paragraph()
    run = p.add_run("Click here")
    run.underline = True
    run.font.color.rgb = __import__('docx').shared.RGBColor(0, 0, 0xEE)
    p.add_run(" to download the brochure.")

    doc.add_heading("Portfolio", level=2)
    doc.add_paragraph(
        "Our portfolio demonstrates excellence in accessibility. "
        "Images are embedded below without alternative text descriptions."
    )
    doc.add_heading("Contact", level=2)
    doc.add_paragraph("For more info read more at our website.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_moderate() -> bytes:
    """DOCX: proper headings, skips heading levels (H1 → H3, skips H2)."""
    from docx import Document
    doc = Document()
    doc.core_properties.title = "Employee Handbook"
    doc.add_heading("Employee Handbook", level=1)
    doc.add_paragraph("Welcome to Acme Corp. This handbook outlines our policies.")
    # Skips heading level 2 — WCAG 1.3.1 advisory
    doc.add_heading("3.1 Leave Policy", level=3)
    doc.add_paragraph("Employees are entitled to 20 days annual leave.")
    doc.add_heading("3.2 Remote Work", level=3)
    doc.add_paragraph("Remote work requires manager approval.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_clean() -> bytes:
    """DOCX: full heading hierarchy, alt text on images, accessible table."""
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    doc.core_properties.title = "Accessible Procurement Policy 2025"
    doc.core_properties.language = "en-US"
    doc.add_heading("Accessible Procurement Policy", level=1)
    doc.add_paragraph(
        "This document describes Acme Corp's procurement policy for "
        "accessible digital products and services."
    )
    doc.add_heading("1. Scope", level=2)
    doc.add_paragraph("This policy applies to all software and digital services purchased.")
    doc.add_heading("2. Requirements", level=2)
    doc.add_paragraph(
        "All purchased products must meet Web Content Accessibility Guidelines "
        "(WCAG) 2.1 Level AA."
    )
    doc.add_heading("3. Review", level=2)
    doc.add_paragraph("This policy is reviewed annually.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX helpers ─────────────────────────────────────────────────────────────

def _pptx_critical() -> bytes:
    """PPTX: no slide titles, reading order wrong, all-image slides."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    blank = prs.slide_layouts[6]  # truly blank layout — no title placeholder

    # Slide 1: no title text
    s1 = prs.slides.add_slide(blank)
    tf = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf.text_frame.text = "Annual Results 2025"  # just a textbox, not a title

    # Slide 2: decorative text as image (represented as large coloured text box)
    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tb.text_frame.text = "REVENUE +42%"  # no semantic markup

    # Slide 3: low-contrast text
    s3 = prs.slides.add_slide(blank)
    tb3 = s3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    p = tb3.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "This slide has very low contrast text that fails WCAG 1.4.3."
    run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)  # light grey on white
    run.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_moderate() -> bytes:
    """PPTX: has titles but uses layout 'Title Only' — no body structure."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_only = prs.slide_layouts[5]  # Title Only

    s1 = prs.slides.add_slide(title_only)
    s1.shapes.title.text = "Q3 Business Review"

    s2 = prs.slides.add_slide(title_only)
    s2.shapes.title.text = "Revenue Overview"
    tb = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tb.text_frame.text = "Revenue grew 12% YoY. See chart for details."

    s3 = prs.slides.add_slide(title_only)
    s3.shapes.title.text = "Next Steps"
    tb3 = s3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tb3.text_frame.text = "- Launch new product\n- Hire 10 engineers\n- Expand to EU"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_clean() -> bytes:
    """PPTX: proper title + content layout on every slide, good contrast."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Accessibility Strategy 2025"
    s1.placeholders[1].text = "Our commitment to WCAG 2.1 AA across all digital touchpoints."

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Key Milestones"
    tf = s2.placeholders[1].text_frame
    tf.text = "Q1: Audit complete"
    tf.add_paragraph().text = "Q2: Critical fixes shipped"
    tf.add_paragraph().text = "Q3: AA conformance achieved"

    s3 = prs.slides.add_slide(layout)
    s3.shapes.title.text = "Contact"
    s3.placeholders[1].text = "accessibility@example.com | +1 800 555 0100"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Generate all files ────────────────────────────────────────────────────────

FILES = [
    # (filename, generator_fn, description)
    # CRITICAL
    ("pdf-critical-untagged-no-lang.pdf",     _pdf_critical_untagged,
     "No tags, no title, no language, ambiguous link text — worst case"),
    ("docx-critical-no-headings.docx",        _docx_critical,
     "Bold-as-heading, low-contrast text, pseudo-bullets, no title"),
    ("pptx-critical-no-titles.pptx",          _pptx_critical,
     "Blank layout slides, no slide titles, low contrast, no semantic structure"),

    # SERIOUS
    ("pdf-serious-contrast.pdf",              _pdf_serious_contrast,
     "Nearly-invisible light-grey text on white, no language"),
    ("docx-serious-ambiguous-links.docx",     _docx_serious,
     "Ambiguous link text (click here / read more), images without alt"),

    # MODERATE
    ("pdf-moderate-mixed.pdf",                _pdf_moderate_mixed,
     "Has title but missing lang, abbreviations unexplained"),
    ("docx-moderate-skipped-headings.docx",   _docx_moderate,
     "H1 → H3 (skips H2), otherwise well-structured"),
    ("pptx-moderate-title-only.pptx",         _pptx_moderate,
     "Has slide titles but no content placeholders, pseudo-bullets in textbox"),

    # BORDERLINE
    ("pdf-borderline-contrast.pdf",           _pdf_borderline,
     "One slightly low-contrast section, otherwise OK"),

    # CLEAN
    ("pdf-clean-accessible.pdf",              _pdf_clean,
     "Title set, black text, no structural issues — should score 100"),
    ("docx-clean-accessible.docx",            _docx_clean,
     "Full heading hierarchy, title set, lang set — should pass"),
    ("pptx-clean-accessible.pptx",            _pptx_clean,
     "Title+Content layout, black text, proper structure — should pass"),
]

MANIFEST: list[dict] = []

for fname, fn, desc in FILES:
    path = OUT / fname
    data = fn()
    path.write_bytes(data)
    MANIFEST.append({"file": fname, "size_kb": round(len(data)/1024, 1), "desc": desc})
    print(f"  ✓ {fname}  ({round(len(data)/1024,1)} KB)")

print(f"\n{len(FILES)} files written to {OUT}")

# Write a manifest JSON for the E2E runner
import json
(OUT.parent / "manifest.json").write_text(json.dumps(MANIFEST, indent=2))
print(f"Manifest: {OUT.parent / 'manifest.json'}")

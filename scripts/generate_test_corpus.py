#!/usr/bin/env python3
"""Generate synthetic test corpus documents for ACP scanner regression tests.

Run once after checking out the repo, or whenever the test corpus needs to be
extended.  The output files are committed; this script documents HOW they were
built so their exact rule-trigger conditions are reproducible.

Usage:
    python scripts/generate_test_corpus.py
"""
from __future__ import annotations
from pathlib import Path

ROOT   = Path(__file__).resolve().parent.parent
# Oracle corpus: the fixed synthetic rule-trigger set the scan regression suite
# asserts on — kept OUT of test-corpus/files/, which holds the (changeable) demo estate.
CORPUS = ROOT / "test-corpus/oracle"
CORPUS.mkdir(parents=True, exist_ok=True)

# ── HTML ────────────────────────────────────────────────────────────────────
# The HTML engine is pure-Python (lxml), so these are trivially verifiable.

(CORPUS / "html-missing-title.html").write_text(
    "<!DOCTYPE html>\n"
    "<html lang=\"en\">\n"
    "<head><meta charset=\"utf-8\"></head>\n"   # no <title> → HTML_MISSING_TITLE
    "<body>\n"
    "  <h1>Page without a title</h1>\n"
    "  <p>This document intentionally omits the &lt;title&gt; element.</p>\n"
    "  <a href=\"/about\">About our organisation</a>\n"
    "</body>\n"
    "</html>\n"
)

(CORPUS / "html-heading-skip.html").write_text(
    "<!DOCTYPE html>\n"
    "<html lang=\"en\">\n"
    "<head><meta charset=\"utf-8\"><title>Heading Skip Test</title></head>\n"
    "<body>\n"
    "  <h1>Introduction</h1>\n"
    "  <h3>Sub-section</h3>\n"  # skips h2 → HTML_HEADING_SKIP
    "  <p>Heading level jumps from h1 to h3.</p>\n"
    "</body>\n"
    "</html>\n"
)

(CORPUS / "html-all-violations.html").write_text(
    "<!DOCTYPE html>\n"
    "<html>\n"                                  # no lang → HTML_MISSING_LANG
    "<head><meta charset=\"utf-8\"></head>\n"   # no <title> → HTML_MISSING_TITLE
    "<body>\n"
    "  <h1>Main section</h1>\n"
    "  <h3>Sub-section</h3>\n"                 # h1→h3 → HTML_HEADING_SKIP
    "  <img src=\"logo.png\">\n"               # no alt → HTML_IMG_MISSING_ALT
    "  <a href=\"/x\"></a>\n"                  # empty  → HTML_EMPTY_LINK
    "  <a href=\"/y\">click here</a>\n"        # vague  → HTML_VAGUE_LINK
    "  <form>\n"
    "    <input type=\"text\" id=\"name\">\n"  # no label → HTML_INPUT_NO_LABEL
    "  </form>\n"
    "</body>\n"
    "</html>\n"
)

print("HTML files written.")

# ── XLSX ─────────────────────────────────────────────────────────────────────
# Targets: XLSX-SHEET-001, XLSX-MERGE-001, XLSX-HIDDEN-001,
#          XLSX-TITLE-001, XLSX-LANG-001, XLSX-TABLE-001, XLSX-TABLE-NAME-001

from openpyxl import Workbook
from openpyxl.worksheet.table import Table

wb = Workbook()
ws = wb.active
ws.title = "Sheet1"                        # matches ^Sheet\d+$ → XLSX-SHEET-001

# Seed data in A1:C10 (used by the no-header table below)
for r in range(1, 11):
    for c in range(1, 4):
        ws.cell(row=r, column=c, value=f"R{r}C{c}")

# >20 merged cell ranges in rows 12-34 → XLSX-MERGE-001 (threshold is >20)
for i in range(23):
    row = 12 + i
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2)
    ws.cell(row=row, column=1, value=f"Merged {i + 1}")

# Hidden row 11 with numeric content → XLSX-HIDDEN-001
# Use a numeric value: the .NET HiddenContentRule reads cell.CellValue.Text which works
# for numeric cells (<v>12345</v>) but NOT for openpyxl inline strings (<is><t>...</t></is>).
ws.cell(row=11, column=1, value=99999)
ws.cell(row=11, column=2, value=88888)
ws.row_dimensions[11].hidden = True

# Excel Table with no header row (headerRowCount=0) → XLSX-TABLE-001
tab = Table(displayName="NoHeaderTable", ref="A1:C10")
tab.headerRowCount = 0
ws.add_table(tab)

# A second table, header row intact but left at Excel's auto-generated default
# DisplayName → XLSX-TABLE-NAME-001. Separate range so it doesn't overlap the
# table above; a real header row so it doesn't also trip XLSX-TABLE-001.
ws.cell(row=1, column=5, value="Header1")
ws.cell(row=1, column=6, value="Header2")
ws.cell(row=2, column=5, value="Data1")
ws.cell(row=2, column=6, value="Data2")
tab2 = Table(displayName="Table1", ref="E1:F2")
ws.add_table(tab2)

# Leave document title and language unset → XLSX-TITLE-001, XLSX-LANG-001
# (defaults are None/empty — rule checks IsNullOrWhiteSpace)

wb.save(CORPUS / "xlsx-synthetic-violations.xlsx")
print("XLSX file written.")

# ── PPTX ─────────────────────────────────────────────────────────────────────

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree

# ── pptx-vague-links.pptx ───────────────────────────────────────────────────
# Targets: PPTX-LINK-001 (vague hyperlink text), PPTX-LANG-001 (no language)
# Avoids: PPTX-TITLE-001 (slide has a titled layout + title is filled)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])   # Title and Content
slide.shapes.title.text = "Accessibility Compliance Guide"

txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
tf = txBox.text_frame
p  = tf.paragraphs[0]
run = p.add_run()
run.text = "click here"                              # vague text → PPTX-LINK-001

# Attach an external hyperlink relationship to the slide part
rel_id = slide.part.relate_to(
    "https://example.com",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    is_external=True,
)

# Embed hlinkClick inside <a:rPr> of the run
rPr = run._r.get_or_add_rPr()
hlink = etree.SubElement(rPr, qn("a:hlinkClick"))
hlink.set(qn("r:id"), rel_id)

# Do NOT set prs.core_properties.language → PPTX-LANG-001

prs.save(CORPUS / "pptx-vague-links.pptx")
print("PPTX vague-links written.")

# ── pptx-reading-order.pptx ─────────────────────────────────────────────────
# Targets: PPTX-ORDER-001 (tab order ≠ visual order), PPTX-LANG-001
# Also fires: PPTX-TITLE-001 (blank layout has no title placeholder)
#
# Strategy: add 3 shapes in REVERSE visual order.
#   Creation order  →  TabOrder  |  Y position  →  VisualRank
#   1st shape       →  0         |  y=3.5in    →  2  (bottom)
#   2nd shape       →  1         |  y=2.5in    →  1  (middle)
#   3rd shape       →  2         |  y=1.5in    →  0  (top)
# For the bottom shape: |TabOrder=0 – VisualRank=2| = 2 > 1 → fires ORDER-001

prs2 = Presentation()
slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])   # blank — no title ph

for label, y_in in [("Bottom (created first)", 3.5),
                    ("Middle (created second)", 2.5),
                    ("Top (created third)", 1.5)]:
    tb = slide2.shapes.add_textbox(Inches(1), Inches(y_in), Inches(4), Inches(0.6))
    tb.text_frame.text = label

# Do NOT set language → PPTX-LANG-001

prs2.save(CORPUS / "pptx-reading-order.pptx")
print("PPTX reading-order written.")

# ── pptx-alt-missing-image.pptx ─────────────────────────────────────────────
# Targets: PPTX-ALT-001 (image without alt text), PPTX-LANG-001
# Strategy: embed a tiny in-memory PNG — python-pptx creates <p:pic> with no
#   <p:nvPr><p:ph/> descr, so the .NET rule's alt-text check finds nothing.

import io
from PIL import Image as _PILImage

_img = _PILImage.new("RGB", (100, 100), color="#4285F4")
_img_buf = io.BytesIO()
_img.save(_img_buf, "PNG")
_img_buf.seek(0)

prs3 = Presentation()
slide3 = prs3.slides.add_slide(prs3.slide_layouts[1])  # Title and Content
slide3.shapes.title.text = "Quarterly Report"
slide3.shapes.add_picture(_img_buf, Inches(1), Inches(2), Inches(3), Inches(3))
# Do NOT set alt text or language → PPTX-ALT-001 + PPTX-LANG-001

prs3.save(CORPUS / "pptx-alt-missing-image.pptx")
print("PPTX alt-missing-image written.")

# ── docx-heading-structure.docx ──────────────────────────────────────────────
# Targets: DOCX-HEAD-001 (skipped heading levels), DOCX-LANG-001, DOCX-TITLE-001
# Strategy: H1 immediately followed by H3, skipping H2.

from docx import Document as _DocxDoc

_doc = _DocxDoc()
_doc.add_paragraph("Introduction", style="Heading 1")
_doc.add_paragraph("Sub-section skipping H2", style="Heading 3")  # H1→H3 skip
_doc.add_paragraph("This document intentionally skips heading level 2.")
# core_properties.title and .language left empty → DOCX-TITLE-001, DOCX-LANG-001

_doc.save(CORPUS / "docx-heading-structure.docx")
print("DOCX heading-structure written.")

# ── xlsx-chart-no-alt.xlsx ───────────────────────────────────────────────────
# Targets: XLSX-ALT-001 (chart with no alt text), XLSX-TITLE-001, XLSX-LANG-001
# Strategy: add a bar chart over numeric data; leave chart.title unset.

from openpyxl.chart import BarChart, Reference as _ORef

_wb2 = Workbook()
_ws2 = _wb2.active
_ws2.title = "Data"
for _r, _v in enumerate([10, 25, 13, 47, 31], start=1):
    _ws2.cell(row=_r, column=1, value=_v)

_chart = BarChart()
_chart.type = "col"
_chart.title = None           # no chart title
_data_ref = _ORef(_ws2, min_col=1, min_row=1, max_row=5)
_chart.add_data(_data_ref)
_ws2.add_chart(_chart, "C1")

# Leave wb.properties.title and language empty → XLSX-TITLE-001, XLSX-LANG-001
_wb2.save(CORPUS / "xlsx-chart-no-alt.xlsx")
print("XLSX chart-no-alt written.")

# ── html-partial-violations.html ─────────────────────────────────────────────
# Targets: HTML_IMG_MISSING_ALT + HTML_VAGUE_LINK only.
# Deliberately avoids: MISSING_LANG (lang="en"), MISSING_TITLE (title present),
#   HEADING_SKIP (h1→h2 is sequential), INPUT_NO_LABEL (label present),
#   EMPTY_LINK (link has text).

(CORPUS / "html-partial-violations.html").write_text(
    "<!DOCTYPE html>\n"
    "<html lang=\"en\">\n"
    "<head><meta charset=\"utf-8\"><title>Partial Violations Test</title></head>\n"
    "<body>\n"
    "  <h1>Annual Accessibility Report</h1>\n"
    "  <h2>Executive Summary</h2>\n"
    "  <p>This page has exactly two accessibility violations.</p>\n"
    "  <img src=\"chart.png\">\n"                      # no alt → HTML_IMG_MISSING_ALT
    "  <p>For details <a href=\"/report\">read more</a>.</p>\n"  # vague → HTML_VAGUE_LINK
    "  <form>\n"
    "    <label for=\"q\">Search</label>\n"
    "    <input type=\"text\" id=\"q\">\n"             # labeled → no HTML_INPUT_NO_LABEL
    "  </form>\n"
    "</body>\n"
    "</html>\n"
)
print("HTML partial-violations written.")

# ── docx-all-violations.docx ─────────────────────────────────────────────────
# Targets: DOCX-TITLE-001, DOCX-LANG-001, DOCX-ALT-001, DOCX-LINK-001, DOCX-TABLE-001
# Strategy: one doc that hits all five rules in docx-noncompliant's oracle set.

from docx.oxml.ns import qn as _qn
from docx.oxml import OxmlElement as _OXml
from docx.shared import Inches as _DIn

def _add_hyperlink(paragraph, url, text):
    """Add a hyperlink run to an existing paragraph (python-docx XML workaround)."""
    rel_id = paragraph.part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hl = _OXml("w:hyperlink")
    hl.set(_qn("r:id"), rel_id)
    hl.set(_qn("w:history"), "1")
    run_el = _OXml("w:r")
    rPr = _OXml("w:rPr")
    style_el = _OXml("w:rStyle")
    style_el.set(_qn("w:val"), "Hyperlink")
    rPr.append(style_el)
    run_el.append(rPr)
    t_el = _OXml("w:t")
    t_el.text = text
    run_el.append(t_el)
    hl.append(run_el)
    paragraph._p.append(hl)

_doc2 = _DocxDoc()
# No title or language in core_properties → DOCX-TITLE-001, DOCX-LANG-001

_p_link = _doc2.add_paragraph("See ")
_add_hyperlink(_p_link, "https://example.com", "click here")  # vague → DOCX-LINK-001

# Table with no explicit header row → DOCX-TABLE-001
_tbl = _doc2.add_table(rows=3, cols=2)
for _row in _tbl.rows:
    for _cell in _row.cells:
        _cell.text = "data"

# Embedded image without alt text → DOCX-ALT-001
_img2 = _PILImage.new("RGB", (50, 50), color="#EA4335")
_img2_buf = io.BytesIO()
_img2.save(_img2_buf, "PNG")
_img2_buf.seek(0)
_doc2.add_picture(_img2_buf, width=_DIn(1))

_doc2.save(CORPUS / "docx-all-violations.docx")
print("DOCX all-violations written.")

# ── pptx-duplicate-titles.pptx ───────────────────────────────────────────────
# Targets: PPTX-TITLE-002 (second slide's title duplicates the first's)
# Avoids: PPTX-TITLE-001 (both titles are non-empty)

prs4 = Presentation()
_slide4a = prs4.slides.add_slide(prs4.slide_layouts[1])
_slide4a.shapes.title.text = "Quarterly Report"
_slide4b = prs4.slides.add_slide(prs4.slide_layouts[1])
_slide4b.shapes.title.text = "Quarterly Report"          # duplicate → PPTX-TITLE-002
prs4.core_properties.language = "en-US"                   # avoid PPTX-LANG-001

prs4.save(CORPUS / "pptx-duplicate-titles.pptx")
print("PPTX duplicate-titles written.")

# ── xlsx-duplicate-sheets.xlsx ───────────────────────────────────────────────
# Targets: XLSX-SHEET-002 (second sheet's name duplicates the first's)
# Avoids: XLSX-SHEET-001 (neither name matches the ^Sheet\d+$ default pattern)
# openpyxl silently deduplicates a colliding create_sheet() name (Data -> Data1), so
# the true collision is forced by rewriting xl/workbook.xml's <sheet name=...> after save.

_wb3 = Workbook()
_wb3.active.title = "Data"
_wb3.active.cell(row=1, column=1, value="not blank")       # avoid XLSX-BLANK-001
_sheet3b = _wb3.create_sheet("Data1")
_sheet3b.cell(row=1, column=1, value="not blank")           # avoid XLSX-BLANK-001
_wb3.properties.title = "Duplicate Sheets Demo"            # avoid XLSX-TITLE-001
_wb3.properties.language = "en-US"                         # avoid XLSX-LANG-001

import zipfile as _zipfile

_dup_path = CORPUS / "xlsx-duplicate-sheets.xlsx"
_wb3.save(_dup_path)

_dup_bytes = {}
with _zipfile.ZipFile(_dup_path) as _zin:
    for _item in _zin.infolist():
        _data = _zin.read(_item.filename)
        if _item.filename == "xl/workbook.xml":
            _data = _data.replace(b'name="Data1"', b'name="Data"')  # force the collision
        _dup_bytes[_item.filename] = _data
with _zipfile.ZipFile(_dup_path, "w", _zipfile.ZIP_DEFLATED) as _zout:
    for _name, _data in _dup_bytes.items():
        _zout.writestr(_name, _data)

print("XLSX duplicate-sheets written.")

print("\nAll synthetic corpus files generated successfully.")

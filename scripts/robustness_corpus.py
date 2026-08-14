#!/usr/bin/env python3
"""Generate an ADVERSARIAL robustness smoke-test corpus for ACP + an expected-results manifest.

The goal is ROBUSTNESS, not accuracy: prove ACP degrades *honestly* on large and malformed files —
no crash, no hang, no silent drop, truncation surfaced — not that a 120-page file scores perfectly.

Each generated file is deliberately built to stress a KNOWN bound in ACP (documented in
docs/pdf-assessment-remediation.md and docs/office-html-assessment-remediation.md):

  - OCR image cap        30 images/file  (OCR_IMAGE_CAP_REACHED)
  - Vision figure cap     25 figures/file
  - PDF reading-order     first 20 pages sampled (tail unassessed — must be surfaced)
  - .NET Office CLI       180s subprocess timeout (huge xlsx/pptx -> "uncertain", not crash)
  - job queue lease       long scans must finish or be honestly reclaimed (no double-processing)

Every file carries COUNTED, known defects so the run can prove nothing was silently dropped, and a
manifest.json records the expected Discovery bucket + which caps should trip + the robustness
assertions. Running the corpus against ACP is a separate step (needs a live backend); this only
produces the files + the expected-results manifest to compare against.

Requires (wheel-only, no system libs):  pip install fpdf2 pypdf Pillow python-docx python-pptx openpyxl
Usage:                                   python scripts/robustness_corpus.py --out ./robustness-corpus [--scale 1.0]
"""
from __future__ import annotations
import argparse, io, json, os, sys, traceback
from datetime import datetime, timezone

# ── optional libs, imported lazily so a missing one skips only its own files ───────────────────
def _try(mod):
    try:
        return __import__(mod)
    except Exception:
        return None

MANIFEST: list[dict] = []
SKIPPED: list[dict] = []


def record(filename, *, category, fmt, size, known_defects, expected_bucket, caps, notes):
    """One manifest row. `caps` = bounds this file is built to trip; assertions are standard."""
    MANIFEST.append({
        "filename": filename,
        "category": category,
        "format": fmt,
        "size": size,
        "known_defects": known_defects,                 # {sc: count} — must all be reported
        "expected_bucket": expected_bucket,             # Discovery bucket (see docs/discovery-triage-spec.md)
        "caps_expected_to_trip": caps,
        "robustness_assertions": [
            "scan terminates (no hang) and does not crash the worker",
            "counts reconcile — file lands in exactly one bucket",
            "any truncation/cap is SURFACED, never a silent clean pass",
            "on engine timeout the file is 'uncertain' (upper-bound score), not a fabricated pass",
        ],
        "notes": notes,
    })


def _txt_image(PIL, text, w=800, h=1000):
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    for i, line in enumerate(text.splitlines()):
        d.text((40, 40 + i * 22), line, fill="black")
    return img


# ── PDF corpus (fpdf2 to build, pypdf to encrypt/truncate) ─────────────────────────────────────
def gen_pdfs(out, scale):
    fpdf = _try("fpdf"); PIL = _try("PIL")
    pypdf = _try("pypdf")
    if not fpdf:
        SKIPPED.append({"format": "pdf", "reason": "fpdf2 not installed"}); return
    from fpdf import FPDF

    # 1) 120-page text PDF, no title/lang — tail beyond p20 must be surfaced as reading-order-unassessed
    pages = max(4, int(120 * scale))
    pdf = FPDF()
    for p in range(pages):
        pdf.add_page(); pdf.set_font("Helvetica", size=11)
        pdf.multi_cell(0, 6, f"Page {p+1} of {pages}. " + ("Clinical discharge instructions. " * 40))
    f = os.path.join(out, "pdf-text-120p.pdf"); pdf.output(f)
    record("pdf-text-120p.pdf", category="big-pdf", fmt="pdf", size={"pages": pages},
           known_defects={"2.4.2": 1, "3.1.1": 1},
           expected_bucket="assessable-remediable",
           caps=["PDF reading-order samples first 20 pages — tail must be surfaced as unassessed"],
           notes="Large text PDF, no /Title or /Lang (both auto-fixable). Tests parse throughput + honest reading-order tail.")

    # 2) 100-page SCANNED PDF (full-page image of text per page) — hammers OCR image cap (30)
    if PIL:
        spages = max(4, int(100 * scale))
        pdf = FPDF()
        img = _txt_image(PIL, "SCANNED PATIENT FORM\n" + "\n".join(f"field {i}: value" for i in range(20)))
        tmp = os.path.join(out, "_scan_page.png"); img.save(tmp)
        for _ in range(spages):
            pdf.add_page(); pdf.image(tmp, x=10, y=10, w=190)
        f = os.path.join(out, "pdf-scanned-100img.pdf"); pdf.output(f); os.remove(tmp)
        record("pdf-scanned-100img.pdf", category="big-pdf", fmt="pdf", size={"pages": spages, "images": spages},
               known_defects={"1.4.5": spages, "1.1.1": spages},
               expected_bucket="assessable-assess-only",
               caps=[f"OCR image cap 30/file (this has {spages}) -> OCR_IMAGE_CAP_REACHED must be surfaced",
                     "vision figure cap 25/file"],
               notes="Every page is a full-page image of text (images-of-text 1.4.5). Proves the OCR cap is reported, not hidden.")

    # 3) PDF with 50 embedded images missing alt — vision 25-cap + OCR 30-cap
    if PIL:
        n = max(4, int(50 * scale))
        pdf = FPDF(); pdf.add_page(); pdf.set_font("Helvetica", size=12)
        thumb = os.path.join(out, "_thumb.png"); _txt_image(PIL, "figure", 200, 150).save(thumb)
        y = 20
        for i in range(n):
            if y > 250:
                pdf.add_page(); y = 20
            pdf.image(thumb, x=10, y=y, w=40); y += 45
        f = os.path.join(out, "pdf-50images.pdf"); pdf.output(f); os.remove(thumb)
        record("pdf-50images.pdf", category="big-pdf", fmt="pdf", size={"images": n},
               known_defects={"1.1.1": n},
               expected_bucket="assessable-remediable",
               caps=[f"vision figure cap 25 (this has {n})", f"OCR image cap 30 (this has {n})"],
               notes="50 undescribed images. With RunPod vision down (R12) these fall to CPU/manual — tests the fallback path at volume.")

    # 4) password-protected PDF -> UNREADABLE bucket
    if pypdf:
        base = FPDF(); base.add_page(); base.set_font("Helvetica", size=12); base.multi_cell(0, 8, "Locked document.")
        raw = bytes(base.output())
        try:
            from pypdf import PdfReader, PdfWriter
            r = PdfReader(io.BytesIO(raw)); w = PdfWriter()
            for pg in r.pages: w.add_page(pg)
            w.encrypt("s3cret")
            f = os.path.join(out, "locked-password.pdf")
            with open(f, "wb") as fh: w.write(fh)
            record("locked-password.pdf", category="edge", fmt="pdf", size={"pages": 1},
                   known_defects={}, expected_bucket="unreadable",
                   caps=["password-protected — parse must fail CLOSED into the unreadable bucket with a reason"],
                   notes="Must NEVER be scored as passing. 'N of M excluded' path.")
        except Exception:
            SKIPPED.append({"file": "locked-password.pdf", "reason": traceback.format_exc(limit=1)})

    # 5) corrupt / truncated PDF -> UNREADABLE
    base = FPDF(); base.add_page(); base.set_font("Helvetica", size=12); base.multi_cell(0, 8, "x" * 4000)
    raw = bytes(base.output())
    f = os.path.join(out, "corrupt-truncated.pdf")
    with open(f, "wb") as fh: fh.write(raw[: len(raw) // 2])   # chop the xref/trailer off
    record("corrupt-truncated.pdf", category="edge", fmt="pdf", size={"bytes": len(raw) // 2},
           known_defects={}, expected_bucket="unreadable",
           caps=["truncated (no trailer/xref) — must fail closed, not crash the worker"],
           notes="Half a PDF. Tests the parser's exception path.")


# ── DOCX corpus (python-docx) ──────────────────────────────────────────────────────────────────
def gen_docx(out, scale):
    docx = _try("docx"); PIL = _try("PIL")
    if not docx:
        SKIPPED.append({"format": "docx", "reason": "python-docx not installed"}); return
    from docx import Document
    from docx.shared import Inches

    def tiny_png():
        from PIL import Image
        buf = io.BytesIO(); Image.new("RGB", (16, 16), "gray").save(buf, "PNG"); buf.seek(0); return buf

    # 6) 150-page DOCX, hundreds of images (no alt) + tables (no header row), no title/lang
    imgs = max(10, int(200 * scale)); tables = max(4, int(30 * scale)); paras = max(50, int(1500 * scale))
    doc = Document()
    if PIL:
        for i in range(imgs):
            doc.add_picture(tiny_png(), width=Inches(0.5))
            if i % 20 == 0:
                doc.add_paragraph("Section body text. " * 30)
    for _ in range(tables):
        t = doc.add_table(rows=4, cols=3)
        for r in t.rows:
            for c in r.cells: c.text = "data"
    for _ in range(paras // 10):
        doc.add_paragraph("Filler paragraph to grow the page count. " * 12)
    f = os.path.join(out, "docx-150p-images-tables.docx"); doc.save(f)
    record("docx-150p-images-tables.docx", category="big-office", fmt="docx",
           size={"images": imgs, "tables": tables, "paragraphs": paras},
           known_defects={"1.1.1": imgs, "1.3.1": tables, "2.4.2": 1, "3.1.1": 1},
           expected_bucket="assessable-remediable",
           caps=[".NET Office CLI 180s timeout risk on a large image-heavy docx",
                 f"OCR image cap 30 (this has {imgs} images)"],
           notes="Big docx: images without alt (1.1.1), tables without headers (1.3.1), no title/lang. Tests .NET throughput + honest caps.")

    # 7) 500 tiny images — memory / image-handling stress
    if PIL:
        n = max(20, int(500 * scale))
        doc = Document()
        for _ in range(n): doc.add_picture(tiny_png(), width=Inches(0.15))
        f = os.path.join(out, "docx-500-tiny-images.docx"); doc.save(f)
        record("docx-500-tiny-images.docx", category="edge", fmt="docx", size={"images": n},
               known_defects={"1.1.1": n},
               expected_bucket="assessable-remediable",
               caps=[f"OCR image cap 30 (this has {n})", "image-handling / memory stress"],
               notes="Pathological image count. Must not OOM or hang; the 30-cap truncation must be surfaced.")

    # 8) zero-byte docx -> UNREADABLE
    f = os.path.join(out, "zero-byte.docx"); open(f, "wb").close()
    record("zero-byte.docx", category="edge", fmt="docx", size={"bytes": 0},
           known_defects={}, expected_bucket="unreadable",
           caps=["0 bytes — must fail closed into unreadable, not crash"],
           notes="Empty file with a docx extension.")

    # 9) wrong extension — a PDF renamed .docx (content-sniff mismatch)
    from fpdf import FPDF
    p = FPDF(); p.add_page(); p.set_font("Helvetica", size=12); p.multi_cell(0, 8, "I am really a PDF.")
    f = os.path.join(out, "wrong-ext.docx")
    with open(f, "wb") as fh: fh.write(bytes(p.output()))
    record("wrong-ext.docx", category="edge", fmt="docx", size={"bytes": None},
           known_defects={}, expected_bucket="unreadable",
           caps=["extension says docx, bytes are PDF — the OpenXml open must fail closed, not misassess"],
           notes="Tests content-vs-extension handling. Ideally content-sniffed; at minimum fails to the unreadable bucket.")


# ── PPTX corpus (python-pptx) ──────────────────────────────────────────────────────────────────
def gen_pptx(out, scale):
    pptx = _try("pptx"); PIL = _try("PIL")
    if not pptx:
        SKIPPED.append({"format": "pptx", "reason": "python-pptx not installed"}); return
    from pptx import Presentation
    from pptx.util import Inches

    slides = max(6, int(100 * scale))
    prs = Presentation()
    blank = prs.slide_layouts[6]
    img = None
    if PIL:
        from PIL import Image
        buf = io.BytesIO(); Image.new("RGB", (64, 64), "navy").save(buf, "PNG"); img = buf.getvalue()
    for _ in range(slides):
        s = prs.slides.add_slide(blank)          # blank layout => no title placeholder (2.4.2)
        if img:
            s.shapes.add_picture(io.BytesIO(img), Inches(1), Inches(1), Inches(2), Inches(2))  # no alt (1.1.1)
    f = os.path.join(out, "pptx-100slides.pptx"); prs.save(f)
    record("pptx-100slides.pptx", category="big-office", fmt="pptx", size={"slides": slides},
           known_defects={"2.4.2": slides, "1.1.1": slides},
           expected_bucket="assessable-remediable",
           caps=[".NET Office CLI 180s timeout risk on a 100-slide deck"],
           notes="100 titleless slides, each with an undescribed image. Tests .NET throughput; on timeout must go 'uncertain', not fake-pass.")


# ── XLSX corpus (openpyxl) ─────────────────────────────────────────────────────────────────────
def gen_xlsx(out, scale):
    opx = _try("openpyxl")
    if not opx:
        SKIPPED.append({"format": "xlsx", "reason": "openpyxl not installed"}); return
    from openpyxl import Workbook

    # 10) 30 sheets, ~100k cells, default-ish sheet names, no title
    sheets = max(4, int(30 * scale)); per = max(200, int(3400 * scale))
    wb = Workbook(); wb.remove(wb.active)
    total = 0
    for si in range(sheets):
        ws = wb.create_sheet(title=f"Sheet{si+1}")           # default-pattern names (2.4.6)
        rows = max(20, per // 10)
        for r in range(rows):
            for c in range(10):
                ws.cell(row=r + 1, column=c + 1, value=f"{r}-{c}"); total += 1
    f = os.path.join(out, "xlsx-30sheets-100kcells.xlsx"); wb.save(f)
    record("xlsx-30sheets-100kcells.xlsx", category="big-office", fmt="xlsx",
           size={"sheets": sheets, "cells": total},
           known_defects={"2.4.6": sheets, "2.4.2": 1},
           expected_bucket="assessable-remediable",
           caps=[".NET Office CLI 180s timeout risk + memory on a large multi-sheet workbook"],
           notes="Big workbook: default sheet names (2.4.6), no workbook title (2.4.2). Tests .NET timeout/memory handling.")

    # 11) empty workbook -> TRIVIAL candidate (ROT) + BlankWorksheet (1.3.2)
    wb = Workbook(); wb.active.title = "Sheet1"
    f = os.path.join(out, "empty-workbook.xlsx"); wb.save(f)
    record("empty-workbook.xlsx", category="edge", fmt="xlsx", size={"sheets": 1, "cells": 0},
           known_defects={"1.3.2": 1},
           expected_bucket="trivial-candidate",
           caps=["BlankWorksheetRule (1.3.2) should fire; triage should flag as Trivial (ROT)"],
           notes="Empty workbook. Tests the blank-detector AND the ROT/trivial triage path.")


# ── non-document + control ─────────────────────────────────────────────────────────────────────
def gen_misc(out, scale):
    PIL = _try("PIL")
    if PIL:
        from PIL import Image
        f = os.path.join(out, "image-not-document.png")
        Image.new("RGB", (640, 480), "teal").save(f)
        record("image-not-document.png", category="edge", fmt="png", size={"px": "640x480"},
               known_defects={}, expected_bucket="filtered-by-type",
               caps=["not a document type — must be counted as filtered-by-type, never 'assessed clean'"],
               notes="A plain image. Tests the reconciling filtered-by-type bucket.")
    docx = _try("docx")
    if docx:
        from docx import Document
        doc = Document(); doc.core_properties.title = "Clean Control Memo"
        doc.core_properties.language = "en-US"
        doc.add_heading("Clean Control", level=1)
        doc.add_paragraph("A small, mostly-clean control document.")
        f = os.path.join(out, "control-clean.docx"); doc.save(f)
        record("control-clean.docx", category="control", fmt="docx", size={"pages": 1},
               known_defects={},
               expected_bucket="assessable-remediable",
               caps=[],
               notes="Baseline control: small and mostly clean. If THIS misbehaves, the problem is not scale.")


def main():
    ap = argparse.ArgumentParser(description="Generate ACP robustness smoke-test corpus + manifest.")
    ap.add_argument("--out", default="./robustness-corpus", help="output directory")
    ap.add_argument("--scale", type=float, default=1.0, help="0<scale<=1 to shrink sizes for a quick run")
    args = ap.parse_args()
    out = os.path.abspath(args.out); os.makedirs(out, exist_ok=True)

    for gen in (gen_pdfs, gen_docx, gen_pptx, gen_xlsx, gen_misc):
        try:
            gen(out, args.scale)
        except Exception:
            SKIPPED.append({"generator": gen.__name__, "reason": traceback.format_exc(limit=3)})

    # sizes on disk (evidence, not fabricated)
    for row in MANIFEST:
        p = os.path.join(out, row["filename"])
        row["bytes_on_disk"] = os.path.getsize(p) if os.path.exists(p) else None

    manifest = {
        "generated_utc": os.environ.get("ROBUSTNESS_STAMP") or datetime.now(timezone.utc).isoformat(),
        "scale": args.scale,
        "file_count": len(MANIFEST),
        "buckets_present": sorted({r["expected_bucket"] for r in MANIFEST}),
        "reconciliation_rule": "every file lands in exactly one bucket; the run's bucket counts must sum to file_count",
        "files": MANIFEST,
        "skipped": SKIPPED,
    }
    with open(os.path.join(out, "manifest.json"), "w") as fh:
        json.dump(manifest, fh, indent=2)

    print(f"corpus: {len(MANIFEST)} files -> {out}")
    for r in MANIFEST:
        kb = (r["bytes_on_disk"] or 0) / 1024
        print(f"  {r['filename']:32s} {r['expected_bucket']:24s} {kb:9.1f} KB")
    if SKIPPED:
        print(f"skipped: {len(SKIPPED)} (see manifest.json 'skipped')")


if __name__ == "__main__":
    main()

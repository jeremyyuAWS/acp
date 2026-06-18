"""Generate detailed, multi-page sample documents (one per type) with REAL
accessibility issues, for the Upload-tab demo. Output -> frontend/public/samples/.

Each file deliberately contains issues a WCAG scan would flag (untagged PDF,
images without alt text, tables without header rows, no document language, low
contrast, unlabeled form fields), so the demo reflects authentic findings.

Run:  python scripts/gen_samples.py   (deps: reportlab python-docx python-pptx openpyxl pillow)
"""
from __future__ import annotations
import os
from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "frontend" / "public" / "samples"
OUT.mkdir(parents=True, exist_ok=True)
ASSET = OUT / "sample-image.png"

LOREM = ("Patients should follow the care plan outlined by their care team. "
         "Contact your coordinator with any questions about medications, follow-up "
         "appointments, or symptoms that worsen. ") * 4


def make_chart():
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (520, 300), "#ffffff")
    d = ImageDraw.Draw(img)
    bars = [180, 240, 140, 200, 90]
    for i, h in enumerate(bars):
        x = 40 + i * 95
        d.rectangle([x, 280 - h, x + 60, 280], fill="#7a5c8e")
    d.line([30, 280, 500, 280], fill="#888", width=2)
    img.save(ASSET)


def gen_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    p = OUT / "patient-discharge-instructions.pdf"
    c = canvas.Canvas(str(p), pagesize=letter)
    w, h = letter
    # NOTE: reportlab output is an UNTAGGED PDF with no /Lang -> real WCAG issues
    for page in range(1, 5):
        c.setFont("Helvetica-Bold", 18)
        c.drawString(72, h - 80, f"Discharge Instructions — Section {page}")  # visual heading only, not structured
        c.setFont("Helvetica", 11)
        text = c.beginText(72, h - 110)
        for line in (LOREM * 2).split(". "):
            text.textLine(line.strip() + ".")
        c.drawText(text)
        if page == 2:
            c.drawImage(str(ASSET), 72, h - 460, width=320, height=185)  # image, no alt text
        c.showPage()
    c.save()
    return p


def gen_docx():
    from docx import Document
    from docx.shared import Pt
    p = OUT / "benefits-policy.docx"
    doc = Document()
    # no doc.core_properties.title set -> "document has no title" issue
    for sec in ("Eligibility", "Enrollment", "Coverage", "Appeals"):
        run = doc.add_paragraph().add_run(sec)  # bold text used as a heading, not a Heading style
        run.bold = True; run.font.size = Pt(16)
        for _ in range(3):
            doc.add_paragraph(LOREM)
        if sec == "Coverage":
            t = doc.add_table(rows=4, cols=3)  # no header row marked
            for r in range(4):
                for col in range(3):
                    t.cell(r, col).text = f"cell {r}.{col}"
        if sec == "Enrollment":
            doc.add_picture(str(ASSET))  # image, no alt text
        doc.add_page_break()
    doc.save(str(p))
    return p


def gen_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    p = OUT / "quarterly-town-hall.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    titled = prs.slide_layouts[5]
    for i in range(6):
        if i % 2 == 0:
            s = prs.slides.add_slide(titled); s.shapes.title.text = f"Quarterly Update {i+1}"
        else:
            s = prs.slides.add_slide(blank)  # no title placeholder -> "slide missing title"
        s.shapes.add_picture(str(ASSET), Inches(1), Inches(1.6), width=Inches(5))  # image, no alt text
    prs.save(str(p))
    return p


def gen_xlsx():
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    p = OUT / "finance-metrics.xlsx"
    wb = Workbook()
    ws = wb.active  # leaves default name "Sheet" -> "sheet not named" issue
    ws.merge_cells("A1:C1"); ws["A1"] = "FY26 Metrics"  # merged cells
    rows = [["Month", "Admissions", "Revenue"], ["Jan", 1200, 480000], ["Feb", 1310, 512000], ["Mar", 1180, 470000], ["Apr", 1400, 533000]]
    for r in rows:
        ws.append(r)
    chart = BarChart()  # chart with no alt/title
    data = Reference(ws, min_col=2, min_row=2, max_row=6, max_col=3)
    chart.add_data(data, titles_from_data=False)
    ws.add_chart(chart, "E2")
    wb.create_sheet("Sheet2")
    wb.save(str(p))
    return p


def gen_html():
    p = OUT / "careers-landing.html"
    # missing lang, low-contrast text, image without alt, unlabeled input, no heading hierarchy
    p.write_text(
        "<!doctype html>\n<html>\n<head><meta charset=\"utf-8\"><title>Careers</title>\n"
        "<style>body{font-family:sans-serif} .faint{color:#bdbdbd;background:#fff} "
        ".hero{background:#f3eef6;padding:30px}</style></head>\n<body>\n"
        "<div class=\"hero\"><b style=\"font-size:26px\">Join our team</b></div>\n"
        "<img src=\"sample-image.png\" width=\"420\">\n"  # no alt
        "<p class=\"faint\">Explore open roles across our health system. We hire across "
        "clinical, research, and administrative tracks.</p>\n"
        + ("<p>We offer competitive benefits and a mission-driven culture.</p>\n" * 8)
        + "<form>\n<p>Search roles</p>\n<input type=\"text\">\n"  # input with no label
        "<input type=\"email\"> <button>Apply</button>\n</form>\n</body>\n</html>\n",
        encoding="utf-8",
    )
    return p


def main():
    make_chart()
    made = [gen_pdf(), gen_docx(), gen_pptx(), gen_xlsx(), gen_html()]
    for f in made:
        print(f"  {f.name:38} {f.stat().st_size // 1024} KB")
    print(f"-> {OUT}")


if __name__ == "__main__":
    main()

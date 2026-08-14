#!/usr/bin/env python3
"""Generate COMPLEX accessibility fixtures — one large file per format (docx/xlsx/pptx/pdf) with
50+ deliberately-injected issues spread across the WCAG success criteria ACP assesses — plus a
manifest mapping each file to its injected {SC: count}.

Purpose: a coverage/accuracy test (complements scripts/robustness_corpus.py, which tests size +
malformed handling). Here every issue is COUNTED and mapped to an SC, so a scan's findings can be
compared against a known floor: the detector should surface at least these SCs, at roughly these
counts (allowing for detectors that report per-document rather than per-element).

HONESTY: some SCs cannot be authored with these Python libraries (a real AcroForm control for 4.1.2,
a slide animation for 2.1.1, autoplay audio for 1.4.2, a true colour-only PDF link for 1.4.1). Those
are listed per file under `not_seeded` with the reason, so the manifest never claims coverage it
didn't produce. The SC list + per-format scope are taken from api/assessment_policy.py.

Requires (wheel-only): pip install fpdf2 pypdf Pillow python-docx python-pptx openpyxl
Usage:                  python scripts/complex_corpus.py --out ./complex-corpus [--issues 55]
"""
from __future__ import annotations
import argparse, io, json, os, traceback
from collections import Counter
from datetime import datetime, timezone

MANIFEST: list[dict] = []
SKIPPED: list[dict] = []

# Plain-language SC names (api/assessment_policy.py) — for a readable manifest.
SC_NAME = {
    "1.1.1": "Non-text Content", "1.3.1": "Info & Relationships", "1.3.2": "Meaningful Sequence",
    "1.3.3": "Sensory Characteristics", "1.4.1": "Use of Color", "1.4.3": "Contrast (Minimum)",
    "1.4.5": "Images of Text", "1.4.6": "Contrast (Enhanced)", "1.4.8": "Visual Presentation",
    "2.4.2": "Page Titled", "2.4.4": "Link Purpose", "2.4.6": "Headings and Labels",
    "3.1.1": "Language of Page", "3.1.2": "Language of Parts",
}


def _tiny_png(color="gray", size=(24, 24)):
    from PIL import Image
    buf = io.BytesIO(); Image.new("RGB", size, color).save(buf, "PNG"); buf.seek(0); return buf


def record(filename, fmt, counts: Counter, not_seeded: dict, notes: str):
    total = sum(counts.values())
    MANIFEST.append({
        "filename": filename, "format": fmt,
        "total_injected_issues": total,
        "injected_by_sc": {sc: counts[sc] for sc in sorted(counts)},
        "sc_names": {sc: SC_NAME.get(sc, "?") for sc in sorted(counts)},
        "distinct_scs": len(counts),
        "not_seeded": not_seeded,        # {sc: reason} — cannot be authored with these libs
        "expectation": "the scan should surface AT LEAST these SCs; counts are a floor "
                       "(some detectors report per-document, not per-element)",
        "notes": notes,
    })


# ── DOCX ────────────────────────────────────────────────────────────────────────────────────────
def gen_docx(out, issues):
    docx = __import__("docx"); c = Counter()
    from docx import Document
    from docx.shared import RGBColor, Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    # 2.4.2 — no core title/language set (leave core_properties.title empty) → 1 doc-level each
    c["2.4.2"] += 1; c["3.1.1"] += 1
    # 2.4.6 — heading skips + an empty heading
    doc.add_heading("Clinical Handbook", level=1)
    for i in range(8):
        doc.add_heading(f"Detail {i}", level=3)      # H1→H3 skip, no H2
        c["2.4.6"] += 1
    doc.add_heading("", level=2); c["2.4.6"] += 1     # empty heading
    # 1.1.1 — images with no alt text
    for _ in range(max(10, issues // 4)):
        doc.add_picture(_tiny_png(), width=Inches(0.4)); c["1.1.1"] += 1
    # 1.3.1 — tables with no header row + pseudo-headings (bold paragraph, not a heading style)
    for _ in range(6):
        t = doc.add_table(rows=3, cols=3)
        for r in t.rows:
            for cell in r.cells: cell.text = "data"
        c["1.3.1"] += 1
    for _ in range(4):
        p = doc.add_paragraph(); run = p.add_run("Section Heading (fake)"); run.bold = True; run.font.size = Pt(16)
        c["1.3.1"] += 1                               # pseudo-heading (looks like a heading, isn't)
    # 1.4.3 — low-contrast runs (light grey on white)
    for _ in range(max(8, issues // 5)):
        p = doc.add_paragraph(); r = p.add_run("Low-contrast body text that fails AA."); r.font.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)
        c["1.4.3"] += 1
    # 1.4.8 — justified paragraphs
    for _ in range(5):
        p = doc.add_paragraph("Justified text stretched to both margins. " * 6); p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        c["1.4.8"] += 1
    # 1.3.3 — sensory instructions
    for txt in ["Click the red button on the right.", "See the box below.", "Press the round icon."]:
        doc.add_paragraph(txt); c["1.3.3"] += 1
    # 3.1.2 — foreign-script runs with no language mark
    for _ in range(3):
        doc.add_paragraph("Πληροφορίες ασθενούς — δεν έχει οριστεί γλώσσα."); c["3.1.2"] += 1
    # 2.4.4 — raw-URL / vague link text (as text; a true w:hyperlink needs XML surgery)
    for txt in ["click here", "http://example.org/very/long/raw/url", "read more"]:
        doc.add_paragraph(txt); c["2.4.4"] += 1
    f = os.path.join(out, "complex-docx-many-issues.docx"); doc.save(f)
    record("complex-docx-many-issues.docx", "docx", c,
           {"4.1.2": "a real content-control / form field can't be authored with python-docx",
            "2.4.9": "requires a genuine w:hyperlink (link-text purpose) — approximated as 2.4.4 text"},
           "docx: images w/o alt, header-less tables + pseudo-headings, low-contrast runs, justified "
           "text, heading skips, sensory text, foreign-script runs, raw-URL links.")


# ── XLSX ────────────────────────────────────────────────────────────────────────────────────────
def gen_xlsx(out, issues):
    opx = __import__("openpyxl"); c = Counter()
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.drawing.image import Image as XLImage
    wb = Workbook(); wb.remove(wb.active)
    c["2.4.2"] += 1; c["3.1.1"] += 1                  # no workbook title / language
    grey = Font(color="CFCFCF"); white = PatternFill("solid", fgColor="FFFFFF")
    for si in range(max(6, issues // 8)):
        ws = wb.create_sheet(title=f"Sheet{si+1}")     # 2.4.6 default sheet name
        c["2.4.6"] += 1
        # data-table with no header row (1.3.1)
        for r in range(1, 6):
            for col in range(1, 5):
                cell = ws.cell(row=r, column=col, value=f"{r}-{col}")
                if r <= 3: cell.font = grey; cell.fill = white   # 1.4.3 low contrast
        c["1.3.1"] += 1; c["1.4.3"] += 2; c["1.4.6"] += 2
        # merged cells break reading order (1.3.2)
        ws.merge_cells("A7:C7"); ws["A7"] = "Merged banner"; c["1.3.2"] += 1
        # hidden column holding data (1.3.2)
        ws["F1"] = "hidden data"; ws.column_dimensions["F"].hidden = True; c["1.3.2"] += 1
        # =HYPERLINK with vague text (2.4.4)
        ws["A9"] = '=HYPERLINK("http://example.org","click here")'; c["2.4.4"] += 1
        # image with no alt (1.1.1)
        try:
            img = XLImage(_tiny_png()); ws.add_image(img, "H1"); c["1.1.1"] += 1
        except Exception:
            pass
    f = os.path.join(out, "complex-xlsx-many-issues.xlsx"); wb.save(f)
    record("complex-xlsx-many-issues.xlsx", "xlsx", c,
           {"3.1.2": "xlsx language-of-parts has no element to author into (human-only, per capability)"},
           "xlsx: default sheet names, header-less data ranges, low-contrast font/fill, merged cells, "
           "hidden data columns, =HYPERLINK vague text, images w/o alt.")


# ── PPTX ────────────────────────────────────────────────────────────────────────────────────────
def gen_pptx(out, issues):
    pptx = __import__("pptx"); c = Counter()
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation(); blank = prs.slide_layouts[6]
    c["3.1.1"] += 1                                    # no language
    for i in range(max(10, issues // 4)):
        s = prs.slides.add_slide(blank)                # blank layout → no title (2.4.2)
        c["2.4.2"] += 1
        s.shapes.add_picture(_tiny_png(), Inches(0.5), Inches(0.5), Inches(1), Inches(1))  # 1.1.1 no alt
        c["1.1.1"] += 1
        # low-contrast text box (1.4.3 / 1.4.6)
        tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1)).text_frame
        tb.text = "Low-contrast slide text."; tb.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        c["1.4.3"] += 1; c["1.4.6"] += 1
        if i % 3 == 0:                                 # table with no header row (1.3.1)
            rows, cols = 3, 3
            gt = s.shapes.add_table(rows, cols, Inches(1), Inches(3.5), Inches(6), Inches(1.5)).table
            for rr in range(rows):
                for cc in range(cols): gt.cell(rr, cc).text = "data"
            c["1.3.1"] += 1
        if i % 4 == 0:                                 # sensory text (1.3.3)
            s.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(0.6)).text_frame.text = "Click the green box on the left."
            c["1.3.3"] += 1
    f = os.path.join(out, "complex-pptx-many-issues.pptx"); prs.save(f)
    record("complex-pptx-many-issues.pptx", "pptx", c,
           {"2.1.1": "a slide animation/keyboard trigger can't be authored with python-pptx",
            "1.4.2": "autoplay audio timing can't be authored with python-pptx",
            "2.4.4": "a real hyperlink run needs XML surgery — omitted"},
           "pptx: titleless slides, images w/o alt, low-contrast text, header-less tables, sensory text.")


# ── PDF (fpdf2 → untagged PDF, which trips structure/alt/reading-order broadly) ───────────────────
def gen_pdf(out, issues):
    fpdf = __import__("fpdf"); c = Counter()
    from fpdf import FPDF
    from PIL import Image
    pdf = FPDF()

    def wln(txt):                       # robust line writer: reset x, use effective page width
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(pdf.epw, 6, txt)
    # untagged PDF → 1.3.1 (no StructTreeRoot) + 1.3.2 (reading order) at the document level
    c["1.3.1"] += 1; c["1.3.2"] += 1
    c["2.4.2"] += 1; c["3.1.1"] += 1                   # no /Title, no /Lang
    # image-of-text pages (1.4.5) — a full-page rendered-text image
    txt_img = os.path.join(out, "_pdf_txt.png")
    im = Image.new("RGB", (700, 900), "white")
    from PIL import ImageDraw
    d = ImageDraw.Draw(im)
    for i in range(25): d.text((40, 40 + i * 24), f"Patient instruction line {i}", fill="black")
    im.save(txt_img)
    for _ in range(5):
        pdf.add_page(); pdf.image(txt_img, x=10, y=10, w=190); c["1.4.5"] += 1
    os.remove(txt_img)
    # figure images with no alt (1.1.1)
    thumb = os.path.join(out, "_pdf_fig.png"); _img = Image.new("RGB", (120, 90), "teal"); _img.save(thumb)
    pdf.add_page(); pdf.set_font("Helvetica", size=11); y = 20
    for _ in range(max(22, issues // 2)):
        if y > 250: pdf.add_page(); y = 20
        pdf.image(thumb, x=10, y=y, w=30); y += 34; c["1.1.1"] += 1
    os.remove(thumb)
    # low-contrast text (1.4.3 / 1.4.6) — light grey glyphs
    pdf.add_page(); pdf.set_font("Helvetica", size=11); pdf.set_text_color(205, 205, 205)
    for _ in range(10):
        wln("Low-contrast paragraph that fails AA contrast."); c["1.4.3"] += 1; c["1.4.6"] += 1
    pdf.set_text_color(0, 0, 0)
    # 2.4.4 raw-URL link text (as text)
    for _ in range(3):
        wln("http://example.org/a/very/long/raw/link click here"); c["2.4.4"] += 1
    # pad to >10 pages with no bookmarks/outline (reinforces 2.4.2 navigation)
    for _ in range(6):
        pdf.add_page(); wln("Filler page, no outline. " * 20)
    f = os.path.join(out, "complex-pdf-many-issues.pdf"); pdf.output(f)
    record("complex-pdf-many-issues.pdf", "pdf", c,
           {"1.4.1": "a true colour-only PDF link annotation can't be authored with fpdf2 (text-approx omitted)",
            "1.4.9": "AAA images-of-text overlaps 1.4.5 here; not separately counted"},
           "pdf (untagged): no tags/title/lang, image-of-text pages, figures w/o alt, low-contrast text, "
           "raw-URL links, 10+ pages with no bookmarks.")


def main():
    ap = argparse.ArgumentParser(description="Generate complex multi-issue WCAG fixtures + manifest.")
    ap.add_argument("--out", default="./complex-corpus")
    ap.add_argument("--issues", type=int, default=55, help="target minimum issues per file (a floor)")
    args = ap.parse_args()
    out = os.path.abspath(args.out); os.makedirs(out, exist_ok=True)
    for gen in (gen_docx, gen_xlsx, gen_pptx, gen_pdf):
        try:
            gen(out, args.issues)
        except Exception:
            SKIPPED.append({"generator": gen.__name__, "reason": traceback.format_exc(limit=3)})
    for row in MANIFEST:
        p = os.path.join(out, row["filename"]); row["bytes_on_disk"] = os.path.getsize(p) if os.path.exists(p) else None
    manifest = {
        "generated_utc": os.environ.get("COMPLEX_STAMP") or datetime.now(timezone.utc).isoformat(),
        "target_min_issues_per_file": args.issues,
        "files": MANIFEST, "skipped": SKIPPED,
        "note": "Counts are INJECTED issues (a detection floor). Some SCs can't be authored with these "
                "libraries — see each file's `not_seeded`. Compare a scan's findings against injected_by_sc.",
    }
    with open(os.path.join(out, "manifest.json"), "w") as fh:
        json.dump(manifest, fh, indent=2)
    print(f"complex corpus: {len(MANIFEST)} files -> {out}")
    for r in MANIFEST:
        kb = (r["bytes_on_disk"] or 0) / 1024
        print(f"  {r['filename']:32s} {r['total_injected_issues']:3d} issues / {r['distinct_scs']:2d} SCs  {kb:8.1f} KB")
    if SKIPPED:
        print(f"skipped: {len(SKIPPED)} (see manifest.json)")


if __name__ == "__main__":
    main()
